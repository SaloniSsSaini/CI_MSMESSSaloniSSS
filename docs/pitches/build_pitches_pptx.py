"""
Generate PowerPoint pitch decks (Investor / MSME / Bank) using the repo's
existing markdown narratives + SVG product visuals under docs/pitches/assets/.
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS_SVG_DIR = ROOT / "assets"
ASSETS_PNG_DIR = ROOT / "assets_png"
OUT_DIR = ROOT / "pptx"


def _emu(inches: float) -> int:
    # python-pptx expects EMU (914400 per inch)
    return int(inches * 914400)


def ensure_dirs() -> None:
    ASSETS_PNG_DIR.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def svg_to_png(svg_path: Path, png_path: Path, *, width_px: int = 2200) -> None:
    # PPT renders PNGs reliably; SVG support is inconsistent across versions.
    # Some repo SVGs may include stray control characters or CP1252 punctuation,
    # which makes them invalid XML for parsers like ElementTree.
    raw = svg_path.read_bytes()
    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        text = raw.decode("cp1252")

    def _ok(ch: str) -> bool:
        o = ord(ch)
        if ch in ("\t", "\n", "\r"):
            return True
        if o < 0x20:
            return False
        # XML 1.0 disallows C1 control chars too.
        if 0x7F <= o <= 0x9F:
            return False
        return True

    cleaned = "".join(ch for ch in text if _ok(ch)).encode("utf-8")

    cairosvg.svg2png(
        bytestring=cleaned,
        write_to=str(png_path),
        output_width=width_px,
        background_color="white",
    )


def convert_all_svgs() -> dict[str, Path]:
    """
    Returns a mapping from asset stem to PNG path, e.g.:
      {"product-dashboard-web": ".../assets_png/product-dashboard-web.png"}
    """
    pngs: dict[str, Path] = {}
    for svg in sorted(ASSETS_SVG_DIR.glob("*.svg")):
        png = ASSETS_PNG_DIR / f"{svg.stem}.png"
        try:
            svg_to_png(svg, png)
        except Exception as e:  # noqa: BLE001 - surfaced with context
            raise RuntimeError(f"Failed converting SVG to PNG: {svg}") from e
        pngs[svg.stem] = png
    return pngs


@dataclass(frozen=True)
class Theme:
    title: RGBColor = RGBColor(0x11, 0x11, 0x11)
    body: RGBColor = RGBColor(0x22, 0x22, 0x22)
    muted: RGBColor = RGBColor(0x66, 0x66, 0x66)
    accent: RGBColor = RGBColor(0x18, 0xA0, 0x7B)  # teal/green


THEME = Theme()


def make_presentation() -> Presentation:
    prs = Presentation()
    # Force 16:9 widescreen (13.333 x 7.5 in)
    prs.slide_width = _emu(13.333)
    prs.slide_height = _emu(7.5)
    return prs


def _set_run_font(run, *, size_pt: int, color: RGBColor, bold: bool = False, name: str = "Calibri"):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(
    prs: Presentation,
    *,
    top_title: str,
    line1: str,
    line2: str,
    footer: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Centered title block (template-like)
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(11.333)
    height = Inches(2.8)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = top_title
    _set_run_font(r, size_pt=44, color=THEME.title, bold=False)

    for text, size, color in [(line1, 20, THEME.body), (line2, 20, THEME.body)]:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = text
        _set_run_font(r2, size_pt=size, color=color, bold=False)

    if footer:
        f = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.8), Inches(0.35))
        ft = f.text_frame
        ft.clear()
        fp = ft.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        fr = fp.add_run()
        fr.text = footer
        _set_run_font(fr, size_pt=12, color=THEME.muted)


def add_section_slide(prs: Presentation, *, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Accent band
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.2),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = THEME.accent
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.2), Inches(12.0), Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _set_run_font(r, size_pt=34, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(12.0), Inches(1.0))
        tf2 = tb2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        _set_run_font(r2, size_pt=22, color=THEME.body)


def add_bullets_slide(prs: Presentation, *, title: str, bullets: list[str], note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(12.0), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, size_pt=32, color=THEME.title, bold=True)

    body = slide.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(11.6), Inches(5.3))
    btf = body.text_frame
    btf.clear()
    btf.word_wrap = True

    for i, b in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.level = 0
        para.text = b
        para.font.size = Pt(22)
        para.font.name = "Calibri"
        para.font.color.rgb = THEME.body

    if note:
        n = slide.shapes.add_textbox(Inches(1.1), Inches(6.95), Inches(11.6), Inches(0.4))
        ntf = n.text_frame
        ntf.clear()
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = note
        _set_run_font(nr, size_pt=14, color=THEME.muted)


def add_image_slide(
    prs: Presentation,
    *,
    title: str,
    image_path: Path,
    caption: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(12.0), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, size_pt=32, color=THEME.title, bold=True)

    # Image centered, large
    slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.6), width=Inches(11.333))

    if caption:
        c = slide.shapes.add_textbox(Inches(1.0), Inches(7.02), Inches(11.333), Inches(0.4))
        ctf = c.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = caption
        _set_run_font(cr, size_pt=14, color=THEME.muted)


def add_two_col_bullets_image(
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    image_path: Path,
    image_caption: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(12.0), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, size_pt=30, color=THEME.title, bold=True)

    # Left bullets
    left = slide.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(5.6), Inches(5.7))
    ltf = left.text_frame
    ltf.clear()
    ltf.word_wrap = True
    for i, b in enumerate(bullets):
        para = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        para.text = b
        para.level = 0
        para.font.size = Pt(20)
        para.font.name = "Calibri"
        para.font.color.rgb = THEME.body

    # Right image
    slide.shapes.add_picture(str(image_path), Inches(6.8), Inches(1.75), width=Inches(6.1))

    if image_caption:
        c = slide.shapes.add_textbox(Inches(6.8), Inches(7.02), Inches(6.1), Inches(0.35))
        ctf = c.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = image_caption
        _set_run_font(cr, size_pt=12, color=THEME.muted)


def build_investor_deck(pngs: dict[str, Path]) -> Path:
    prs = make_presentation()
    add_title_slide(
        prs,
        top_title="MAP Day Update",
        line1="Carbon Intelligence",
        line2="Investor Pitch",
        footer="Confidential",
    )

    add_section_slide(prs, title="The Story", subtitle="Data → Action → Finance for MSMEs and Banks")

    add_bullets_slide(
        prs,
        title="One-liner",
        bullets=[
            "Carbon Intelligence helps MSMEs measure emissions, reduce costs, and prove progress—then uses that verified data to unlock better-priced green finance for banks.",
        ],
    )

    add_two_col_bullets_image(
        prs,
        title="The problem (why now)",
        bullets=[
            "MSMEs face rising sustainability requests from buyers and lenders.",
            "Operational data is fragmented (emails, folders, spreadsheets).",
            "Banks have an ESG data gap: inconsistent SME data, expensive verification.",
            "Result: compliance is painful, decarbonization is slow, finance is mispriced.",
        ],
        image_path=pngs["platform-loop"],
        image_caption="We compound: measure → improve → prove → finance",
    )

    add_image_slide(
        prs,
        title="Product overview (web experience)",
        image_path=pngs["product-dashboard-web"],
        caption="Card-based dashboard: score, quick stats, trends, and actions",
    )

    add_section_slide(prs, title="Product Proof", subtitle="What exists in the UI today")

    add_image_slide(prs, title="Guided carbon assessment", image_path=pngs["product-carbon-assessment"])
    add_image_slide(prs, title="Document intelligence (evidence layer)", image_path=pngs["product-document-management"])
    add_image_slide(prs, title="Reporting (shareable outcomes)", image_path=pngs["product-reporting"])
    add_image_slide(prs, title="Green loans (conversion into finance)", image_path=pngs["product-green-loans"])

    add_section_slide(prs, title="Go-to-market + Model", subtitle="Wedge now, expand into finance workflows")

    add_bullets_slide(
        prs,
        title="Market wedge and expansion",
        bullets=[
            "Wedge: MSME carbon assessment + reporting + ROI-based recommendations (fast time-to-value).",
            "Expand: evidence/audit packs, supplier requests, benchmarking, continuous monitoring.",
            "Integrate: bank/fintech green products (eligibility, pricing, monitoring).",
        ],
    )

    add_bullets_slide(
        prs,
        title="Business model",
        bullets=[
            "SaaS (MSME): subscription per company/site, tiered by features.",
            "Usage: per report export, per processed document, per audit pack.",
            "Bank revenue: platform fee for origination enablement + monitoring (or referral/origination economics).",
            "Future: services marketplace + carbon credit rails.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Defensibility",
        bullets=[
            "Workflow lock-in: once reporting + documents + recommendations are embedded, switching costs rise.",
            "Data moat: labeled activity + document data improves benchmarks and underwriting signals over time.",
            "Distribution: banks and anchor enterprises can onboard MSMEs via portfolio/supply-chain programs.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Milestones to prove next (fundraising-ready)",
        bullets=[
            "Repeatable acquisition via bank or supply-chain onboarding channel.",
            "Retention/expansion: documents processed, reports exported, actions implemented.",
            "Outcome metrics: ₹ cost savings and kg CO₂ reduction per MSME cohort.",
            "Bank proof: faster credit files and improved portfolio coverage/monitoring.",
        ],
    )

    add_bullets_slide(
        prs,
        title="The ask",
        bullets=[
            "Capital to accelerate distribution partnerships and product depth in evidence + green lending workflows.",
            "Warm intros to banks/anchor enterprises to run portfolio or supply-chain pilots.",
        ],
        note="Replace this slide with your specific raise amount, runway, and milestones.",
    )

    out = OUT_DIR / "Carbon_Intelligence_Investor_Pitch.pptx"
    prs.save(out)
    return out


def build_msme_deck(pngs: dict[str, Path]) -> Path:
    prs = make_presentation()
    add_title_slide(
        prs,
        top_title="MAP Day Update",
        line1="Carbon Intelligence",
        line2="Pitch for MSMEs",
        footer="Customer deck",
    )

    add_image_slide(
        prs,
        title="The promise",
        image_path=pngs["product-dashboard-web"],
        caption="Measure footprint, cut costs, and generate reports—without spreadsheets",
    )

    add_bullets_slide(
        prs,
        title="What you get (clear benefits)",
        bullets=[
            "Win more customers: share sustainability reports + evidence when asked.",
            "Reduce waste: recommendations prioritized by impact, savings, and payback.",
            "Get audit-ready faster: organize bills/invoices with a clean evidence trail.",
            "Access green finance: convert performance into eligibility and better pricing.",
        ],
    )

    add_section_slide(prs, title="How it works", subtitle="A simple 4-step flow")

    add_two_col_bullets_image(
        prs,
        title="1) Guided carbon assessment",
        bullets=[
            "Answer a short step-by-step questionnaire across energy, transport, materials, water/waste, and processes.",
            "Receive total CO₂, category split, scopes (1/2/3), and a carbon score.",
        ],
        image_path=pngs["product-carbon-assessment"],
    )

    add_two_col_bullets_image(
        prs,
        title="2) Recommendations with ROI",
        bullets=[
            "Not generic tips: each action includes priority, estimated savings (₹), payback (months), and CO₂ reduction.",
            "Pick quick wins first, plan upgrades next (solar, efficiency, process improvements).",
        ],
        image_path=pngs["product-reporting"],
        image_caption="Track actions and results in reporting",
    )

    add_two_col_bullets_image(
        prs,
        title="3) Keep documents report-ready",
        bullets=[
            "Upload bills/invoices once; auto-extract key fields and compute carbon signals.",
            "Detect duplicates to reduce errors and create an audit-friendly evidence pack.",
            "Stop hunting for files when customers or auditors request proof.",
        ],
        image_path=pngs["product-document-management"],
    )

    add_two_col_bullets_image(
        prs,
        title="4) Generate reports you can share",
        bullets=[
            "Export reports for customers, internal reviews, tenders, and compliance.",
            "Show progress over time: footprint trend + actions implemented.",
        ],
        image_path=pngs["product-reporting"],
    )

    add_section_slide(prs, title="Green finance", subtitle="Convert sustainability into capital")

    add_two_col_bullets_image(
        prs,
        title="Green loan enablement",
        bullets=[
            "Plan a reduction project (solar, machinery upgrade, waste systems).",
            "Check eligibility up front and see transparent terms.",
            "Submit a cleaner application using verified activity + plan data.",
        ],
        image_path=pngs["product-green-loans"],
    )

    add_bullets_slide(
        prs,
        title="MSME ROI framing",
        bullets=[
            "Time saved: less manual reporting and document hunting.",
            "Cost savings: energy/fuel/process improvements with clear payback.",
            "Revenue enablement: meet buyer requirements and improve trust.",
            "Better financing access: qualify for green loans and incentives.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Next steps",
        bullets=[
            "Pick your sector and baseline period (last 3–12 months).",
            "Upload your latest bills/invoices and complete the assessment.",
            "Review prioritized actions and export your first customer-ready report.",
        ],
    )

    out = OUT_DIR / "Carbon_Intelligence_MSME_Pitch.pptx"
    prs.save(out)
    return out


def build_bank_deck(pngs: dict[str, Path]) -> Path:
    prs = make_presentation()
    add_title_slide(
        prs,
        top_title="MAP Day Update",
        line1="Carbon Intelligence",
        line2="Pitch for Banks",
        footer="Customer deck",
    )

    add_image_slide(
        prs,
        title="The promise",
        image_path=pngs["platform-loop"],
        caption="Turn SME sustainability data into underwriting + reporting advantage",
    )

    add_bullets_slide(
        prs,
        title="Why banks use this (value drivers)",
        bullets=[
            "Grow green lending efficiently: guided onboarding + evidence pack reduces friction.",
            "Better risk signals: score, emissions trajectory, action adoption, document-backed activity data.",
            "Portfolio visibility: standardized outputs across MSMEs for dashboards and disclosures.",
            "Lower cost-to-serve: fewer document chases, fewer data quality issues, faster files.",
        ],
    )

    add_section_slide(prs, title="Bank-ready outputs", subtitle="Standardized, comparable, defensible")

    add_two_col_bullets_image(
        prs,
        title="1) Verified evidence layer (documents)",
        bullets=[
            "MSMEs upload bills/invoices; we structure fields and maintain an audit trail.",
            "Duplicate detection improves cleanliness and reduces fraud/over-claim risk.",
            "Creates a reusable sustainability evidence pack per borrower.",
        ],
        image_path=pngs["product-document-management"],
    )

    add_two_col_bullets_image(
        prs,
        title="2) Carbon + ESG outputs (comparable across customers)",
        bullets=[
            "Total CO₂ + category breakdown; scopes 1/2/3 split.",
            "Carbon score + trend over time; recommended actions with payback.",
            "Consistent schema enables portfolio rollups and benchmarking.",
        ],
        image_path=pngs["product-reporting"],
    )

    add_two_col_bullets_image(
        prs,
        title="3) Green loan enablement (origination workflow)",
        bullets=[
            "Eligibility checks + transparent term preview based on policy rules and carbon score.",
            "Cleaner application files reduce back-and-forth and speed approvals.",
            "Ongoing monitoring ties disbursement impact to measurable reductions.",
        ],
        image_path=pngs["product-green-loans"],
    )

    add_section_slide(prs, title="Deployment models", subtitle="3 common ways banks roll this out")

    add_bullets_slide(
        prs,
        title="Rollout options",
        bullets=[
            "Model A — Bank-led onboarding: offer to MSME customers to prepare green-loan-ready profiles.",
            "Model B — Co-lending / partner channel: embed into partner program and share standardized outputs.",
            "Model C — Portfolio monitoring: enroll existing borrowers for ongoing capture and early-warning signals.",
        ],
    )

    add_bullets_slide(
        prs,
        title="KPIs this can improve",
        bullets=[
            "Acquisition: more qualified green loan leads per RM/branch.",
            "Approval speed: fewer document gaps; cleaner eligibility checks.",
            "Portfolio coverage: % of MSME book with standardized sustainability data.",
            "Monitoring: emissions trend + action adoption as engagement/risk signals.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Pilot proposal (suggested)",
        bullets=[
            "Select 50–200 MSMEs across 2–3 sectors; onboard via RM-led or digital journey.",
            "Collect baseline assessment + documents; generate standardized outputs and evidence packs.",
            "Measure time-to-file improvement, coverage, and eligibility conversion for green products.",
        ],
        note="This slide is designed to be customized with your bank’s target segment and product policy rules.",
    )

    out = OUT_DIR / "Carbon_Intelligence_Bank_Pitch.pptx"
    prs.save(out)
    return out


def build_partners_channels_deck(pngs: dict[str, Path]) -> Path:
    prs = make_presentation()
    add_title_slide(
        prs,
        top_title="MAP Day Update",
        line1="Carbon Intelligence",
        line2="Key Employees / Partners / Channels",
        footer="Internal + partner deck",
    )

    add_section_slide(prs, title="Why this exists", subtitle="Build the data → action → finance loop at scale")

    add_two_col_bullets_image(
        prs,
        title="The mission",
        bullets=[
            "Help MSMEs measure emissions, reduce costs, and prove progress with audit-ready evidence.",
            "Turn verified sustainability data into a financial primitive for green lending and portfolio reporting.",
            "Make sustainability adoption practical: ROI-first recommendations, simple workflows, fast time-to-value.",
        ],
        image_path=pngs["platform-loop"],
        image_caption="Our compounding loop: measure → improve → prove → finance",
    )

    add_image_slide(
        prs,
        title="What we’ve built (product snapshot)",
        image_path=pngs["product-dashboard-web"],
        caption="Web dashboard: score, trends, documents, reporting, and green finance workflows",
    )

    add_section_slide(prs, title="Who this deck is for", subtitle="Three audiences; one coherent program")

    add_bullets_slide(
        prs,
        title="Audience and outcomes",
        bullets=[
            "Key Employees: join the team and own core product + distribution outcomes.",
            "Partners: embed/extend the platform (consultancies, auditors, implementers, industry bodies).",
            "Channels: drive repeatable MSME onboarding via banks, anchor enterprises, or ecosystem partners.",
        ],
    )

    add_section_slide(prs, title="Where we win", subtitle="Reliable evidence + simple UX + finance conversion")

    add_two_col_bullets_image(
        prs,
        title="Evidence layer (documents) = defensibility",
        bullets=[
            "Bills/invoices become structured data with audit trail.",
            "Duplicate detection reduces errors and strengthens trust.",
            "Creates repeatable ‘evidence packs’ for buyers, auditors, and lenders.",
        ],
        image_path=pngs["product-document-management"],
    )

    add_two_col_bullets_image(
        prs,
        title="Reporting that stakeholders actually consume",
        bullets=[
            "Standardized outputs: total CO₂, category split, scope 1/2/3, score and trend.",
            "Exports for buyers, internal teams, and bank/portfolio reporting.",
            "Action adoption tracked over time (progress, payback, impact).",
        ],
        image_path=pngs["product-reporting"],
    )

    add_two_col_bullets_image(
        prs,
        title="Finance conversion (green lending workflow)",
        bullets=[
            "Eligibility and term preview driven by policy rules + verified data.",
            "Cleaner files reduce back-and-forth and speed approvals.",
            "Monitoring ties capital deployment to measurable reduction outcomes.",
        ],
        image_path=pngs["product-green-loans"],
    )

    add_section_slide(prs, title="Key employees", subtitle="What we’re hiring for (and why it’s exciting)")

    add_bullets_slide(
        prs,
        title="Roles we typically need (adapt as needed)",
        bullets=[
            "Product + Engineering: workflows, evidence layer, integrations, reporting/export.",
            "Data/ML: extraction quality, dedup, carbon estimation models, benchmarks.",
            "Growth + Partnerships: bank/anchor programs, onboarding playbooks, channel enablement.",
            "Customer Success: repeatable onboarding, retention/expansion, outcome reporting (₹ savings + CO₂).",
        ],
        note="Replace this slide with your exact open roles and seniority.",
    )

    add_bullets_slide(
        prs,
        title="Why a key employee should join",
        bullets=[
            "Mission + measurable outcomes: help MSMEs cut costs while meeting compliance pressure.",
            "Real product in place: ship on a working workflow system (not a concept).",
            "Distribution leverage: banks and anchor enterprises can onboard thousands via programs.",
            "Compounding advantage: better evidence → better data → better decisions → better finance conversion.",
        ],
    )

    add_section_slide(prs, title="Partners", subtitle="Who we work with and what they gain")

    add_bullets_slide(
        prs,
        title="Partner types (examples)",
        bullets=[
            "Sustainability consultants / implementers (baseline → roadmap → execution).",
            "Auditors / verifiers (evidence packs + standardized outputs).",
            "Industry associations (member onboarding programs).",
            "Solution providers (solar/efficiency/waste) who need qualified demand + proof.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Partner value proposition",
        bullets=[
            "Faster delivery: standardized workflows and reporting reduce manual effort.",
            "Higher trust: document-backed evidence improves auditability and buyer/lender confidence.",
            "More revenue: partners can bundle services around recommendations and implementation.",
            "Clear packaging: tiered engagement templates (assessment → action plan → monitoring).",
        ],
    )

    add_section_slide(prs, title="Channels", subtitle="Repeatable acquisition and onboarding")

    add_bullets_slide(
        prs,
        title="Channel models (how leads come in)",
        bullets=[
            "Bank channel: RMs / digital journeys for green lending and portfolio monitoring.",
            "Anchor enterprise channel: supply-chain compliance programs (vendor onboarding).",
            "Ecosystem channel: associations, platforms, and service providers referring MSMEs.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Channel motion (recommended)",
        bullets=[
            "Step 1: MSME invite → lightweight onboarding + baseline assessment.",
            "Step 2: Document upload → evidence pack + data cleanliness checks.",
            "Step 3: Report export → buyer/bank-ready output (shareable).",
            "Step 4: Action plan → track ROI + implementation status.",
            "Step 5: Finance conversion (optional) → eligibility + application workflow.",
        ],
    )

    add_bullets_slide(
        prs,
        title="What channels need from us (enablement kit)",
        bullets=[
            "Co-branded landing + onboarding journey.",
            "Sector-specific templates (inputs, recommended actions, ROI assumptions).",
            "Partner portal: cohort tracking, export packs, and progress dashboards.",
            "Training: talk track, demo script, objection handling, and ROI calculator.",
        ],
    )

    add_section_slide(prs, title="How we measure success", subtitle="Simple KPIs per audience")

    add_bullets_slide(
        prs,
        title="KPIs",
        bullets=[
            "Acquisition: invited → activated MSMEs; CAC by channel; time-to-first-report.",
            "Engagement: documents processed; exports; action adoption rate.",
            "Outcomes: ₹ savings realized; kg CO₂ reduced; score improvement.",
            "Finance: eligibility rate; application completion; approval speed; portfolio coverage.",
        ],
    )

    add_bullets_slide(
        prs,
        title="Next steps (choose your path)",
        bullets=[
            "Key employees: share your role fit + a 30/60/90 plan you’d execute.",
            "Partners: propose 1–2 packaged offerings you can deliver using the platform.",
            "Channels: select a pilot cohort size and define the onboarding + reporting workflow.",
        ],
        note="This deck is designed to be customized per audience and partner type.",
    )

    out = OUT_DIR / "Carbon_Intelligence_Key_Employees_Partners_Channels_Pitch.pptx"
    prs.save(out)
    return out


def main() -> None:
    ensure_dirs()
    pngs = convert_all_svgs()

    outputs = [
        build_investor_deck(pngs),
        build_msme_deck(pngs),
        build_bank_deck(pngs),
        build_partners_channels_deck(pngs),
    ]

    print("Generated:")
    for p in outputs:
        print(f"- {p}")


if __name__ == "__main__":
    # Avoid cairo DLL issues by ensuring we only run in a normal env.
    os.environ.setdefault("CAIROSVG_LOGLEVEL", "ERROR")
    main()

