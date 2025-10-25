#!/usr/bin/env python3
"""
Final refined 5-slide pitch presentation about mist solutions for air and water borne infections
Designed for IdeaOne Hackathon 2024 - Maximum visual impact and professional presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_FILL
import os

def create_final_mist_solutions_pitch():
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define professional color scheme
    primary_blue = RGBColor(0, 102, 153)      # Deep blue
    secondary_green = RGBColor(0, 153, 76)    # Health green
    accent_orange = RGBColor(255, 102, 0)     # Action orange
    text_dark = RGBColor(44, 62, 80)          # Dark text
    background_light = RGBColor(248, 249, 250) # Light background
    success_green = RGBColor(40, 167, 69)     # Success green
    warning_red = RGBColor(220, 53, 69)       # Warning red
    light_blue = RGBColor(173, 216, 230)      # Light blue
    light_green = RGBColor(144, 238, 144)     # Light green
    white = RGBColor(255, 255, 255)           # White
    
    # Slide 1: The Problem - Maximum visual impact
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background with gradient effect
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_light
    
    # Top accent bar with gradient
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.4))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = warning_red
    
    # Central problem icon
    problem_circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.5), Inches(2.5), Inches(2.5))
    problem_circle.fill.solid()
    problem_circle.fill.fore_color.rgb = warning_red
    problem_circle.line.fill.background()
    
    # Problem icon text
    problem_text = slide1.shapes.add_textbox(Inches(5.6), Inches(2.3), Inches(2.3), Inches(0.9))
    problem_frame = problem_text.text_frame
    problem_p = problem_frame.paragraphs[0]
    problem_p.text = "⚠️\nCRISIS"
    problem_p.font.size = Pt(24)
    problem_p.font.bold = True
    problem_p.font.color.rgb = white
    problem_p.alignment = PP_ALIGN.CENTER
    
    # Main title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "THE GLOBAL HEALTH CRISIS"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = warning_red
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Air & Water Borne Infections Threaten Global Health"
    subtitle_p.font.size = Pt(26)
    subtitle_p.font.color.rgb = text_dark
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Statistics with enhanced visual boxes
    # Stat 1
    stat1_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(3), Inches(1.8))
    stat1_box.fill.solid()
    stat1_box.fill.fore_color.rgb = light_blue
    stat1_box.line.color.rgb = primary_blue
    stat1_box.line.width = Pt(3)
    
    stat1_text = slide1.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(2.8), Inches(1.2))
    stat1_frame = stat1_text.text_frame
    stat1_p = stat1_frame.paragraphs[0]
    stat1_p.text = "3.2M\nDEATHS\nANNUALLY"
    stat1_p.font.size = Pt(20)
    stat1_p.font.bold = True
    stat1_p.font.color.rgb = primary_blue
    stat1_p.alignment = PP_ALIGN.CENTER
    
    # Stat 2
    stat2_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(4.2), Inches(3), Inches(1.8))
    stat2_box.fill.solid()
    stat2_box.fill.fore_color.rgb = light_green
    stat2_box.line.color.rgb = secondary_green
    stat2_box.line.width = Pt(3)
    
    stat2_text = slide1.shapes.add_textbox(Inches(3.9), Inches(4.5), Inches(2.8), Inches(1.2))
    stat2_frame = stat2_text.text_frame
    stat2_p = stat2_frame.paragraphs[0]
    stat2_p.text = "$1.2T\nHEALTHCARE\nCOSTS"
    stat2_p.font.size = Pt(20)
    stat2_p.font.bold = True
    stat2_p.font.color.rgb = secondary_green
    stat2_p.alignment = PP_ALIGN.CENTER
    
    # Stat 3
    stat3_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(4.2), Inches(3), Inches(1.8))
    stat3_box.fill.solid()
    stat3_box.fill.fore_color.rgb = RGBColor(255, 218, 185)  # Light orange
    stat3_box.line.color.rgb = accent_orange
    stat3_box.line.width = Pt(3)
    
    stat3_text = slide1.shapes.add_textbox(Inches(7.2), Inches(4.5), Inches(2.8), Inches(1.2))
    stat3_frame = stat3_text.text_frame
    stat3_p = stat3_frame.paragraphs[0]
    stat3_p.text = "90%\nGLOBAL\nPOPULATION"
    stat3_p.font.size = Pt(20)
    stat3_p.font.bold = True
    stat3_p.font.color.rgb = accent_orange
    stat3_p.alignment = PP_ALIGN.CENTER
    
    # Stat 4
    stat4_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(4.2), Inches(2.5), Inches(1.8))
    stat4_box.fill.solid()
    stat4_box.fill.fore_color.rgb = RGBColor(255, 182, 193)  # Light pink
    stat4_box.line.color.rgb = warning_red
    stat4_box.line.width = Pt(3)
    
    stat4_text = slide1.shapes.add_textbox(Inches(10.5), Inches(4.5), Inches(2.3), Inches(1.2))
    stat4_frame = stat4_text.text_frame
    stat4_p = stat4_frame.paragraphs[0]
    stat4_p.text = "15-20%\nLIVESTOCK\nMORTALITY"
    stat4_p.font.size = Pt(18)
    stat4_p.font.bold = True
    stat4_p.font.color.rgb = warning_red
    stat4_p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "IdeaOne Hackathon 2024 | Healthcare Innovation Challenge"
    footer_p.font.size = Pt(16)
    footer_p.font.color.rgb = primary_blue
    footer_p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Solution Overview - Enhanced with visual elements
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = background_light
    
    # Top accent bar
    accent_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.4))
    accent_bar2.fill.solid()
    accent_bar2.fill.fore_color.rgb = secondary_green
    
    # Title
    title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(1))
    title2_frame = title2_box.text_frame
    title2_p = title2_frame.paragraphs[0]
    title2_p.text = "OUR REVOLUTIONARY SOLUTION"
    title2_p.font.size = Pt(38)
    title2_p.font.bold = True
    title2_p.font.color.rgb = primary_blue
    title2_p.alignment = PP_ALIGN.CENTER
    
    # Solution overview
    solution_box = slide2.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(0.8))
    solution_frame = solution_box.text_frame
    solution_p = solution_frame.paragraphs[0]
    solution_p.text = "Advanced Mist Technology for Air & Water Disinfection"
    solution_p.font.size = Pt(26)
    solution_p.font.color.rgb = text_dark
    solution_p.alignment = PP_ALIGN.CENTER
    
    # Technology icon
    tech_circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(3.2), Inches(2.5), Inches(2.5))
    tech_circle.fill.solid()
    tech_circle.fill.fore_color.rgb = secondary_green
    tech_circle.line.fill.background()
    
    tech_text = slide2.shapes.add_textbox(Inches(5.6), Inches(4), Inches(2.3), Inches(0.9))
    tech_frame = tech_text.text_frame
    tech_p = tech_frame.paragraphs[0]
    tech_p.text = "MIST\nTECH"
    tech_p.font.size = Pt(24)
    tech_p.font.bold = True
    tech_p.font.color.rgb = white
    tech_p.alignment = PP_ALIGN.CENTER
    
    # Key features in enhanced visual boxes
    # Left feature box
    left_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6), Inches(6), Inches(1.2))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = light_blue
    left_box.line.color.rgb = primary_blue
    left_box.line.width = Pt(3)
    
    left_text = slide2.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(5.8), Inches(0.8))
    left_frame = left_text.text_frame
    left_p = left_frame.paragraphs[0]
    left_p.text = "🔬 99.99% Pathogen Elimination | ⚡ IoT Smart Control | 🌍 Dual Air & Water Applications"
    left_p.font.size = Pt(16)
    left_p.font.bold = True
    left_p.font.color.rgb = primary_blue
    left_p.alignment = PP_ALIGN.CENTER
    
    # Right feature box
    right_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(6), Inches(6), Inches(1.2))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = light_green
    right_box.line.color.rgb = secondary_green
    right_box.line.width = Pt(3)
    
    right_text = slide2.shapes.add_textbox(Inches(6.6), Inches(6.2), Inches(5.8), Inches(0.8))
    right_frame = right_text.text_frame
    right_p = right_frame.paragraphs[0]
    right_p.text = "💰 60% Cost Reduction | 📈 85% Infection Decrease | 🏥 Zero Health Effects"
    right_p.font.size = Pt(16)
    right_p.font.bold = True
    right_p.font.color.rgb = secondary_green
    right_p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Technology Deep Dive - Enhanced with visual elements
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background3 = slide3.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = background_light
    
    # Top accent bar
    accent_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.4))
    accent_bar3.fill.solid()
    accent_bar3.fill.fore_color.rgb = accent_orange
    
    # Title
    title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(1))
    title3_frame = title3_box.text_frame
    title3_p = title3_frame.paragraphs[0]
    title3_p.text = "BREAKTHROUGH TECHNOLOGY"
    title3_p.font.size = Pt(38)
    title3_p.font.bold = True
    title3_p.font.color.rgb = primary_blue
    title3_p.alignment = PP_ALIGN.CENTER
    
    # Air disinfection section with enhanced visual box
    air_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(6), Inches(3.8))
    air_box.fill.solid()
    air_box.fill.fore_color.rgb = light_blue
    air_box.line.color.rgb = primary_blue
    air_box.line.width = Pt(4)
    
    air_title = slide3.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(5.8), Inches(0.6))
    air_title_frame = air_title.text_frame
    air_title_p = air_title_frame.paragraphs[0]
    air_title_p.text = "🌬️ AIR DISINFECTION"
    air_title_p.font.size = Pt(22)
    air_title_p.font.bold = True
    air_title_p.font.color.rgb = primary_blue
    air_title_p.alignment = PP_ALIGN.CENTER
    
    air_content = slide3.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(5.8), Inches(2.7))
    air_content_frame = air_content.text_frame
    air_content_p = air_content_frame.paragraphs[0]
    air_content_text = """• Droplet Size: 1-5 micrometers
• Coverage: 100-1000 m² per unit
• Flow Rate: 50-500 L/min
• Power: 50-200W per unit
• Agents: Hypochlorous acid, H₂O₂

Mechanism:
✓ Mist captures airborne pathogens
✓ Antimicrobial agents neutralize threats
✓ Continuous air circulation
✓ Real-time quality monitoring"""
    air_content_p.text = air_content_text
    air_content_p.font.size = Pt(14)
    air_content_p.font.color.rgb = text_dark
    
    # Water disinfection section with enhanced visual box
    water_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.2), Inches(6), Inches(3.8))
    water_box.fill.solid()
    water_box.fill.fore_color.rgb = light_green
    water_box.line.color.rgb = secondary_green
    water_box.line.width = Pt(4)
    
    water_title = slide3.shapes.add_textbox(Inches(6.6), Inches(2.4), Inches(5.8), Inches(0.6))
    water_title_frame = water_title.text_frame
    water_title_p = water_title_frame.paragraphs[0]
    water_title_p.text = "💧 WATER DISINFECTION"
    water_title_p.font.size = Pt(22)
    water_title_p.font.bold = True
    water_title_p.font.color.rgb = secondary_green
    water_title_p.alignment = PP_ALIGN.CENTER
    
    water_content = slide3.shapes.add_textbox(Inches(6.6), Inches(3.1), Inches(5.8), Inches(2.7))
    water_content_frame = water_content.text_frame
    water_content_p = water_content_frame.paragraphs[0]
    water_content_text = """• Treatment Rate: 100-1000 L/hour
• Efficiency: 99.99% pathogen elimination
• Residual: 0.1-0.5 ppm safe levels
• Power: 100-500W
• Maintenance: 6-month intervals

Process:
✓ Pre-filtration removes particles
✓ Ultrasonic mist generation
✓ UV-C irradiation during formation
✓ Antimicrobial agent injection
✓ Post-treatment filtration"""
    water_content_p.text = water_content_text
    water_content_p.font.size = Pt(14)
    water_content_p.font.color.rgb = text_dark
    
    # Innovation highlights with special styling
    innovation_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6.2), Inches(11.33), Inches(1))
    innovation_box.fill.solid()
    innovation_box.fill.fore_color.rgb = RGBColor(255, 248, 220)  # Light yellow
    innovation_box.line.color.rgb = accent_orange
    innovation_box.line.width = Pt(4)
    
    innovation_text = slide3.shapes.add_textbox(Inches(1.1), Inches(6.4), Inches(11.13), Inches(0.6))
    innovation_frame = innovation_text.text_frame
    innovation_p = innovation_frame.paragraphs[0]
    innovation_p.text = "🚀 INNOVATION: Smart IoT Control | Variable Frequency Transducers | Zero Ozone Production | Seamless Integration"
    innovation_p.font.size = Pt(18)
    innovation_p.font.bold = True
    innovation_p.font.color.rgb = accent_orange
    innovation_p.alignment = PP_ALIGN.CENTER
    
    # Slide 4: Market Impact & Applications - Enhanced with visual elements
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background4 = slide4.background
    fill4 = background4.fill
    fill4.solid()
    fill4.fore_color.rgb = background_light
    
    # Top accent bar
    accent_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.4))
    accent_bar4.fill.solid()
    accent_bar4.fill.fore_color.rgb = success_green
    
    # Title
    title4_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(1))
    title4_frame = title4_box.text_frame
    title4_p = title4_frame.paragraphs[0]
    title4_p.text = "MARKET IMPACT & APPLICATIONS"
    title4_p.font.size = Pt(38)
    title4_p.font.bold = True
    title4_p.font.color.rgb = primary_blue
    title4_p.alignment = PP_ALIGN.CENTER
    
    # Market size with special styling
    market_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(1.4))
    market_box.fill.solid()
    market_box.fill.fore_color.rgb = RGBColor(255, 218, 185)  # Light orange
    market_box.line.color.rgb = accent_orange
    market_box.line.width = Pt(5)
    
    market_text = slide4.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.13), Inches(0.8))
    market_frame = market_text.text_frame
    market_p = market_frame.paragraphs[0]
    market_p.text = "💰 TOTAL ADDRESSABLE MARKET: $32.4 BILLION"
    market_p.font.size = Pt(30)
    market_p.font.bold = True
    market_p.font.color.rgb = accent_orange
    market_p.alignment = PP_ALIGN.CENTER
    
    # Applications grid with enhanced visual boxes
    # Healthcare sector
    health_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(6), Inches(2.8))
    health_box.fill.solid()
    health_box.fill.fore_color.rgb = light_blue
    health_box.line.color.rgb = primary_blue
    health_box.line.width = Pt(3)
    
    health_title = slide4.shapes.add_textbox(Inches(0.6), Inches(4), Inches(5.8), Inches(0.5))
    health_title_frame = health_title.text_frame
    health_title_p = health_title_frame.paragraphs[0]
    health_title_p.text = "🏥 HEALTHCARE SECTOR"
    health_title_p.font.size = Pt(20)
    health_title_p.font.bold = True
    health_title_p.font.color.rgb = primary_blue
    health_title_p.alignment = PP_ALIGN.CENTER
    
    health_content = slide4.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(5.8), Inches(1.8))
    health_content_frame = health_content.text_frame
    health_content_p = health_content_frame.paragraphs[0]
    health_content_text = """• Hospitals & clinics (2.5M facilities)
• Operating rooms & ICUs
• Emergency departments
• Home healthcare
• 85% reduction in infections
• 70% decrease in illness rates"""
    health_content_p.text = health_content_text
    health_content_p.font.size = Pt(14)
    health_content_p.font.color.rgb = text_dark
    
    # Agricultural sector
    agri_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3.8), Inches(6), Inches(2.8))
    agri_box.fill.solid()
    agri_box.fill.fore_color.rgb = light_green
    agri_box.line.color.rgb = secondary_green
    agri_box.line.width = Pt(3)
    
    agri_title = slide4.shapes.add_textbox(Inches(6.6), Inches(4), Inches(5.8), Inches(0.5))
    agri_title_frame = agri_title.text_frame
    agri_title_p = agri_title_frame.paragraphs[0]
    agri_title_p.text = "🏭 AGRICULTURAL SECTOR"
    agri_title_p.font.size = Pt(20)
    agri_title_p.font.bold = True
    agri_title_p.font.color.rgb = secondary_green
    agri_title_p.alignment = PP_ALIGN.CENTER
    
    agri_content = slide4.shapes.add_textbox(Inches(6.6), Inches(4.6), Inches(5.8), Inches(1.8))
    agri_content_frame = agri_content.text_frame
    agri_content_p = agri_content_frame.paragraphs[0]
    agri_content_text = """• Livestock operations (570M farms)
• Poultry & dairy farms
• Aquaculture facilities
• Veterinary clinics
• 40% reduction in costs
• 60% decrease in mortality"""
    agri_content_p.text = agri_content_text
    agri_content_p.font.size = Pt(14)
    agri_content_p.font.color.rgb = text_dark
    
    # Key metrics box
    metrics_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6.8), Inches(11.33), Inches(0.6))
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Very light blue
    metrics_box.line.color.rgb = primary_blue
    metrics_box.line.width = Pt(3)
    
    metrics_text = slide4.shapes.add_textbox(Inches(1.1), Inches(6.9), Inches(11.13), Inches(0.4))
    metrics_frame = metrics_text.text_frame
    metrics_p = metrics_frame.paragraphs[0]
    metrics_p.text = "KEY METRICS: 25% increase in production efficiency | 60% decrease in animal mortality | 70% decrease in respiratory illness"
    metrics_p.font.size = Pt(16)
    metrics_p.font.bold = True
    metrics_p.font.color.rgb = primary_blue
    metrics_p.alignment = PP_ALIGN.CENTER
    
    # Slide 5: Call to Action & Next Steps - Enhanced with visual elements
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background5 = slide5.background
    fill5 = background5.fill
    fill5.solid()
    fill5.fore_color.rgb = background_light
    
    # Top accent bar
    accent_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.4))
    accent_bar5.fill.solid()
    accent_bar5.fill.fore_color.rgb = primary_blue
    
    # Title
    title5_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(1))
    title5_frame = title5_box.text_frame
    title5_p = title5_frame.paragraphs[0]
    title5_p.text = "JOIN THE REVOLUTION"
    title5_p.font.size = Pt(38)
    title5_p.font.bold = True
    title5_p.font.color.rgb = primary_blue
    title5_p.alignment = PP_ALIGN.CENTER
    
    # Value proposition
    value_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(1))
    value_box.fill.solid()
    value_box.fill.fore_color.rgb = RGBColor(255, 248, 220)  # Light yellow
    value_box.line.color.rgb = accent_orange
    value_box.line.width = Pt(4)
    
    value_text = slide5.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(11.13), Inches(0.6))
    value_frame = value_text.text_frame
    value_p = value_frame.paragraphs[0]
    value_p.text = "Transform Global Health Through Innovative Mist Technology"
    value_p.font.size = Pt(26)
    value_p.font.bold = True
    value_p.font.color.rgb = accent_orange
    value_p.alignment = PP_ALIGN.CENTER
    
    # Partnership opportunities
    partner_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.4), Inches(6), Inches(3.2))
    partner_box.fill.solid()
    partner_box.fill.fore_color.rgb = light_blue
    partner_box.line.color.rgb = primary_blue
    partner_box.line.width = Pt(3)
    
    partner_title = slide5.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(5.8), Inches(0.5))
    partner_title_frame = partner_title.text_frame
    partner_title_p = partner_title_frame.paragraphs[0]
    partner_title_p.text = "🤝 PARTNERSHIP OPPORTUNITIES"
    partner_title_p.font.size = Pt(20)
    partner_title_p.font.bold = True
    partner_title_p.font.color.rgb = primary_blue
    partner_title_p.alignment = PP_ALIGN.CENTER
    
    partner_content = slide5.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.2))
    partner_content_frame = partner_content.text_frame
    partner_content_p = partner_content_frame.paragraphs[0]
    partner_content_text = """• Pilot program participation
• Technology development partnership
• Investment & funding opportunities
• Distribution & sales partnerships
• Research & development collaboration

💡 WHY PARTNER WITH US
• Proven 99.99% effectiveness
• $32.4B market opportunity
• Strong competitive advantages
• Experienced team & advisors"""
    partner_content_p.text = partner_content_text
    partner_content_p.font.size = Pt(14)
    partner_content_p.font.color.rgb = text_dark
    
    # Next steps
    steps_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3.4), Inches(6), Inches(3.2))
    steps_box.fill.solid()
    steps_box.fill.fore_color.rgb = light_green
    steps_box.line.color.rgb = secondary_green
    steps_box.line.width = Pt(3)
    
    steps_title = slide5.shapes.add_textbox(Inches(6.6), Inches(3.6), Inches(5.8), Inches(0.5))
    steps_title_frame = steps_title.text_frame
    steps_title_p = steps_title_frame.paragraphs[0]
    steps_title_p.text = "📋 IMMEDIATE NEXT STEPS"
    steps_title_p.font.size = Pt(20)
    steps_title_p.font.bold = True
    steps_title_p.font.color.rgb = secondary_green
    steps_title_p.alignment = PP_ALIGN.CENTER
    
    steps_content = slide5.shapes.add_textbox(Inches(6.6), Inches(4.2), Inches(5.8), Inches(2.2))
    steps_content_frame = steps_content.text_frame
    steps_content_p = steps_content_frame.paragraphs[0]
    steps_content_text = """1. Schedule technology demonstration
2. Discuss pilot program opportunities
3. Explore partnership possibilities
4. Review investment opportunities
5. Join our advisory board

📞 CONTACT INFORMATION
Email: partnerships@misthealth.com
Phone: +1-555-MIST-HEALTH
Website: www.misthealth.com"""
    steps_content_p.text = steps_content_text
    steps_content_p.font.size = Pt(14)
    steps_content_p.font.color.rgb = text_dark
    
    # Impact statement with special styling
    impact_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6.8), Inches(11.33), Inches(0.6))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Very light green
    impact_box.line.color.rgb = success_green
    impact_box.line.width = Pt(4)
    
    impact_text = slide5.shapes.add_textbox(Inches(1.1), Inches(6.9), Inches(11.13), Inches(0.4))
    impact_frame = impact_text.text_frame
    impact_p = impact_frame.paragraphs[0]
    impact_p.text = "Together, we can prevent millions of respiratory infections, save billions in healthcare costs, and create a healthier future for all."
    impact_p.font.size = Pt(18)
    impact_p.font.bold = True
    impact_p.font.color.rgb = success_green
    impact_p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_path = "/workspace/mist_gen/final_mist_solutions_pitch.pptx"
    prs.save(output_path)
    print(f"Final PowerPoint presentation created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_final_mist_solutions_pitch()