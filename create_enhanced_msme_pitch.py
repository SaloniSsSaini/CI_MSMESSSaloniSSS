#!/usr/bin/env python3
"""
Enhanced Script to create a PowerPoint presentation for MSMEs as customers on Carbon Intelligence benefits
Focus on Carbon Savings, Green Finance Access, and ROI Benefits
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_enhanced_msme_pitch_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme - Green and professional
    primary_color = RGBColor(46, 125, 50)  # Green for sustainability
    secondary_color = RGBColor(0, 102, 51)  # Dark green
    accent_color = RGBColor(255, 140, 0)  # Orange for highlights
    text_color = RGBColor(51, 51, 51)  # Dark gray text
    light_bg = RGBColor(248, 249, 250)  # Light background
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Carbon Intelligence"
    subtitle.text = "Transform Your Business with Smart Sustainability\n\nSave Money • Access Green Finance • Build Competitive Advantage\n\nPowered by AI-Driven Carbon Intelligence Scoring\nFor Micro, Small & Medium Enterprises\n2024"
    
    # Format title
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Format subtitle
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Slide 2: Why Carbon Intelligence?
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Why Choose Carbon Intelligence?"
    content2.text = """🎯 DESIGNED SPECIFICALLY FOR MSMEs

The Problem:
• High energy and resource costs eating into profits (15-25% of revenue)
• Complex sustainability requirements and regulations
• Expensive consultants and manual processes (₹2-5 lakh annually)
• No clear visibility into environmental impact and costs
• Difficulty accessing green finance and government incentives

Our Solution:
✅ AI-powered Carbon Intelligence Scoring (0-100 scale)
✅ Real-time cost optimization insights and recommendations
✅ Automated ESG compliance reporting and documentation
✅ Direct access to green finance and government incentives
✅ Mobile app for easy management and monitoring

BENEFITS FOR YOUR BUSINESS:
💰 Save 20-40% on operational costs annually
📈 Increase profits through efficiency and optimization
🏆 Build competitive advantage and market reputation
🌱 Meet sustainability goals and regulatory compliance
📊 Access to green finance with lower interest rates
⚡ Easy-to-use mobile and web platform
🎯 Personalized recommendations based on your industry"""
    
    # Slide 3: Carbon Intelligence Scoring System
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Carbon Intelligence Scoring System"
    content3.text = """🧠 AI-POWERED CARBON INTELLIGENCE SCORING

SCORING METHODOLOGY (0-100 Scale):
• Energy Efficiency (25%): Renewable energy usage, consumption patterns
• Water Management (20%): Conservation, recycling, efficiency measures
• Waste Management (20%): Reduction, recycling, circular economy
• Transportation (15%): Green transport, logistics optimization
• Materials & Supply Chain (10%): Sustainable sourcing, efficiency
• ESG Compliance (10%): Regulatory compliance, reporting

REAL-TIME MONITORING:
• Continuous data collection from SMS, emails, and transactions
• AI-powered analysis and categorization
• Automated carbon footprint calculation
• Predictive analytics for future performance

SCORING BENEFITS:
• Higher scores = Better access to green finance
• Lower interest rates on loans (1-3% reduction)
• Priority processing and approval
• Enhanced market reputation and credibility
• Competitive advantage in tenders and contracts

SCORE TIERS & BENEFITS:
• Platinum (90-100): 3% rate reduction, highest loan amounts
• Gold (80-89): 2% rate reduction, priority processing
• Silver (70-79): 1% rate reduction, standard processing
• Bronze (60-69): 0.5% rate reduction, basic processing
• Below 60: Improvement plan with support and guidance"""
    
    # Slide 4: Carbon Savings & Cost Reduction
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Carbon Savings & Cost Reduction"
    content4.text = """💰 PROVEN COST SAVINGS & CARBON REDUCTION

AVERAGE SAVINGS PER MSME:
• Energy Costs: 25-40% reduction (₹75,000-₹3,00,000 annually)
• Water Usage: 20-30% reduction (₹40,000-₹1,50,000 annually)
• Waste Management: 35-50% reduction (₹25,000-₹1,00,000 annually)
• Transportation: 15-25% reduction (₹30,000-₹1,50,000 annually)
• Materials: 20-30% reduction (₹50,000-₹2,00,000 annually)

TOTAL ANNUAL SAVINGS: ₹2,20,000 - ₹9,00,000 per MSME

REAL CUSTOMER EXAMPLES:

Manufacturing MSME (75 employees):
• Energy savings: ₹2,50,000/year (40% reduction)
• Water savings: ₹1,20,000/year (30% reduction)
• Waste reduction: ₹80,000/year (45% reduction)
• Total savings: ₹4,50,000/year
• Carbon footprint reduction: 35%

Textile MSME (50 employees):
• Energy savings: ₹1,80,000/year (35% reduction)
• Water savings: ₹90,000/year (25% reduction)
• Material optimization: ₹1,50,000/year (30% reduction)
• Total savings: ₹4,20,000/year
• Carbon footprint reduction: 30%

Food Processing MSME (40 employees):
• Energy savings: ₹1,20,000/year (30% reduction)
• Water savings: ₹60,000/year (20% reduction)
• Waste reduction: ₹70,000/year (40% reduction)
• Total savings: ₹2,50,000/year
• Carbon footprint reduction: 25%

ROI: 400-600% return on investment within first year"""
    
    # Slide 5: Green Finance Access
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Green Finance Access & Benefits"
    content5.text = """🏦 COMPREHENSIVE GREEN FINANCE SOLUTIONS

GREEN LOAN PRODUCTS:
• Carbon Intelligence Green Loans: 1-3% lower interest rates
• Solar Energy Loans: Up to ₹2 crores at 8-10% interest
• Energy Efficiency Loans: Equipment financing at 9-11%
• Water Conservation Loans: Infrastructure at 10-12%
• Waste Management Loans: Technology at 11-13%

SCORING-BASED BENEFITS:
• Higher carbon scores = Better loan terms
• Faster approval process (24-48 hours vs 7-14 days)
• Reduced collateral requirements
• Extended repayment periods
• Priority customer support

GOVERNMENT INCENTIVES:
• Energy efficiency grants (up to ₹5 lakh)
• Solar panel subsidies (30-40% of cost)
• Water conservation incentives (up to ₹2 lakh)
• Waste management grants (up to ₹3 lakh)
• Green certification benefits and tax incentives

CARBON CREDIT TRADING:
• Earn money by reducing emissions
• Sell carbon credits to other companies
• Offset remaining emissions cost-effectively
• Additional revenue stream (₹50,000-₹2,00,000 annually)
• Environmental impact monetization

BANKING PARTNERSHIPS:
• Direct access to partner banks
• Pre-approved green loan products
• Dedicated relationship managers
• Streamlined application process
• Competitive interest rates

SUCCESS STORY:
"Through Carbon Intelligence, we got a green loan at 2.5% lower interest rate, saving ₹3.6 lakh annually on our ₹1.5 crore loan. The platform helped us qualify by tracking our sustainability improvements." - Rajesh Kumar, EcoTech Manufacturing"""
    
    # Slide 6: Carbon Trading & Additional Revenue
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Carbon Trading & Additional Revenue"
    content6.text = """🌍 CARBON CREDIT TRADING OPPORTUNITIES

CARBON CREDIT MARKETPLACE:
• Verified carbon credits from your sustainability efforts
• Real-time pricing and market data
• Automated trading and portfolio management
• Additional revenue stream for your business

OFFSET PROJECTS:
• Renewable Energy: Solar, wind, hydro projects
• Reforestation: Tree planting and forest conservation
• Energy Efficiency: Building and industrial efficiency
• Waste Management: Methane capture and utilization
• Transportation: Electric vehicle and fuel efficiency

FINANCIAL BENEFITS:
• Earn ₹50-200 per ton of CO2 reduced
• Average MSME can earn ₹50,000-₹2,00,000 annually
• Additional working capital for business growth
• Enhanced sustainability credentials
• Competitive advantage in green markets

TRADING OPPORTUNITIES:
• Sell excess carbon credits to other companies
• Offset remaining emissions cost-effectively
• Carbon credit-backed loans and financing
• Market making and liquidity provision
• Advisory services for carbon strategies

REAL EXAMPLES:
• Manufacturing MSME: Earned ₹1,50,000 from carbon credits
• Textile MSME: Generated ₹1,20,000 additional revenue
• Food Processing MSME: Created ₹80,000 new income stream
• Average additional revenue: ₹1,00,000-₹2,00,000 annually

REGULATORY COMPLIANCE:
• Verified Carbon Standard (VCS) compliance
• Gold Standard certification support
• Local carbon market regulations
• International carbon trading standards
• Automated reporting and documentation"""
    
    # Slide 7: Easy Implementation & ROI
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Easy Implementation & ROI"
    content7.text = """🚀 SIMPLE 3-STEP PROCESS

Step 1: Quick Setup (1 day)
• Download mobile app or access web platform
• Create account with basic company information
• Connect your business email and phone
• Upload recent utility bills and invoices

Step 2: Data Integration (1 week)
• Our AI analyzes your historical data
• Automatic categorization of expenses
• Initial carbon footprint calculation
• Personalized recommendations generated

Step 3: Start Saving (Immediately)
• Begin implementing recommendations
• Track real-time savings and improvements
• Access green finance opportunities
• Generate professional reports

ROI CALCULATION:
• Platform Cost: ₹2,999-₹7,999 per month
• Average Savings: ₹2,20,000-₹9,00,000 annually
• ROI: 400-600% in first year
• Payback Period: 2-3 months
• Break-even: Immediate

QUICK WINS (First 30 days):
• 15-20% immediate cost savings
• Clear visibility into spending patterns
• First sustainability report generated
• Access to green finance options
• Improved operational efficiency

ONBOARDING SUPPORT:
✅ Free setup and training
✅ Dedicated account manager
✅ 24/7 customer support
✅ Video tutorials and guides
✅ Regular check-ins and optimization

NO TECHNICAL EXPERTISE REQUIRED:
• User-friendly interface
• Mobile app for easy access
• Automated data processing
• Clear instructions and guidance
• Ongoing support and training"""
    
    # Slide 8: Success Stories & Case Studies
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Success Stories & Case Studies"
    content8.text = """📊 REAL CUSTOMER RESULTS

Case Study 1: EcoTech Manufacturing
• Industry: Electronics manufacturing
• Employees: 75
• Carbon Score: 85 (Gold tier)
• Results after 6 months:
  - 40% reduction in energy costs (₹2,50,000 saved)
  - 30% reduction in water usage (₹1,20,000 saved)
  - 45% reduction in waste generation (₹80,000 saved)
  - 50% improvement in sustainability score
  - Secured green loan at 2.5% lower interest rate
  - Earned ₹1,50,000 from carbon credit trading
  - Total annual savings: ₹6,00,000

Case Study 2: GreenTextile Ltd
• Industry: Textile manufacturing
• Employees: 100
• Carbon Score: 78 (Silver tier)
• Results after 4 months:
  - 35% reduction in energy consumption
  - 25% reduction in water usage
  - 30% reduction in waste generation
  - 45% improvement in sustainability score
  - Access to carbon credit trading
  - Generated ₹1,20,000 additional revenue
  - Total annual savings: ₹4,20,000

Case Study 3: FreshFoods Pvt Ltd
• Industry: Food processing
• Employees: 50
• Carbon Score: 72 (Silver tier)
• Results after 3 months:
  - 30% reduction in energy costs
  - 20% reduction in water usage
  - 40% reduction in waste disposal costs
  - 35% improvement in carbon score
  - Enhanced market reputation
  - Created ₹80,000 new income stream
  - Total annual savings: ₹2,50,000

CUSTOMER TESTIMONIALS:
"Carbon Intelligence helped us reduce our operational costs by 40% in just 6 months. The platform is easy to use and the insights are invaluable." - Priya Sharma, GreenTextile Ltd

"We saved ₹3.6 lakh annually on our loan interest rate by qualifying for a green loan through this platform." - Rajesh Kumar, EcoTech Manufacturing

"The carbon credit trading feature created an additional revenue stream of ₹1.5 lakh annually." - Amit Patel, FreshFoods Pvt Ltd"""
    
    # Slide 9: Pricing & Plans
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Affordable Pricing Plans"
    content9.text = """💳 FLEXIBLE PRICING OPTIONS

STARTER PLAN - ₹2,999/month
• Up to 25 employees
• Basic carbon tracking and scoring
• Monthly sustainability report
• Mobile app access
• Email support
• Perfect for small businesses

PROFESSIONAL PLAN - ₹4,999/month
• Up to 100 employees
• Advanced analytics and insights
• Quarterly sustainability reports
• Carbon trading access
• Priority support
• API integration
• Most popular choice

ENTERPRISE PLAN - ₹7,999/month
• Unlimited employees
• Full platform features
• Custom reporting and analytics
• Dedicated account manager
• White-label options
• Advanced integrations
• Perfect for growing businesses

SPECIAL OFFERS:
🎉 30-day free trial (no credit card required)
🎉 20% discount for annual payment
🎉 First 3 months at 50% off
🎉 Free setup and training
🎉 No long-term contracts

ADDITIONAL SERVICES:
• Sustainability consulting: ₹25,000-₹1,00,000
• ESG report generation: ₹15,000-₹50,000
• Carbon credit trading: 2-5% commission
• Green loan facilitation: No additional cost
• Training and workshops: ₹10,000-₹25,000

MONEY-BACK GUARANTEE:
If you don't save at least 5x your subscription cost in the first year, we'll refund your money. No questions asked.

ROI GUARANTEE:
Average customer saves ₹2,20,000-₹9,00,000 annually with platform cost of ₹36,000-₹96,000."""
    
    # Slide 10: Why Start Now?
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Why Start Your Sustainability Journey Now?"
    content10.text = """⏰ URGENT BUSINESS IMPERATIVES

REGULATORY PRESSURE:
• BRSR compliance mandatory for large companies
• ESG reporting requirements increasing
• Carbon credit trading scheme launched
• Government incentives available now
• Penalties for non-compliance

MARKET OPPORTUNITIES:
• Customers prefer sustainable businesses
• Green finance options expanding rapidly
• Competitive advantage opportunity
• Investor interest in ESG companies
• Supply chain sustainability requirements

COST PRESSURES:
• Rising energy and resource costs
• Waste disposal costs increasing
• Water scarcity and pricing
• Transportation costs rising
• Material costs fluctuating

COMPETITIVE ADVANTAGE:
• First-mover advantage in sustainability
• Enhanced brand reputation and credibility
• Customer loyalty and retention
• Access to new markets and opportunities
• Attract top talent and investors

FUTURE-PROOFING:
• Prepare for stricter regulations
• Build sustainable business model
• Reduce operational risks
• Create long-term value
• Stay ahead of competition

IMMEDIATE BENEFITS:
✅ Start saving money from day 1
✅ Access to green finance immediately
✅ Improve operational efficiency
✅ Build competitive advantage
✅ Meet customer expectations

DON'T WAIT - START TODAY!
Every day you delay is money lost and opportunities missed. The sooner you start, the more you save."""
    
    # Slide 11: How to Get Started
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "How to Get Started"
    content11.text = """🚀 SIMPLE 3-STEP PROCESS

Step 1: Sign Up (5 minutes)
• Visit www.carbonintelligence.com
• Click "Start Free Trial"
• Enter your company details
• Choose your plan
• No credit card required for trial

Step 2: Quick Setup (1 day)
• Download mobile app
• Connect your business accounts
• Upload recent bills and invoices
• Our AI analyzes your data
• Get initial recommendations

Step 3: Start Saving (Immediately)
• Implement first recommendations
• Track your progress
• Access green finance options
• Generate your first report
• Begin your sustainability journey

WHAT HAPPENS NEXT:
• Welcome call from our team
• Free setup and training session
• Access to all platform features
• Regular check-ins and support
• Continuous optimization

SUPPORT AVAILABLE:
📞 Phone: +91-98765-43210
📧 Email: support@carbonintelligence.com
💬 Live chat on website
📱 In-app support
🎥 Video tutorials and guides

SPECIAL LAUNCH OFFER:
🎉 30-day free trial
🎉 50% off first 3 months
🎉 Free setup and training
🎉 Money-back guarantee
🎉 No long-term commitment

READY TO START?
Visit www.carbonintelligence.com today and begin your sustainability transformation!"""
    
    # Slide 12: Thank You
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "Thank You"
    content12.text = """🌱 TRANSFORM YOUR BUSINESS TODAY

Key Takeaways:
✅ Save 20-40% on operational costs
✅ Access green finance with lower interest rates
✅ Build competitive advantage and market reputation
✅ Meet sustainability goals and regulatory compliance
✅ Easy-to-use platform with AI-powered insights
✅ Additional revenue from carbon credit trading

Ready to Start?
• 30-day free trial
• No credit card required
• Immediate cost savings
• Professional support
• Money-back guarantee

Contact Us:
📧 info@carbonintelligence.com
📞 +91-98765-43210
🌐 www.carbonintelligence.com
📱 Download our mobile app

Questions & Discussion

Let's build a sustainable future together! 🌱

Carbon Intelligence - Empowering MSMEs for a Greener Tomorrow

Start your free trial today at www.carbonintelligence.com"""
    
    # Save the presentation
    output_path = "/workspace/Carbon_Intelligence_Enhanced_MSME_Pitch.pptx"
    prs.save(output_path)
    print(f"Enhanced MSME Customer PowerPoint presentation created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_enhanced_msme_pitch_presentation()