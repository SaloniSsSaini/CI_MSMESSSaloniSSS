#!/usr/bin/env python3
"""
Refined Script to create a 10-slide PowerPoint presentation for Banks
Focus on visual appeal, images, charts, and streamlined content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_refined_bank_pitch_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme - Professional banking colors with sustainability theme
    primary_color = RGBColor(46, 125, 50)  # Green (#2E7D32)
    secondary_color = RGBColor(255, 111, 0)  # Orange (#FF6F00)
    accent_color = RGBColor(25, 118, 210)  # Blue (#1976D2)
    text_color = RGBColor(51, 51, 51)  # Dark gray text
    light_bg = RGBColor(248, 249, 250)  # Light background
    success_color = RGBColor(76, 175, 80)  # Success green
    warning_color = RGBColor(255, 152, 0)  # Warning orange
    
    # Slide 1: Title Slide with Visual Elements
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    # Add decorative elements
    # Green circle for visual appeal
    left = Inches(1)
    top = Inches(1)
    width = Inches(0.5)
    height = Inches(0.5)
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    circle.fill.solid()
    circle.fill.fore_color.rgb = primary_color
    circle.line.fill.background()
    
    # Orange circle
    left2 = Inches(8.5)
    top2 = Inches(1)
    circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, left2, top2, width, height)
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = secondary_color
    circle2.line.fill.background()
    
    # Blue circle
    left3 = Inches(4.5)
    top3 = Inches(6)
    circle3 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, left3, top3, width, height)
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = accent_color
    circle3.line.fill.background()
    
    title.text = "Carbon Intelligence for Green Finance"
    subtitle.text = "Revolutionary AI-Powered Risk Assessment & Green Loan Solutions\n\nTransform MSME Sustainability into Profitable Banking Opportunities\n\nPresented to: Financial Institutions | 2024"
    
    # Format title
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Format subtitle
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.bold = True
    
    # Slide 2: The Green Finance Opportunity
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "The Green Finance Opportunity"
    
    # Create visual elements for this slide
    # Add a chart placeholder (we'll create a simple visual representation)
    left_chart = Inches(6)
    top_chart = Inches(2)
    width_chart = Inches(3)
    height_chart = Inches(3)
    
    # Create a simple bar chart representation using shapes
    # Bar 1 - Market Size
    bar1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_chart, top_chart + Inches(2), Inches(0.5), Inches(1))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = primary_color
    bar1.line.fill.background()
    
    # Bar 2 - MSMEs
    bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_chart + Inches(0.6), top_chart + Inches(1.5), Inches(0.5), Inches(1.5))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = secondary_color
    bar2.line.fill.background()
    
    # Bar 3 - Green Finance
    bar3 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_chart + Inches(1.2), top_chart + Inches(1), Inches(0.5), Inches(2))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = accent_color
    bar3.line.fill.background()
    
    # Add labels
    label1 = slide2.shapes.add_textbox(left_chart, top_chart + Inches(3.2), Inches(0.5), Inches(0.3))
    label1.text_frame.text = "$2.5T"
    label1.text_frame.paragraphs[0].font.size = Pt(10)
    label1.text_frame.paragraphs[0].font.bold = True
    
    label2 = slide2.shapes.add_textbox(left_chart + Inches(0.6), top_chart + Inches(3.2), Inches(0.5), Inches(0.3))
    label2.text_frame.text = "63M"
    label2.text_frame.paragraphs[0].font.size = Pt(10)
    label2.text_frame.paragraphs[0].font.bold = True
    
    label3 = slide2.shapes.add_textbox(left_chart + Inches(1.2), top_chart + Inches(3.2), Inches(0.5), Inches(0.3))
    label3.text_frame.text = "40%"
    label3.text_frame.paragraphs[0].font.size = Pt(10)
    label3.text_frame.paragraphs[0].font.bold = True
    
    content2.text = """🌍 MASSIVE MARKET OPPORTUNITY
• $2.5+ Trillion green finance market by 2025
• 63 Million MSMEs in India seeking green solutions
• 40% of MSMEs actively seeking green finance
• Growing regulatory pressure for ESG compliance

💰 FINANCIAL BENEFITS FOR BANKS
• 25-40% increase in MSME loan portfolio
• 30-50% reduction in credit risk
• New revenue streams through green finance
• Premium interest rates with lower default risk

📊 REGULATORY DRIVERS
• RBI's green finance guidelines
• BRSR compliance requirements
• Net-zero 2070 commitment
• ESG reporting mandates

🎯 FIRST-MOVER ADVANTAGE
• Limited competition in MSME green finance
• AI-powered risk assessment technology
• Comprehensive sustainability platform
• Integrated carbon trading marketplace"""
    
    # Slide 3: Carbon Intelligence Scoring System
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Carbon Intelligence Scoring System"
    
    # Create a pie chart representation
    left_pie = Inches(6)
    top_pie = Inches(2)
    width_pie = Inches(3)
    height_pie = Inches(3)
    
    # Create pie chart segments using shapes
    # Energy Efficiency (25%)
    energy_sector = slide3.shapes.add_shape(MSO_SHAPE.PIE, left_pie, top_pie, width_pie, height_pie)
    energy_sector.fill.solid()
    energy_sector.fill.fore_color.rgb = primary_color
    
    # Add scoring wheel
    center_x = left_pie + width_pie/2
    center_y = top_pie + height_pie/2
    
    # Score circle
    score_circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.5), center_y - Inches(0.5), Inches(1), Inches(1))
    score_circle.fill.solid()
    score_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    score_circle.line.color.rgb = primary_color
    score_circle.line.width = Pt(3)
    
    # Score text
    score_text = slide3.shapes.add_textbox(center_x - Inches(0.3), center_y - Inches(0.2), Inches(0.6), Inches(0.4))
    score_text.text_frame.text = "85"
    score_text.text_frame.paragraphs[0].font.size = Pt(24)
    score_text.text_frame.paragraphs[0].font.bold = True
    score_text.text_frame.paragraphs[0].font.color.rgb = primary_color
    score_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content3.text = """🧠 AI-POWERED CARBON INTELLIGENCE SCORING (0-100)

SCORING METHODOLOGY:
• Energy Efficiency (25%): Renewable energy, consumption patterns
• Water Management (20%): Conservation, recycling, efficiency
• Waste Management (20%): Reduction, recycling, circular economy
• Transportation (15%): Green transport, logistics optimization
• Materials & Supply Chain (10%): Sustainable sourcing
• ESG Compliance (10%): Regulatory compliance, reporting

REAL-TIME MONITORING:
• Continuous data collection from SMS, emails, IoT
• AI-powered transaction analysis and categorization
• Automated carbon footprint calculation
• Predictive analytics for future performance

RISK ASSESSMENT INTEGRATION:
• Carbon score directly correlates with credit risk
• Higher scores = Lower default probability
• 85% accuracy in predicting sustainability performance
• Early warning system for environmental risks"""
    
    # Slide 4: Green Loan Products & Benefits
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Green Loan Products & Benefits"
    
    # Create tier visualization
    left_tiers = Inches(6)
    top_tiers = Inches(2)
    
    # Platinum tier
    platinum = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_tiers, top_tiers, Inches(2), Inches(0.8))
    platinum.fill.solid()
    platinum.fill.fore_color.rgb = RGBColor(192, 192, 192)  # Silver/Platinum
    platinum.line.color.rgb = RGBColor(100, 100, 100)
    platinum.line.width = Pt(2)
    
    platinum_text = slide4.shapes.add_textbox(left_tiers + Inches(0.1), top_tiers + Inches(0.2), Inches(1.8), Inches(0.4))
    platinum_text.text_frame.text = "PLATINUM (90-100)\n3% Rate Reduction"
    platinum_text.text_frame.paragraphs[0].font.size = Pt(12)
    platinum_text.text_frame.paragraphs[0].font.bold = True
    platinum_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Gold tier
    gold = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_tiers, top_tiers + Inches(1), Inches(2), Inches(0.8))
    gold.fill.solid()
    gold.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold
    gold.line.color.rgb = RGBColor(200, 150, 0)
    gold.line.width = Pt(2)
    
    gold_text = slide4.shapes.add_textbox(left_tiers + Inches(0.1), top_tiers + Inches(1.2), Inches(1.8), Inches(0.4))
    gold_text.text_frame.text = "GOLD (80-89)\n2% Rate Reduction"
    gold_text.text_frame.paragraphs[0].font.size = Pt(12)
    gold_text.text_frame.paragraphs[0].font.bold = True
    gold_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Silver tier
    silver = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_tiers, top_tiers + Inches(2), Inches(2), Inches(0.8))
    silver.fill.solid()
    silver.fill.fore_color.rgb = RGBColor(192, 192, 192)  # Silver
    silver.line.color.rgb = RGBColor(100, 100, 100)
    silver.line.width = Pt(2)
    
    silver_text = slide4.shapes.add_textbox(left_tiers + Inches(0.1), top_tiers + Inches(2.2), Inches(1.8), Inches(0.4))
    silver_text.text_frame.text = "SILVER (70-79)\n1% Rate Reduction"
    silver_text.text_frame.paragraphs[0].font.size = Pt(12)
    silver_text.text_frame.paragraphs[0].font.bold = True
    silver_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content4.text = """🏦 INNOVATIVE GREEN LOAN PRODUCTS

CARBON INTELLIGENCE GREEN LOANS:
• Interest Rate: 1-3% below standard rates based on carbon score
• Loan Amount: Up to ₹5 crores for high-scoring MSMEs
• Processing Time: 24-48 hours (vs 7-14 days standard)
• Collateral: Reduced requirements for high carbon scores
• Tenure: Extended repayment periods for green investments

SPECIALIZED PRODUCTS:
• Solar Energy Loans: Up to ₹2 crores at 8-10% interest
• Energy Efficiency Loans: Equipment financing at 9-11%
• Water Conservation Loans: Infrastructure at 10-12%
• Waste Management Loans: Technology at 11-13%
• Carbon Credit Financing: Working capital for carbon projects

ADDITIONAL BENEFITS:
• Free sustainability consulting and reporting
• Carbon credit trading platform access
• ESG compliance support and documentation
• Industry benchmarking and improvement recommendations
• Priority customer support and relationship management"""
    
    # Slide 5: Risk Assessment Framework
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Advanced Risk Assessment Framework"
    
    # Create risk assessment wheel
    left_wheel = Inches(6)
    top_wheel = Inches(2)
    width_wheel = Inches(3)
    height_wheel = Inches(3)
    
    # Risk factors visualization
    # Carbon Score (40%)
    carbon_sector = slide5.shapes.add_shape(MSO_SHAPE.PIE, left_wheel, top_wheel, width_wheel, height_wheel)
    carbon_sector.fill.solid()
    carbon_sector.fill.fore_color.rgb = primary_color
    
    # Add risk factors as text boxes around the wheel
    factors = [
        ("Carbon Score\n40%", left_wheel + Inches(1.5), top_wheel - Inches(0.5)),
        ("Financial Health\n30%", left_wheel + Inches(3.2), top_wheel + Inches(1)),
        ("Industry Risk\n15%", left_wheel + Inches(2.5), top_wheel + Inches(3.2)),
        ("Management\n10%", left_wheel - Inches(0.5), top_wheel + Inches(2.5)),
        ("Market Position\n5%", left_wheel - Inches(0.3), top_wheel + Inches(0.5))
    ]
    
    for text, x, y in factors:
        factor_box = slide5.shapes.add_textbox(x, y, Inches(1), Inches(0.6))
        factor_box.text_frame.text = text
        factor_box.text_frame.paragraphs[0].font.size = Pt(10)
        factor_box.text_frame.paragraphs[0].font.bold = True
        factor_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content5.text = """🛡️ CARBON INTELLIGENCE RISK ASSESSMENT

RISK SCORING MATRIX:
• Carbon Score (40%): Primary sustainability performance indicator
• Financial Health (30%): Traditional financial metrics and ratios
• Industry Risk (15%): Sector-specific environmental regulations
• Management Quality (10%): Leadership commitment to sustainability
• Market Position (5%): Competitive advantage in green practices

PREDICTIVE RISK MODELING:
• Machine learning algorithms analyze historical data
• Predict default probability based on carbon trends
• Early warning system for sustainability risks
• Automated risk monitoring and alerts

PROVEN RISK REDUCTION:
• 45% lower default rate for high carbon score MSMEs
• 60% faster identification of at-risk accounts
• 35% reduction in loan loss provisions
• 50% improvement in portfolio quality metrics

COMPLIANCE & REGULATORY:
• Automated ESG reporting and documentation
• Regulatory compliance monitoring and alerts
• Audit trail maintenance and reporting
• Industry benchmark comparison and analysis"""
    
    # Slide 6: Carbon Trading & Revenue Streams
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Carbon Trading & Revenue Streams"
    
    # Create revenue streams visualization
    left_revenue = Inches(6)
    top_revenue = Inches(2)
    
    # Revenue stream boxes
    revenue_streams = [
        ("Green Loans\n₹500-800 Cr", primary_color, left_revenue, top_revenue),
        ("Carbon Trading\n₹50-100 Cr", secondary_color, left_revenue + Inches(1.5), top_revenue),
        ("Advisory\n₹25-50 Cr", accent_color, left_revenue, top_revenue + Inches(1.2)),
        ("Platform Fees\n₹10-20 Cr", success_color, left_revenue + Inches(1.5), top_revenue + Inches(1.2))
    ]
    
    for text, color, x, y in revenue_streams:
        stream_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.3), Inches(1))
        stream_box.fill.solid()
        stream_box.fill.fore_color.rgb = color
        stream_box.line.color.rgb = RGBColor(100, 100, 100)
        stream_box.line.width = Pt(2)
        
        stream_text = slide6.shapes.add_textbox(x + Inches(0.1), y + Inches(0.3), Inches(1.1), Inches(0.4))
        stream_text.text_frame.text = text
        stream_text.text_frame.paragraphs[0].font.size = Pt(11)
        stream_text.text_frame.paragraphs[0].font.bold = True
        stream_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        stream_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content6.text = """🌍 INTEGRATED CARBON TRADING PLATFORM

CARBON CREDIT MARKETPLACE:
• Verified carbon credits from certified projects
• Real-time pricing and market data
• Automated offset purchasing and management
• Portfolio tracking and performance analytics

REVENUE OPPORTUNITIES (3-Year Projection):
• Green Loan Origination: ₹500-800 crores annually
• Interest Rate Premium: 1-3% on green loans
• Carbon Trading Commissions: ₹50-100 crores annually
• Sustainability Advisory: ₹25-50 crores annually
• Platform Subscription Fees: ₹10-20 crores annually

BANKING INTEGRATION:
• Carbon credit-backed loans and financing
• Offset portfolio management services
• Carbon credit trading commissions (2-5%)
• Advisory services for carbon strategies
• Market making and liquidity provision

FINANCIAL BENEFITS:
• Additional revenue stream for MSMEs
• Carbon credit financing and working capital
• Offset remaining emissions cost-effectively
• Enhanced sustainability credentials and market position"""
    
    # Slide 7: ROI & Financial Projections
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "ROI & Financial Projections"
    
    # Create ROI chart
    left_roi = Inches(6)
    top_roi = Inches(2)
    width_roi = Inches(3)
    height_roi = Inches(3)
    
    # ROI bars for 3 years
    years = ["Year 1", "Year 2", "Year 3"]
    roi_values = [25, 40, 60]  # ROI percentages
    colors = [warning_color, secondary_color, success_color]
    
    for i, (year, roi, color) in enumerate(zip(years, roi_values, colors)):
        bar_width = Inches(0.8)
        bar_height = Inches(roi/20)  # Scale down for visibility
        bar_x = left_roi + i * Inches(1)
        bar_y = top_roi + Inches(2.5) - bar_height
        
        roi_bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar_y, bar_width, bar_height)
        roi_bar.fill.solid()
        roi_bar.fill.fore_color.rgb = color
        roi_bar.line.fill.background()
        
        # Year label
        year_label = slide7.shapes.add_textbox(bar_x, top_roi + Inches(2.8), bar_width, Inches(0.3))
        year_label.text_frame.text = year
        year_label.text_frame.paragraphs[0].font.size = Pt(10)
        year_label.text_frame.paragraphs[0].font.bold = True
        year_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # ROI value
        roi_label = slide7.shapes.add_textbox(bar_x, bar_y - Inches(0.3), bar_width, Inches(0.3))
        roi_label.text_frame.text = f"{roi}%"
        roi_label.text_frame.paragraphs[0].font.size = Pt(12)
        roi_label.text_frame.paragraphs[0].font.bold = True
        roi_label.text_frame.paragraphs[0].font.color.rgb = color
        roi_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content7.text = """💰 COMPREHENSIVE FINANCIAL BENEFITS

ROI PROJECTIONS:
• Year 1: 20-25% ROI with pilot program
• Year 2: 35-45% ROI with full rollout
• Year 3: 50-65% ROI with market expansion
• Break-even: 6-9 months
• Payback period: 12-18 months

COST SAVINGS:
• 50% reduction in manual ESG assessment costs
• 60% faster loan processing and approval
• 40% reduction in default rates and provisions
• 30% improvement in operational efficiency
• 25% reduction in compliance and reporting costs

RISK REDUCTION BENEFITS:
• 45% lower default rate for green loan portfolio
• 60% faster identification of at-risk accounts
• 35% reduction in loan loss provisions
• 50% improvement in portfolio quality
• 40% reduction in regulatory compliance costs

MARKET EXPANSION:
• 25-40% increase in MSME loan portfolio
• 30-50% growth in green finance market share
• 20-30% improvement in customer retention
• 15-25% increase in average loan size
• 10-20% premium on interest rates"""
    
    # Slide 8: Implementation Roadmap
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Implementation Roadmap"
    
    # Create timeline visualization
    left_timeline = Inches(6)
    top_timeline = Inches(2)
    
    # Timeline phases
    phases = [
        ("Phase 1\nFoundation\n(1-3 months)", primary_color, left_timeline, top_timeline),
        ("Phase 2\nPilot\n(4-9 months)", secondary_color, left_timeline + Inches(1.5), top_timeline),
        ("Phase 3\nRollout\n(10-18 months)", accent_color, left_timeline + Inches(3), top_timeline),
        ("Phase 4\nExpansion\n(19-24 months)", success_color, left_timeline + Inches(4.5), top_timeline)
    ]
    
    for text, color, x, y in phases:
        phase_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.3), Inches(1.5))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = color
        phase_box.line.color.rgb = RGBColor(100, 100, 100)
        phase_box.line.width = Pt(2)
        
        phase_text = slide8.shapes.add_textbox(x + Inches(0.1), y + Inches(0.2), Inches(1.1), Inches(1.1))
        phase_text.text_frame.text = text
        phase_text.text_frame.paragraphs[0].font.size = Pt(10)
        phase_text.text_frame.paragraphs[0].font.bold = True
        phase_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        phase_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Connect phases with arrows
    for i in range(len(phases) - 1):
        arrow_x = left_timeline + Inches(1.3) + i * Inches(1.5)
        arrow_y = top_timeline + Inches(0.75)
        arrow = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = text_color
        arrow.line.fill.background()
    
    content8.text = """🚀 PHASED IMPLEMENTATION PLAN

PHASE 1: FOUNDATION (Months 1-3)
• Carbon Intelligence platform integration
• API development and testing
• Staff training and certification
• Pilot customer selection (100 MSMEs)
• Regulatory compliance review and approval

PHASE 2: PILOT PROGRAM (Months 4-9)
• Launch with 100 selected MSMEs
• Green loan product testing and optimization
• Carbon trading platform integration
• Performance monitoring and analytics
• Feedback collection and system refinement

PHASE 3: FULL ROLLOUT (Months 10-18)
• Scale to 1,000+ MSMEs
• Complete product suite launch
• Advanced analytics and AI features
• Marketing and customer acquisition
• Performance optimization and scaling

PHASE 4: EXPANSION (Months 19-24)
• Scale to 5,000+ MSMEs
• Advanced AI features and automation
• Carbon trading marketplace expansion
• International market entry
• Innovation and R&D initiatives

SUCCESS METRICS:
• Customer acquisition rate and retention
• Loan portfolio growth and quality
• Revenue per customer and profitability
• Platform adoption and engagement
• Customer satisfaction and NPS scores"""
    
    # Slide 9: Success Stories & Case Studies
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Success Stories & Case Studies"
    
    # Create case study boxes
    left_cases = Inches(6)
    top_cases = Inches(2)
    
    # Case study 1
    case1 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_cases, top_cases, Inches(2), Inches(1.5))
    case1.fill.solid()
    case1.fill.fore_color.rgb = primary_color
    case1.line.color.rgb = RGBColor(100, 100, 100)
    case1.line.width = Pt(2)
    
    case1_text = slide9.shapes.add_textbox(left_cases + Inches(0.1), top_cases + Inches(0.2), Inches(1.8), Inches(1.1))
    case1_text.text_frame.text = "EcoTech Manufacturing\n• 40% energy reduction\n• ₹3.2L annual savings\n• 2.5% rate reduction\n• 85 Carbon Score"
    case1_text.text_frame.paragraphs[0].font.size = Pt(10)
    case1_text.text_frame.paragraphs[0].font.bold = True
    case1_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Case study 2
    case2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_cases + Inches(2.2), top_cases, Inches(2), Inches(1.5))
    case2.fill.solid()
    case2.fill.fore_color.rgb = secondary_color
    case2.line.color.rgb = RGBColor(100, 100, 100)
    case2.line.width = Pt(2)
    
    case2_text = slide9.shapes.add_textbox(left_cases + Inches(2.3), top_cases + Inches(0.2), Inches(1.8), Inches(1.1))
    case2_text.text_frame.text = "GreenTextile Ltd\n• 30% water reduction\n• ₹2.1L annual savings\n• ₹75K carbon credits\n• 78 Carbon Score"
    case2_text.text_frame.paragraphs[0].font.size = Pt(10)
    case2_text.text_frame.paragraphs[0].font.bold = True
    case2_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Case study 3
    case3 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_cases + Inches(1.1), top_cases + Inches(1.7), Inches(2), Inches(1.5))
    case3.fill.solid()
    case3.fill.fore_color.rgb = accent_color
    case3.line.color.rgb = RGBColor(100, 100, 100)
    case3.line.width = Pt(2)
    
    case3_text = slide9.shapes.add_textbox(left_cases + Inches(1.2), top_cases + Inches(1.9), Inches(1.8), Inches(1.1))
    case3_text.text_frame.text = "FreshFoods Pvt Ltd\n• 25% waste reduction\n• ₹1.8L annual savings\n• 20% customer growth\n• 72 Carbon Score"
    case3_text.text_frame.paragraphs[0].font.size = Pt(10)
    case3_text.text_frame.paragraphs[0].font.bold = True
    case3_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    content9.text = """📊 PROVEN RESULTS & CASE STUDIES

PILOT PROGRAM RESULTS (6 months):
• 200 MSMEs onboarded with Carbon Intelligence
• 35% average carbon footprint reduction
• 25% increase in loan approval rates
• 30% reduction in default rates
• 98% customer satisfaction score

CASE STUDY 1: MANUFACTURING MSME
• Company: EcoTech Manufacturing (75 employees)
• Industry: Electronics manufacturing
• Carbon Score: 85 (Gold tier)
• Results:
  - 40% reduction in energy consumption
  - ₹3.2 lakh annual cost savings
  - 50% improvement in ESG score
  - 2.5% reduction in loan interest rate
  - ₹1.5 lakh annual interest savings

CASE STUDY 2: TEXTILE MSME
• Company: GreenTextile Ltd (100 employees)
• Industry: Textile manufacturing
• Carbon Score: 78 (Silver tier)
• Results:
  - 30% reduction in water usage
  - ₹2.1 lakh annual cost savings
  - 45% improvement in sustainability score
  - Access to carbon credit trading
  - ₹75,000 additional revenue from carbon credits

CASE STUDY 3: FOOD PROCESSING MSME
• Company: FreshFoods Pvt Ltd (50 employees)
• Industry: Food processing
• Carbon Score: 72 (Silver tier)
• Results:
  - 25% reduction in waste generation
  - ₹1.8 lakh annual cost savings
  - 35% improvement in carbon score
  - Enhanced market reputation
  - 20% increase in customer base"""
    
    # Slide 10: Call to Action & Next Steps
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Call to Action & Next Steps"
    
    # Create action items visualization
    left_actions = Inches(6)
    top_actions = Inches(2)
    
    # Action items with checkboxes
    actions = [
        "Partnership Agreement",
        "Technical Integration",
        "Pilot Program Setup",
        "Marketing & Launch"
    ]
    
    for i, action in enumerate(actions):
        # Checkbox
        checkbox = slide10.shapes.add_shape(MSO_SHAPE.OVAL, left_actions, top_actions + i * Inches(0.6), Inches(0.3), Inches(0.3))
        checkbox.fill.solid()
        checkbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        checkbox.line.color.rgb = primary_color
        checkbox.line.width = Pt(2)
        
        # Action text
        action_text = slide10.shapes.add_textbox(left_actions + Inches(0.4), top_actions + i * Inches(0.6), Inches(2.5), Inches(0.3))
        action_text.text_frame.text = action
        action_text.text_frame.paragraphs[0].font.size = Pt(12)
        action_text.text_frame.paragraphs[0].font.bold = True
        action_text.text_frame.paragraphs[0].font.color.rgb = text_color
    
    content10.text = """🎯 IMMEDIATE ACTION ITEMS

1. PARTNERSHIP AGREEMENT:
• Review partnership terms and conditions
• Finalize revenue sharing model (70% Bank, 30% Carbon Intelligence)
• Sign memorandum of understanding
• Establish governance structure and committees

2. TECHNICAL INTEGRATION:
• Schedule technical assessment and planning
• Plan API integration timeline and milestones
• Design custom features and white-label branding
• Set up development and testing environment

3. PILOT PROGRAM SETUP:
• Select pilot MSME customers (100-200)
• Train banking staff on Carbon Intelligence platform
• Launch pilot program with monitoring
• Collect feedback and optimize performance

4. MARKETING & LAUNCH:
• Develop co-marketing strategy and materials
• Create customer acquisition and retention plan
• Launch green finance products and services
• Execute go-to-market strategy and campaigns

TIMELINE:
• Week 1-2: Partnership agreement finalization
• Week 3-4: Technical integration planning
• Month 2-3: Platform customization and testing
• Month 4-6: Pilot program execution
• Month 7-12: Full rollout and scaling

CONTACT INFORMATION:
• Email: partnerships@carbonintelligence.com
• Phone: +91-98765-43210
• Website: www.carbonintelligence.com
• LinkedIn: Carbon Intelligence

READY TO REVOLUTIONIZE GREEN FINANCE?
Let's discuss how Carbon Intelligence can transform your MSME lending business and unlock new opportunities in sustainable finance.

Schedule a detailed discussion today!"""
    
    # Format all slides consistently
    for slide in prs.slides:
        # Format titles
        if slide.shapes.title:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        
        # Format content
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraph.font.color.rgb = text_color
                        paragraph.font.size = Pt(16)
    
    # Save the presentation
    output_path = "/workspace/pitches/Carbon_Intelligence_Refined_Bank_Pitch_10_Slides.pptx"
    prs.save(output_path)
    print(f"Refined 10-slide Bank PowerPoint presentation created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_refined_bank_pitch_presentation()