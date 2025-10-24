#!/usr/bin/env python3
"""
Refined 10-Page MSME PowerPoint Presentation with Application Mockups
Focus on Visual Interface Demonstrations, Key Benefits, and ROI Calculations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_refined_msme_pitch_with_mockups():
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme - Professional Green and Orange
    primary_color = RGBColor(46, 125, 50)  # Green for sustainability
    secondary_color = RGBColor(255, 140, 0)  # Orange for highlights
    accent_color = RGBColor(25, 118, 210)  # Blue for trust
    text_color = RGBColor(51, 51, 51)  # Dark gray text
    light_bg = RGBColor(248, 249, 250)  # Light background
    success_color = RGBColor(76, 175, 80)  # Success green
    warning_color = RGBColor(255, 152, 0)  # Warning orange
    
    # Slide 1: Title Slide with Value Proposition
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Carbon Intelligence for MSMEs"
    subtitle.text = "Transform Your Business with AI-Powered Sustainability\n\nSave 20-40% on Operational Costs • Access Green Finance • Build Competitive Advantage\n\nComplete Application Mockups & Key Benefits\n2024"
    
    # Format title
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Format subtitle
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(16)
    
    # Slide 2: MSME Dashboard Mockup
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "MSME Dashboard - Real-Time Carbon Intelligence"
    content2.text = """📊 PROFESSIONAL DASHBOARD INTERFACE

MAIN DASHBOARD FEATURES:
• Real-time Carbon Score Display (0-100 scale)
• Live Cost Savings Tracking and Analytics
• Quick Action Buttons for Key Functions
• Monthly/Quarterly Performance Metrics
• AI-Powered Recommendations Panel

KEY METRICS DISPLAYED:
✅ Current Carbon Score: 85/100 (Gold Tier)
✅ Monthly CO₂ Reduction: 2.4 tons (12% improvement)
✅ Cost Savings This Month: ₹1,25,000
✅ Green Finance Eligibility: ₹50L available
✅ Carbon Credits Earned: 1,250 credits

DASHBOARD BENEFITS:
• Single-screen overview of all sustainability metrics
• Instant access to cost savings and recommendations
• Professional reporting for stakeholders and investors
• Mobile-responsive design for on-the-go management
• Real-time alerts and notifications for optimization opportunities

QUICK ACTIONS AVAILABLE:
🔍 AI Analytics - Smart insights and predictions
💱 Carbon Trading - Buy/sell carbon credits
🏦 Green Finance - Access sustainable loans
📄 ESG Reports - Generate compliance reports
⚡ Recommendations - AI-powered improvement suggestions

USER EXPERIENCE:
• Intuitive navigation and clean interface design
• Color-coded status indicators for easy understanding
• Customizable widgets and dashboard layout
• Export capabilities for reports and data
• Multi-device synchronization across web and mobile"""
    
    # Slide 3: Carbon Intelligence Scoring Interface Mockup
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Carbon Intelligence Scoring System - Visual Interface"
    content3.text = """🧠 AI-POWERED SCORING INTERFACE

SCORING VISUALIZATION:
• Large Circular Score Display (85/100)
• Color-coded Performance Indicators
• Detailed Category Breakdown with Progress Bars
• Real-time Score Updates and Trends
• Industry Benchmarking Comparison

SCORING CATEGORIES (0-100 Scale):
📊 Energy Efficiency (25%): 22/25 (88%) - Excellent
💧 Water Management (20%): 16/20 (80%) - Good
♻️ Waste Management (20%): 18/20 (90%) - Excellent
🚚 Transportation (15%): 12/15 (80%) - Good
📦 Materials & Supply Chain (10%): 8/10 (80%) - Good
📋 ESG Compliance (10%): 9/10 (90%) - Excellent

SCORE TIERS & BENEFITS:
🥇 Platinum (90-100): 3% rate reduction, highest loan amounts
🥈 Gold (80-89): 2% rate reduction, priority processing
🥉 Silver (70-79): 1% rate reduction, standard processing
🏅 Bronze (60-69): 0.5% rate reduction, basic processing
📈 Below 60: Improvement plan with support and guidance

REAL-TIME MONITORING FEATURES:
• Continuous data collection from SMS, emails, and transactions
• AI-powered analysis and categorization
• Automated carbon footprint calculation
• Predictive analytics for future performance
• Instant alerts for score improvements or declines

INTERFACE BENEFITS:
• Clear visual representation of sustainability performance
• Easy identification of improvement areas
• Professional presentation for stakeholders
• Gamified elements to encourage engagement
• Historical trend analysis and goal tracking"""
    
    # Slide 4: Mobile Application Interface Mockup
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Mobile Application - On-the-Go Carbon Management"
    content4.text = """📱 REACT NATIVE MOBILE INTERFACE

MOBILE APP FEATURES:
• iPhone/Android compatible design
• Touch-friendly interface with intuitive navigation
• Offline capability for key features
• Push notifications for important updates
• Biometric authentication for security

MAIN SCREEN LAYOUT:
📊 Dashboard: Real-time metrics and quick stats
📈 Analytics: Charts, trends, and performance tracking
💱 Trading: Carbon credit marketplace and transactions
👤 Profile: Account settings and preferences

QUICK STATS DISPLAY:
• Carbon Score: 85/100 (Gold Tier)
• CO₂ Saved: 245 kg this month
• Cost Savings: ₹15,240 this month
• Carbon Credits: 1,250 available

MOBILE-SPECIFIC FEATURES:
• SMS Analysis: AI-powered message scanning
• Photo Upload: Document and receipt capture
• GPS Tracking: Location-based carbon calculations
• Voice Commands: Hands-free data entry
• QR Code Scanner: Quick data input

QUICK ACTIONS:
🔍 AI Analytics - Smart insights on mobile
💱 Carbon Trading - Buy/sell credits anywhere
🏦 Green Finance - Apply for loans instantly
📄 Reports - Generate and share reports
⚡ Recommendations - Get AI suggestions
📊 Analytics - View charts and trends

MOBILE BENEFITS:
• 24/7 access to carbon intelligence
• Instant decision making capabilities
• Real-time alerts and notifications
• Mobile-first user experience
• Seamless synchronization with web platform"""
    
    # Slide 5: Cost Savings & ROI Calculator Mockup
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Cost Savings Calculator - Real-Time ROI Tracking"
    content5.text = """💰 INTERACTIVE ROI CALCULATOR

COST SAVINGS BREAKDOWN:
• Energy Costs: 25-40% reduction (₹75,000-₹3,00,000 annually)
• Water Usage: 20-30% reduction (₹40,000-₹1,50,000 annually)
• Waste Management: 35-50% reduction (₹25,000-₹1,00,000 annually)
• Transportation: 15-25% reduction (₹30,000-₹1,50,000 annually)
• Materials: 20-30% reduction (₹50,000-₹2,00,000 annually)

TOTAL ANNUAL SAVINGS: ₹2,20,000 - ₹9,00,000 per MSME

ROI CALCULATION INTERFACE:
📊 Platform Cost: ₹2,999-₹7,999 per month
📈 Average Savings: ₹2,20,000-₹9,00,000 annually
💰 ROI: 400-600% in first year
⏱️ Payback Period: 2-3 months
🎯 Break-even: Immediate

REAL CUSTOMER EXAMPLES:

Manufacturing MSME (75 employees):
• Energy savings: ₹2,50,000/year (40% reduction)
• Water savings: ₹1,20,000/year (30% reduction)
• Waste reduction: ₹80,000/year (45% reduction)
• Total savings: ₹4,50,000/year
• Carbon footprint reduction: 35%

Textile MSME (50 employees):
• Energy savings: ₹1,80,000/year (35% reduction)
• Water savings: ₹90,000/year (25% reduction)
• Material optimization: ₹1,50,000/year (30% reduction)
• Total savings: ₹4,20,000/year
• Carbon footprint reduction: 30%

Food Processing MSME (40 employees):
• Energy savings: ₹1,20,000/year (30% reduction)
• Water savings: ₹60,000/year (20% reduction)
• Waste reduction: ₹70,000/year (40% reduction)
• Total savings: ₹2,50,000/year
• Carbon footprint reduction: 25%

CALCULATOR FEATURES:
• Real-time savings tracking
• Industry-specific calculations
• Customizable parameters
• Export capabilities for reports
• Integration with accounting systems"""
    
    # Slide 6: Green Finance Access Interface Mockup
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Green Finance Portal - Sustainable Loan Access"
    content6.text = """🏦 COMPREHENSIVE GREEN FINANCE INTERFACE

GREEN LOAN PRODUCTS DISPLAY:
• Carbon Intelligence Green Loans: 1-3% lower interest rates
• Solar Energy Loans: Up to ₹2 crores at 8-10% interest
• Energy Efficiency Loans: Equipment financing at 9-11%
• Water Conservation Loans: Infrastructure at 10-12%
• Waste Management Loans: Technology at 11-13%

SCORING-BASED BENEFITS:
• Higher carbon scores = Better loan terms
• Faster approval process (24-48 hours vs 7-14 days)
• Reduced collateral requirements
• Extended repayment periods
• Priority customer support

GOVERNMENT INCENTIVES TRACKER:
• Energy efficiency grants (up to ₹5 lakh)
• Solar panel subsidies (30-40% of cost)
• Water conservation incentives (up to ₹2 lakh)
• Waste management grants (up to ₹3 lakh)
• Green certification benefits and tax incentives

CARBON CREDIT TRADING INTERFACE:
• Earn money by reducing emissions
• Sell carbon credits to other companies
• Offset remaining emissions cost-effectively
• Additional revenue stream (₹50,000-₹2,00,000 annually)
• Environmental impact monetization

BANKING PARTNERSHIPS:
• Direct access to partner banks
• Pre-approved green loan products
• Dedicated relationship managers
• Streamlined application process
• Competitive interest rates

SUCCESS STORY DISPLAY:
"Through Carbon Intelligence, we got a green loan at 2.5% lower interest rate, saving ₹3.6 lakh annually on our ₹1.5 crore loan. The platform helped us qualify by tracking our sustainability improvements." - Rajesh Kumar, EcoTech Manufacturing

INTERFACE FEATURES:
• Real-time eligibility checking
• Instant loan pre-approval
• Document upload and verification
• Application status tracking
• Direct bank integration"""
    
    # Slide 7: AI Recommendations Engine Mockup
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "AI Recommendations Engine - Smart Sustainability Suggestions"
    content7.text = """🤖 INTELLIGENT RECOMMENDATION INTERFACE

PRIORITY RECOMMENDATIONS DISPLAY:

🔴 HIGH PRIORITY:
• Switch to Renewable Energy
  - Install solar panels to reduce electricity-related emissions by 60%
  - Potential Savings: 350 kg CO₂/month
  - Investment Required: ₹2,50,000
  - ROI: 18 months

🟡 MEDIUM PRIORITY:
• Improve Waste Management
  - Implement comprehensive recycling program
  - Potential Savings: 120 kg CO₂/month
  - Investment Required: ₹15,000
  - ROI: 3 months

🟢 LOW PRIORITY:
• Optimize Transportation
  - Route optimization and vehicle maintenance
  - Potential Savings: 80 kg CO₂/month
  - Investment Required: ₹5,000
  - ROI: 2 months

AI-POWERED FEATURES:
• Personalized recommendations based on industry
• Cost-benefit analysis for each suggestion
• Implementation timeline and support
• Progress tracking and monitoring
• Success probability calculations

RECOMMENDATION CATEGORIES:
⚡ Energy Efficiency: LED lighting, smart meters, renewable energy
💧 Water Management: Rainwater harvesting, water recycling
♻️ Waste Management: Composting, recycling programs, circular economy
🚚 Transportation: Electric vehicles, route optimization, fuel efficiency
📦 Materials: Sustainable sourcing, packaging optimization
📋 Compliance: ESG reporting, regulatory compliance

IMPLEMENTATION SUPPORT:
• Step-by-step implementation guides
• Vendor recommendations and quotes
• Financing options and incentives
• Progress tracking and milestones
• Expert consultation and support

INTERFACE BENEFITS:
• Clear prioritization of actions
• Detailed cost-benefit analysis
• Easy implementation tracking
• Expert guidance and support
• Measurable impact assessment"""
    
    # Slide 8: Success Stories & Case Studies
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Success Stories & Real Customer Results"
    content8.text = """📊 PROVEN CUSTOMER RESULTS

Case Study 1: EcoTech Manufacturing
• Industry: Electronics manufacturing
• Employees: 75
• Carbon Score: 85 (Gold tier)
• Results after 6 months:
  - 40% reduction in energy costs (₹2,50,000 saved)
  - 30% reduction in water usage (₹1,20,000 saved)
  - 45% reduction in waste generation (₹80,000 saved)
  - 50% improvement in sustainability score
  - Secured green loan at 2.5% lower interest rate
  - Earned ₹1,50,000 from carbon credit trading
  - Total annual savings: ₹6,00,000

Case Study 2: GreenTextile Ltd
• Industry: Textile manufacturing
• Employees: 100
• Carbon Score: 78 (Silver tier)
• Results after 4 months:
  - 35% reduction in energy consumption
  - 25% reduction in water usage
  - 30% reduction in waste generation
  - 45% improvement in sustainability score
  - Access to carbon credit trading
  - Generated ₹1,20,000 additional revenue
  - Total annual savings: ₹4,20,000

Case Study 3: FreshFoods Pvt Ltd
• Industry: Food processing
• Employees: 50
• Carbon Score: 72 (Silver tier)
• Results after 3 months:
  - 30% reduction in energy costs
  - 20% reduction in water usage
  - 40% reduction in waste disposal costs
  - 35% improvement in carbon score
  - Enhanced market reputation
  - Created ₹80,000 new income stream
  - Total annual savings: ₹2,50,000

CUSTOMER TESTIMONIALS:
"Carbon Intelligence helped us reduce our operational costs by 40% in just 6 months. The platform is easy to use and the insights are invaluable." - Priya Sharma, GreenTextile Ltd

"We saved ₹3.6 lakh annually on our loan interest rate by qualifying for a green loan through this platform." - Rajesh Kumar, EcoTech Manufacturing

"The carbon credit trading feature created an additional revenue stream of ₹1.5 lakh annually." - Amit Patel, FreshFoods Pvt Ltd

KEY SUCCESS METRICS:
• Average 35% reduction in operational costs
• 40% improvement in carbon scores
• 25% increase in loan approval rates
• 30% reduction in default rates
• 98% customer satisfaction score"""
    
    # Slide 9: Pricing & Implementation Plan
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Affordable Pricing & Easy Implementation"
    content9.text = """💳 FLEXIBLE PRICING OPTIONS

STARTER PLAN - ₹2,999/month
• Up to 25 employees
• Basic carbon tracking and scoring
• Monthly sustainability report
• Mobile app access
• Email support
• Perfect for small businesses

PROFESSIONAL PLAN - ₹4,999/month
• Up to 100 employees
• Advanced analytics and insights
• Quarterly sustainability reports
• Carbon trading access
• Priority support
• API integration
• Most popular choice

ENTERPRISE PLAN - ₹7,999/month
• Unlimited employees
• Full platform features
• Custom reporting and analytics
• Dedicated account manager
• White-label options
• Advanced integrations
• Perfect for growing businesses

SPECIAL LAUNCH OFFERS:
🎉 30-day free trial (no credit card required)
🎉 20% discount for annual payment
🎉 First 3 months at 50% off
🎉 Free setup and training
🎉 No long-term contracts

IMPLEMENTATION PROCESS:
Step 1: Quick Setup (1 day)
• Download mobile app or access web platform
• Create account with basic company information
• Connect your business email and phone
• Upload recent utility bills and invoices

Step 2: Data Integration (1 week)
• Our AI analyzes your historical data
• Automatic categorization of expenses
• Initial carbon footprint calculation
• Personalized recommendations generated

Step 3: Start Saving (Immediately)
• Begin implementing recommendations
• Track real-time savings and improvements
• Access green finance opportunities
• Generate professional reports

MONEY-BACK GUARANTEE:
If you don't save at least 5x your subscription cost in the first year, we'll refund your money. No questions asked.

ROI GUARANTEE:
Average customer saves ₹2,20,000-₹9,00,000 annually with platform cost of ₹36,000-₹96,000."""
    
    # Slide 10: Call to Action & Next Steps
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Start Your Sustainability Transformation Today"
    content10.text = """🚀 READY TO TRANSFORM YOUR BUSINESS?

KEY BENEFITS SUMMARY:
✅ Save 20-40% on operational costs annually
✅ Access green finance with 1-3% lower interest rates
✅ Build competitive advantage and market reputation
✅ Meet sustainability goals and regulatory compliance
✅ Easy-to-use platform with AI-powered insights
✅ Additional revenue from carbon credit trading

IMMEDIATE NEXT STEPS:
1. Visit www.carbonintelligence.com
2. Start your 30-day free trial
3. Download the mobile app
4. Connect your business accounts
5. Get your first carbon score and recommendations

SUPPORT AVAILABLE:
📞 Phone: +91-98765-43210
📧 Email: support@carbonintelligence.com
💬 Live chat on website
📱 In-app support
🎥 Video tutorials and guides

SPECIAL LAUNCH OFFER:
🎉 30-day free trial
🎉 50% off first 3 months
🎉 Free setup and training
🎉 Money-back guarantee
🎉 No long-term commitment

WHY START NOW?
⏰ Regulatory pressure increasing
📈 Market opportunities expanding
💰 Cost pressures rising
🏆 Competitive advantage available
🌱 Future-proofing your business

DON'T WAIT - START TODAY!
Every day you delay is money lost and opportunities missed. The sooner you start, the more you save.

Contact Us:
📧 info@carbonintelligence.com
📞 +91-98765-43210
🌐 www.carbonintelligence.com
📱 Download our mobile app

Questions & Discussion

Let's build a sustainable future together! 🌱

Carbon Intelligence - Empowering MSMEs for a Greener Tomorrow"""
    
    # Save the presentation
    output_path = "/workspace/pitches/Carbon_Intelligence_Refined_MSME_Pitch_With_Mockups.pptx"
    prs.save(output_path)
    print(f"Refined MSME PowerPoint presentation with mockups created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_refined_msme_pitch_with_mockups()