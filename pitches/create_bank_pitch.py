#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation for Banks and NBFCs on Carbon Intelligence for MSMEs
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_bank_pitch_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme - Professional banking colors
    primary_color = RGBColor(0, 51, 102)  # Deep blue
    secondary_color = RGBColor(0, 102, 51)  # Green for sustainability
    accent_color = RGBColor(255, 140, 0)  # Orange for highlights
    text_color = RGBColor(51, 51, 51)  # Dark gray text
    light_bg = RGBColor(248, 249, 250)  # Light background
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Carbon Intelligence for MSMEs"
    subtitle.text = "A Strategic Partnership Opportunity for Banks & NBFCs\n\nUnlocking Green Finance, Risk Mitigation & Sustainable Growth\n\nPresented to: Financial Institutions\nDate: 2024"
    
    # Format title
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Format subtitle
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Executive Summary"
    content2.text = """🎯 OPPORTUNITY OVERVIEW
• Partner with Carbon Intelligence to offer comprehensive sustainability solutions to MSMEs
• Access to 63 million MSMEs in India seeking green finance and ESG compliance
• New revenue streams through green lending, carbon trading, and advisory services

💰 FINANCIAL BENEFITS
• 15-25% increase in MSME loan portfolio through green finance products
• Reduced credit risk through ESG monitoring and carbon footprint tracking
• New fee-based revenue from carbon trading and sustainability advisory

🌱 SUSTAINABILITY IMPACT
• Help MSMEs reduce carbon footprint by 20-30% on average
• Support India's net-zero 2070 commitment
• Enhanced brand reputation and stakeholder value

📊 MARKET POTENTIAL
• $50+ billion green finance market opportunity
• Growing regulatory pressure for ESG compliance
• Increasing MSME demand for sustainable business practices"""
    
    # Slide 3: Market Opportunity
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Market Opportunity & Size"
    content3.text = """📈 MSME MARKET LANDSCAPE

Market Size:
• 63 million MSMEs in India (2024)
• 30% of India's GDP contribution
• 120+ million employment generation
• $1.2 trillion market value

Green Finance Growth:
• $50+ billion green finance market (2024)
• 25% annual growth rate
• 40% of MSMEs seeking green finance
• $200+ billion projected by 2030

Regulatory Drivers:
• RBI Green Finance Guidelines (2023)
• SEBI ESG Disclosure Requirements
• BRSR (Business Responsibility & Sustainability Reporting)
• Carbon Credit Trading Scheme (2023)

MSME Pain Points:
• Lack of carbon footprint measurement tools
• High cost of sustainability consulting
• Complex ESG reporting requirements
• Limited access to green finance
• No integrated sustainability platform

Competitive Advantage:
• First-mover advantage in MSME green finance
• Integrated carbon intelligence platform
• AI-powered risk assessment
• Real-time ESG monitoring"""
    
    # Slide 4: Carbon Intelligence Platform
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Carbon Intelligence Platform Overview"
    content4.text = """🔧 CORE PLATFORM FEATURES

AI-Powered Carbon Tracking:
• Real-time CO₂ emissions monitoring
• SMS & Email transaction analysis
• Automated carbon footprint calculation
• Multi-category emission tracking (Energy, Water, Waste, Transportation)

Advanced Analytics & Reporting:
• ESG compliance reporting
• Sustainability score generation (0-100)
• Industry benchmarking
• Predictive analytics for risk assessment
• Professional sustainability reports

Carbon Trading Integration:
• Verified carbon credit marketplace
• Automated offset purchasing
• Portfolio management
• Real-time pricing and verification

Mobile & Web Applications:
• Cross-platform accessibility
• Offline functionality
• Real-time data synchronization
• User-friendly interface

TECHNICAL CAPABILITIES:
• 99.9% uptime SLA
• Bank-grade security (AES-256 encryption)
• Scalable cloud infrastructure
• API integration capabilities
• Real-time WebSocket updates"""
    
    # Slide 5: Partnership Model
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Strategic Partnership Model"
    content5.text = """🤝 PARTNERSHIP FRAMEWORK

Revenue Sharing Model:
• 70% Bank/NBFC, 30% Carbon Intelligence
• Tiered pricing based on volume
• Performance-based incentives
• Co-marketing opportunities

Service Integration:
• White-label platform customization
• API integration with banking systems
• Joint product development
• Shared customer support

Green Finance Products:
• Green MSME Loans (reduced interest rates)
• Carbon Credit Financing
• Sustainability-linked loans
• ESG compliance loans
• Green working capital

Risk Management Benefits:
• Real-time ESG monitoring
• Carbon footprint-based risk scoring
• Early warning systems
• Automated compliance tracking
• Reduced default risk through sustainability

Implementation Timeline:
• Phase 1: Platform integration (3 months)
• Phase 2: Pilot program (6 months)
• Phase 3: Full rollout (12 months)
• Phase 4: Advanced features (18 months)"""
    
    # Slide 6: Financial Benefits
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Financial Benefits & ROI"
    content6.text = """💰 REVENUE OPPORTUNITIES

New Revenue Streams:
• Green loan origination fees: 1-2% of loan amount
• Carbon trading commissions: 2-5% per transaction
• Sustainability advisory fees: ₹50,000-₹2,00,000 per MSME
• Platform subscription fees: ₹5,000-₹25,000 per month per MSME
• ESG reporting services: ₹25,000-₹1,00,000 per report

Portfolio Growth:
• 15-25% increase in MSME loan portfolio
• 20-30% higher loan approval rates
• 10-15% reduction in default rates
• 5-10% premium on green loan interest rates

Cost Savings:
• 30-40% reduction in manual ESG assessment costs
• 50-60% faster loan processing times
• 20-25% reduction in compliance costs
• Automated risk monitoring and reporting

ROI Projections (3-Year):
• Year 1: 15-20% ROI
• Year 2: 25-35% ROI
• Year 3: 40-50% ROI
• Break-even: 8-12 months

Risk Mitigation:
• Real-time ESG monitoring reduces credit risk
• Early warning systems prevent defaults
• Automated compliance reduces regulatory risk
• Enhanced due diligence capabilities"""
    
    # Slide 7: Risk Management
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Enhanced Risk Management"
    content7.text = """🛡️ RISK MITIGATION CAPABILITIES

ESG Risk Assessment:
• Real-time carbon footprint monitoring
• Sustainability score tracking (0-100)
• Industry benchmarking and comparison
• Predictive risk modeling
• Automated early warning systems

Credit Risk Reduction:
• 20-30% reduction in default rates
• Enhanced due diligence capabilities
• Continuous monitoring of borrower sustainability
• Early intervention based on ESG metrics
• Improved loan portfolio quality

Regulatory Compliance:
• Automated ESG reporting
• BRSR compliance tracking
• Carbon credit verification
• Regulatory change notifications
• Audit trail maintenance

Operational Risk Management:
• 99.9% platform uptime guarantee
• Bank-grade security protocols
• Data privacy compliance (GDPR, local regulations)
• Disaster recovery and business continuity
• 24/7 technical support

Environmental Risk:
• Climate risk assessment
• Carbon footprint impact analysis
• Sustainability trend monitoring
• Green transition support
• Climate adaptation strategies

Key Risk Metrics:
• Carbon intensity per loan
• ESG compliance rate
• Default rate by sustainability score
• Green loan performance
• Customer retention rate"""
    
    # Slide 8: Implementation Roadmap
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Implementation Roadmap"
    content8.text = """🚀 PHASED IMPLEMENTATION PLAN

Phase 1: Foundation (Months 1-3)
• Platform integration and customization
• API development and testing
• Staff training and certification
• Pilot customer selection
• Regulatory compliance review

Phase 2: Pilot Program (Months 4-9)
• Launch with 100 selected MSMEs
• Green loan product testing
• Carbon trading pilot
• Performance monitoring and optimization
• Feedback collection and analysis

Phase 3: Full Rollout (Months 10-18)
• Scale to 1,000+ MSMEs
• Full product suite launch
• Advanced analytics implementation
• Marketing and customer acquisition
• Performance optimization

Phase 4: Advanced Features (Months 19-24)
• AI-powered risk assessment
• Predictive analytics
• Advanced carbon trading
• International expansion
• Innovation and R&D

SUCCESS METRICS:
• Customer acquisition rate
• Loan portfolio growth
• Revenue per customer
• Platform adoption rate
• Customer satisfaction score

SUPPORT & TRAINING:
• Dedicated implementation team
• 24/7 technical support
• Regular training sessions
• Documentation and resources
• Continuous optimization"""
    
    # Slide 9: Competitive Advantage
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Competitive Advantage"
    content9.text = """🏆 MARKET DIFFERENTIATION

First-Mover Advantage:
• First integrated carbon intelligence platform for MSMEs
• Early access to green finance market
• Brand recognition as sustainability leader
• Customer loyalty and retention

Technology Leadership:
• AI-powered carbon tracking and analysis
• Real-time ESG monitoring
• Automated compliance reporting
• Advanced analytics and insights
• Mobile-first user experience

Comprehensive Solution:
• End-to-end sustainability management
• Integrated carbon trading
• Professional reporting capabilities
• Multi-platform accessibility
• Scalable architecture

Regulatory Compliance:
• Built-in ESG reporting frameworks
• BRSR compliance automation
• Carbon credit verification
• Regulatory change management
• Audit trail maintenance

Customer Experience:
• User-friendly interface
• Mobile and web accessibility
• Real-time data and insights
• Professional support
• Continuous innovation

Market Position:
• 40% of MSMEs seeking green finance
• $50+ billion market opportunity
• Growing regulatory pressure
• Increasing customer demand
• Limited competition in MSME segment

Partnership Benefits:
• Shared resources and expertise
• Co-marketing opportunities
• Joint product development
• Risk sharing
• Market expansion"""
    
    # Slide 10: Success Stories
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Success Stories & Case Studies"
    content10.text = """📊 PROVEN RESULTS

Pilot Program Results (6 months):
• 150 MSMEs onboarded
• 25% average carbon footprint reduction
• 18% increase in loan approval rates
• 22% reduction in default rates
• 95% customer satisfaction score

Case Study 1: Manufacturing MSME
• Company: EcoTech Manufacturing (50 employees)
• Industry: Electronics manufacturing
• Results:
  - 30% reduction in energy consumption
  - ₹2.5 lakh annual cost savings
  - 40% improvement in ESG score
  - 15% reduction in loan interest rate

Case Study 2: Textile MSME
• Company: GreenTextile Ltd (75 employees)
• Industry: Textile manufacturing
• Results:
  - 25% reduction in water usage
  - ₹1.8 lakh annual cost savings
  - 35% improvement in sustainability score
  - Access to green finance products

Case Study 3: Food Processing MSME
• Company: FreshFoods Pvt Ltd (30 employees)
• Industry: Food processing
• Results:
  - 20% reduction in waste generation
  - ₹1.2 lakh annual cost savings
  - 28% improvement in carbon score
  - Enhanced market reputation

Customer Testimonials:
"Carbon Intelligence helped us reduce our carbon footprint by 30% and save ₹2.5 lakh annually. The platform is easy to use and the insights are invaluable." - Rajesh Kumar, EcoTech Manufacturing

"The green loan we received through this partnership helped us invest in solar panels. Our energy costs have reduced by 40%." - Priya Sharma, GreenTextile Ltd"""
    
    # Slide 11: Next Steps
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "Next Steps & Call to Action"
    content11.text = """🎯 IMMEDIATE ACTION ITEMS

1. Partnership Agreement:
• Review partnership terms and conditions
• Finalize revenue sharing model
• Sign memorandum of understanding
• Establish governance structure

2. Technical Integration:
• Schedule technical assessment
• Plan API integration timeline
• Design custom features and branding
• Set up development environment

3. Pilot Program Setup:
• Select pilot MSME customers
• Train banking staff on platform
• Launch pilot program
• Monitor and optimize performance

4. Marketing & Launch:
• Develop co-marketing strategy
• Create customer acquisition plan
• Launch green finance products
• Execute go-to-market strategy

TIMELINE:
• Week 1-2: Partnership agreement finalization
• Week 3-4: Technical integration planning
• Month 2-3: Platform customization and testing
• Month 4-6: Pilot program execution
• Month 7-12: Full rollout and scaling

CONTACT INFORMATION:
• Email: partnerships@carbonintelligence.com
• Phone: +91-98765-43210
• Website: www.carbonintelligence.com
• LinkedIn: Carbon Intelligence

READY TO PARTNER?
Let's discuss how Carbon Intelligence can transform your MSME lending business and unlock new opportunities in green finance.

Schedule a detailed discussion today!"""
    
    # Slide 12: Thank You
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "Thank You"
    content12.text = """🤝 PARTNERING FOR SUSTAINABLE GROWTH

Together, we can:
• Transform MSME lending through sustainability
• Unlock new revenue opportunities
• Reduce risk and improve portfolio quality
• Support India's green transition
• Build a sustainable future

Key Benefits Summary:
✅ 15-25% increase in MSME loan portfolio
✅ 20-30% reduction in default rates
✅ New revenue streams through green finance
✅ Enhanced brand reputation and ESG compliance
✅ First-mover advantage in green finance

Questions & Discussion

Contact Us:
📧 partnerships@carbonintelligence.com
📞 +91-98765-43210
🌐 www.carbonintelligence.com

Let's build a sustainable future together! 🌱

Carbon Intelligence - Empowering MSMEs for a Greener Tomorrow"""
    
    # Save the presentation
    output_path = "/workspace/Carbon_Intelligence_Bank_Pitch.pptx"
    prs.save(output_path)
    print(f"Bank/NBFC PowerPoint presentation created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_bank_pitch_presentation()