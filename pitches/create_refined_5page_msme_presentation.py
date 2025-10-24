#!/usr/bin/env python3
"""
Create Refined 5-Page MSME Carbon Intelligence Presentation
Enhanced with Visual Elements, Charts, Diagrams, and App Mockups
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_refined_5page_msme_presentation():
    """Create a refined 5-page MSME presentation with enhanced visuals"""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define enhanced color scheme
    primary_green = RGBColor(46, 125, 50)  # #2E7D32 - Deep Green
    secondary_orange = RGBColor(255, 140, 0)  # #FF8C00 - Orange
    accent_blue = RGBColor(25, 118, 210)  # #1976D2 - Blue
    success_green = RGBColor(76, 175, 80)  # #4CAF50 - Success Green
    warning_orange = RGBColor(255, 152, 0)  # #FF9800 - Warning Orange
    dark_gray = RGBColor(66, 66, 66)  # #424242 - Dark Gray
    light_gray = RGBColor(245, 245, 245)  # #F5F5F5 - Light Gray
    white = RGBColor(255, 255, 255)  # White
    
    # Slide 1: Enhanced Title Slide with Visual Elements
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative background elements
    # Large circle in background
    bg_circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-1), Inches(6), Inches(6))
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = RGBColor(240, 248, 255)
    bg_circle.line.fill.background()
    
    # Secondary circle
    bg_circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(4), Inches(4))
    bg_circle2.fill.solid()
    bg_circle2.fill.fore_color.rgb = RGBColor(240, 255, 240)
    bg_circle2.line.fill.background()
    
    # Title
    title = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
    title.text = "🌱 Carbon Intelligence"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1))
    subtitle.text = "AI-Powered Carbon Intelligence Platform for MSMEs"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = dark_gray
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Key benefits with icons
    benefits = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(1.5))
    benefits.text = "💰 Save 20-40% on Costs  •  🏦 Access Green Finance  •  🚀 Build Competitive Advantage"
    benefits.text_frame.paragraphs[0].font.size = Pt(18)
    benefits.text_frame.paragraphs[0].font.color.rgb = accent_blue
    benefits.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact = slide1.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.33), Inches(0.8))
    contact.text = "📧 contact@carbonintelligence.com  •  📱 +91-98765-43210  •  🌐 www.carbonintelligence.com"
    contact.text_frame.paragraphs[0].font.size = Pt(14)
    contact.text_frame.paragraphs[0].font.color.rgb = dark_gray
    contact.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem & Solution with Visual Charts
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    title2.text = "The MSME Sustainability Challenge & Our Solution"
    title2.text_frame.paragraphs[0].font.size = Pt(28)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = primary_green
    title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Problem section with visual elements
    problem_left = Inches(0.5)
    problem_top = Inches(1.5)
    problem_width = Inches(6)
    problem_height = Inches(5)
    
    # Problem background
    problem_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, problem_left, problem_top, problem_width, problem_height)
    problem_bg.fill.solid()
    problem_bg.fill.fore_color.rgb = RGBColor(255, 245, 245)
    problem_bg.line.color.rgb = RGBColor(220, 53, 69)
    problem_bg.line.width = Pt(2)
    
    # Problem title
    prob_title = slide2.shapes.add_textbox(problem_left + Inches(0.2), problem_top + Inches(0.2), problem_width - Inches(0.4), Inches(0.6))
    prob_title.text = "🚨 Critical Problems MSMEs Face"
    prob_title.text_frame.paragraphs[0].font.size = Pt(16)
    prob_title.text_frame.paragraphs[0].font.bold = True
    prob_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 53, 69)
    
    # Problem content
    prob_content = slide2.shapes.add_textbox(problem_left + Inches(0.2), problem_top + Inches(0.8), problem_width - Inches(0.4), Inches(3.5))
    prob_content.text = """• High Operational Costs: 60-70% of revenue
• Limited Green Finance Access
• Complex Compliance Requirements
• Lack of Carbon Data & Tracking
• Competitive Pressure
• Regulatory Risk & Penalties

💰 Financial Impact:
• ₹2-5 lakhs annual losses
• 40% higher operational costs
• Missed green finance opportunities
• ₹50,000 crores market opportunity"""
    prob_content.text_frame.paragraphs[0].font.size = Pt(12)
    prob_content.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Solution section with visual elements
    solution_left = Inches(6.5)
    solution_top = Inches(1.5)
    solution_width = Inches(6)
    solution_height = Inches(5)
    
    # Solution background
    solution_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, solution_left, solution_top, solution_width, solution_height)
    solution_bg.fill.solid()
    solution_bg.fill.fore_color.rgb = RGBColor(245, 255, 245)
    solution_bg.line.color.rgb = success_green
    solution_bg.line.width = Pt(2)
    
    # Solution title
    sol_title = slide2.shapes.add_textbox(solution_left + Inches(0.2), solution_top + Inches(0.2), solution_width - Inches(0.4), Inches(0.6))
    sol_title.text = "✅ Carbon Intelligence Solution"
    sol_title.text_frame.paragraphs[0].font.size = Pt(16)
    sol_title.text_frame.paragraphs[0].font.bold = True
    sol_title.text_frame.paragraphs[0].font.color.rgb = success_green
    
    # Solution content
    sol_content = slide2.shapes.add_textbox(solution_left + Inches(0.2), solution_top + Inches(0.8), solution_width - Inches(0.4), Inches(3.5))
    sol_content.text = """🤖 AI-Powered Carbon Intelligence
• Real-time carbon tracking
• Smart cost optimization
• Green finance access
• Automated compliance
• Carbon trading platform

📱 Multi-Platform Access:
• Web Dashboard
• Mobile App
• API Integration

🎯 Key Benefits:
• 20-40% cost reduction
• 1-3% lower interest rates
• Additional revenue streams
• Future-proof operations"""
    sol_content.text_frame.paragraphs[0].font.size = Pt(12)
    sol_content.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Slide 3: App Mockups & Interface Demo
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    title3.text = "📱 Application Interface Mockups & Key Features"
    title3.text_frame.paragraphs[0].font.size = Pt(28)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = primary_green
    title3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Dashboard mockup
    dashboard_left = Inches(0.5)
    dashboard_top = Inches(1.5)
    dashboard_width = Inches(5.5)
    dashboard_height = Inches(4)
    
    # Dashboard background
    dashboard_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dashboard_left, dashboard_top, dashboard_width, dashboard_height)
    dashboard_bg.fill.solid()
    dashboard_bg.fill.fore_color.rgb = light_gray
    dashboard_bg.line.color.rgb = primary_green
    dashboard_bg.line.width = Pt(2)
    
    # Dashboard header
    dash_header = slide3.shapes.add_textbox(dashboard_left + Inches(0.2), dashboard_top + Inches(0.2), dashboard_width - Inches(0.4), Inches(0.6))
    dash_header.text = "🌱 Carbon Intelligence Dashboard"
    dash_header.text_frame.paragraphs[0].font.size = Pt(14)
    dash_header.text_frame.paragraphs[0].font.bold = True
    dash_header.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    # Carbon score display
    score_bg = slide3.shapes.add_shape(MSO_SHAPE.OVAL, dashboard_left + Inches(0.3), dashboard_top + Inches(1), Inches(1.5), Inches(1.5))
    score_bg.fill.solid()
    score_bg.fill.fore_color.rgb = success_green
    
    score_text = slide3.shapes.add_textbox(dashboard_left + Inches(0.35), dashboard_top + Inches(1.1), Inches(1.4), Inches(1.3))
    score_text.text = "85\nGold\nTier"
    score_text.text_frame.paragraphs[0].font.size = Pt(14)
    score_text.text_frame.paragraphs[0].font.bold = True
    score_text.text_frame.paragraphs[0].font.color.rgb = white
    score_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Metrics boxes
    metrics = [
        ("CO₂ Saved", "245.6 kg", success_green),
        ("Cost Savings", "₹15,240", accent_blue),
        ("Reduction", "12.3%", secondary_orange),
        ("Credits", "1,250", primary_green)
    ]
    
    for i, (label, value, color) in enumerate(metrics):
        metric_left = dashboard_left + Inches(2) + (i % 2) * Inches(1.5)
        metric_top = dashboard_top + Inches(1) + (i // 2) * Inches(0.8)
        
        metric_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, metric_left, metric_top, Inches(1.3), Inches(0.6))
        metric_bg.fill.solid()
        metric_bg.fill.fore_color.rgb = white
        metric_bg.line.color.rgb = color
        metric_bg.line.width = Pt(1)
        
        metric_text = slide3.shapes.add_textbox(metric_left + Inches(0.1), metric_top + Inches(0.1), Inches(1.1), Inches(0.4))
        metric_text.text = f"{value}\n{label}"
        metric_text.text_frame.paragraphs[0].font.size = Pt(9)
        metric_text.text_frame.paragraphs[0].font.bold = True
        metric_text.text_frame.paragraphs[0].font.color.rgb = color
        metric_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Mobile app mockup
    mobile_left = Inches(6.5)
    mobile_top = Inches(1.5)
    mobile_width = Inches(2.5)
    mobile_height = Inches(4)
    
    # Mobile background
    mobile_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mobile_left, mobile_top, mobile_width, mobile_height)
    mobile_bg.fill.solid()
    mobile_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
    mobile_bg.line.color.rgb = dark_gray
    mobile_bg.line.width = Pt(3)
    
    # Mobile screen
    mobile_screen = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mobile_left + Inches(0.1), mobile_top + Inches(0.2), mobile_width - Inches(0.2), mobile_height - Inches(0.4))
    mobile_screen.fill.solid()
    mobile_screen.fill.fore_color.rgb = light_gray
    
    # Mobile content
    mobile_content = slide3.shapes.add_textbox(mobile_left + Inches(0.2), mobile_top + Inches(0.3), mobile_width - Inches(0.4), Inches(3.5))
    mobile_content.text = """🌱 Carbon Intelligence

📊 Dashboard
• Score: 85/100
• CO₂: 245 kg
• Savings: ₹15,240

🔍 Quick Actions
• Analytics
• Trading
• Finance
• Reports

📱 24/7 Access
• Push Notifications
• Offline Mode
• GPS Tracking
• Voice Commands"""
    mobile_content.text_frame.paragraphs[0].font.size = Pt(8)
    mobile_content.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Features list
    features_left = Inches(0.5)
    features_top = Inches(5.8)
    features_width = Inches(12.33)
    features_height = Inches(1.2)
    
    features = slide3.shapes.add_textbox(features_left, features_top, features_width, features_height)
    features.text = "🚀 Key Features: Real-time Carbon Tracking • AI-Powered Analytics • Green Finance Access • Carbon Trading • Mobile App • Automated Compliance • Multi-Platform Sync"
    features.text_frame.paragraphs[0].font.size = Pt(14)
    features.text_frame.paragraphs[0].font.color.rgb = accent_blue
    features.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 4: Financial Impact & ROI with Visual Charts
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    title4.text = "💰 Financial Impact & ROI Analysis"
    title4.text_frame.paragraphs[0].font.size = Pt(28)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = primary_green
    title4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ROI Chart
    chart_left = Inches(0.5)
    chart_top = Inches(1.5)
    chart_width = Inches(6)
    chart_height = Inches(4)
    
    # Chart background
    chart_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, chart_left, chart_top, chart_width, chart_height)
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
    chart_bg.line.color.rgb = accent_blue
    chart_bg.line.width = Pt(2)
    
    # Chart title
    chart_title = slide4.shapes.add_textbox(chart_left + Inches(0.2), chart_top + Inches(0.2), chart_width - Inches(0.4), Inches(0.6))
    chart_title.text = "📊 ROI Progression (3-Year Analysis)"
    chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart_title.text_frame.paragraphs[0].font.bold = True
    chart_title.text_frame.paragraphs[0].font.color.rgb = accent_blue
    
    # ROI bars
    years = ["Year 1", "Year 2", "Year 3"]
    roi_values = [400, 500, 600]
    colors = [success_green, accent_blue, primary_green]
    
    for i, (year, roi, color) in enumerate(zip(years, roi_values, colors)):
        bar_left = chart_left + Inches(0.5) + i * Inches(1.8)
        bar_top = chart_top + Inches(1.2)
        bar_width = Inches(1.2)
        bar_height = Inches(2) * (roi / 600)  # Scale to max 600%
        
        # Bar background
        bar_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top + Inches(2) - bar_height, bar_width, bar_height)
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = color
        bar_bg.line.fill.background()
        
        # Bar label
        bar_label = slide4.shapes.add_textbox(bar_left, bar_top + Inches(2.2), bar_width, Inches(0.4))
        bar_label.text = f"{year}\n{roi}% ROI"
        bar_label.text_frame.paragraphs[0].font.size = Pt(10)
        bar_label.text_frame.paragraphs[0].font.bold = True
        bar_label.text_frame.paragraphs[0].font.color.rgb = dark_gray
        bar_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Cost savings breakdown
    savings_left = Inches(7)
    savings_top = Inches(1.5)
    savings_width = Inches(5.5)
    savings_height = Inches(4)
    
    # Savings background
    savings_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, savings_left, savings_top, savings_width, savings_height)
    savings_bg.fill.solid()
    savings_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
    savings_bg.line.color.rgb = success_green
    savings_bg.line.width = Pt(2)
    
    # Savings title
    savings_title = slide4.shapes.add_textbox(savings_left + Inches(0.2), savings_top + Inches(0.2), savings_width - Inches(0.4), Inches(0.6))
    savings_title.text = "💵 Annual Cost Savings Breakdown"
    savings_title.text_frame.paragraphs[0].font.size = Pt(16)
    savings_title.text_frame.paragraphs[0].font.bold = True
    savings_title.text_frame.paragraphs[0].font.color.rgb = success_green
    
    # Savings categories
    categories = [
        ("⚡ Energy", "25-40%", "₹75K-₹3L", success_green),
        ("💧 Water", "20-30%", "₹40K-₹1.5L", accent_blue),
        ("🗑️ Waste", "35-50%", "₹25K-₹1L", secondary_orange),
        ("🚚 Transport", "15-25%", "₹30K-₹1.5L", primary_green),
        ("📦 Materials", "20-30%", "₹50K-₹2L", warning_orange)
    ]
    
    for i, (category, reduction, savings, color) in enumerate(categories):
        cat_left = savings_left + Inches(0.2)
        cat_top = savings_top + Inches(0.9) + i * Inches(0.6)
        cat_width = savings_width - Inches(0.4)
        cat_height = Inches(0.5)
        
        # Category background
        cat_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cat_left, cat_top, cat_width, cat_height)
        cat_bg.fill.solid()
        cat_bg.fill.fore_color.rgb = white
        cat_bg.line.color.rgb = color
        cat_bg.line.width = Pt(1)
        
        # Category text
        cat_text = slide4.shapes.add_textbox(cat_left + Inches(0.1), cat_top + Inches(0.1), cat_width - Inches(0.2), Inches(0.3))
        cat_text.text = f"{category} {reduction} reduction - {savings} annually"
        cat_text.text_frame.paragraphs[0].font.size = Pt(10)
        cat_text.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Total savings
    total_left = savings_left + Inches(0.2)
    total_top = savings_top + Inches(3.5)
    total_width = savings_width - Inches(0.4)
    total_height = Inches(0.4)
    
    total_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, total_left, total_top, total_width, total_height)
    total_bg.fill.solid()
    total_bg.fill.fore_color.rgb = success_green
    
    total_text = slide4.shapes.add_textbox(total_left + Inches(0.1), total_top + Inches(0.1), total_width - Inches(0.2), Inches(0.2))
    total_text.text = "💰 TOTAL ANNUAL SAVINGS: ₹2.2L - ₹9L per MSME"
    total_text.text_frame.paragraphs[0].font.size = Pt(12)
    total_text.text_frame.paragraphs[0].font.bold = True
    total_text.text_frame.paragraphs[0].font.color.rgb = white
    total_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 5: Success Stories & Call to Action
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    title5.text = "🏆 Success Stories & Ready to Transform?"
    title5.text_frame.paragraphs[0].font.size = Pt(28)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = primary_green
    title5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Success stories section
    stories_left = Inches(0.5)
    stories_top = Inches(1.5)
    stories_width = Inches(6)
    stories_height = Inches(4)
    
    # Stories background
    stories_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, stories_left, stories_top, stories_width, stories_height)
    stories_bg.fill.solid()
    stories_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
    stories_bg.line.color.rgb = accent_blue
    stories_bg.line.width = Pt(2)
    
    # Stories title
    stories_title = slide5.shapes.add_textbox(stories_left + Inches(0.2), stories_top + Inches(0.2), stories_width - Inches(0.4), Inches(0.6))
    stories_title.text = "📊 Real Customer Results"
    stories_title.text_frame.paragraphs[0].font.size = Pt(16)
    stories_title.text_frame.paragraphs[0].font.bold = True
    stories_title.text_frame.paragraphs[0].font.color.rgb = accent_blue
    
    # Case studies
    case_studies = [
        ("EcoTech Manufacturing", "75 employees", "40% cost reduction", "₹6L annual savings", success_green),
        ("GreenTextile Ltd", "100 employees", "35% cost reduction", "₹4.2L annual savings", accent_blue),
        ("FreshFoods Pvt Ltd", "50 employees", "30% cost reduction", "₹2.5L annual savings", secondary_orange)
    ]
    
    for i, (company, employees, reduction, savings, color) in enumerate(case_studies):
        case_left = stories_left + Inches(0.2)
        case_top = stories_top + Inches(0.9) + i * Inches(1)
        case_width = stories_width - Inches(0.4)
        case_height = Inches(0.8)
        
        # Case background
        case_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, case_left, case_top, case_width, case_height)
        case_bg.fill.solid()
        case_bg.fill.fore_color.rgb = white
        case_bg.line.color.rgb = color
        case_bg.line.width = Pt(1)
        
        # Case text
        case_text = slide5.shapes.add_textbox(case_left + Inches(0.1), case_top + Inches(0.1), case_width - Inches(0.2), Inches(0.6))
        case_text.text = f"🏢 {company} ({employees})\n📈 {reduction} • 💰 {savings}"
        case_text.text_frame.paragraphs[0].font.size = Pt(10)
        case_text.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Call to action section
    cta_left = Inches(7)
    cta_top = Inches(1.5)
    cta_width = Inches(5.5)
    cta_height = Inches(4)
    
    # CTA background
    cta_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cta_left, cta_top, cta_width, cta_height)
    cta_bg.fill.solid()
    cta_bg.fill.fore_color.rgb = RGBColor(240, 255, 240)
    cta_bg.line.color.rgb = success_green
    cta_bg.line.width = Pt(2)
    
    # CTA title
    cta_title = slide5.shapes.add_textbox(cta_left + Inches(0.2), cta_top + Inches(0.2), cta_width - Inches(0.4), Inches(0.6))
    cta_title.text = "🚀 Ready to Transform?"
    cta_title.text_frame.paragraphs[0].font.size = Pt(16)
    cta_title.text_frame.paragraphs[0].font.bold = True
    cta_title.text_frame.paragraphs[0].font.color.rgb = success_green
    
    # CTA content
    cta_content = slide5.shapes.add_textbox(cta_left + Inches(0.2), cta_top + Inches(0.9), cta_width - Inches(0.4), Inches(2.5))
    cta_content.text = """✅ Immediate Benefits:
• Start saving from day 1
• Access green finance
• Build competitive advantage
• Future-proof operations

🎁 Special Launch Offer:
• 30-day free trial
• 50% off first 3 months
• Free setup & training
• Money-back guarantee

📞 Get Started Today:
• Visit: www.carbonintelligence.com
• Call: +91-98765-43210
• Email: contact@carbonintelligence.com

💰 ROI Guarantee:
5x subscription cost savings
or your money back!"""
    cta_content.text_frame.paragraphs[0].font.size = Pt(11)
    cta_content.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Bottom action bar
    action_left = Inches(0.5)
    action_top = Inches(5.8)
    action_width = Inches(12.33)
    action_height = Inches(1.2)
    
    action_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, action_left, action_top, action_width, action_height)
    action_bg.fill.solid()
    action_bg.fill.fore_color.rgb = primary_green
    
    action_text = slide5.shapes.add_textbox(action_left + Inches(0.2), action_top + Inches(0.2), action_width - Inches(0.4), Inches(0.8))
    action_text.text = "🌱 Transform Your MSME Today! Join 500+ MSMEs already saving money and reducing their carbon footprint with Carbon Intelligence."
    action_text.text_frame.paragraphs[0].font.size = Pt(16)
    action_text.text_frame.paragraphs[0].font.bold = True
    action_text.text_frame.paragraphs[0].font.color.rgb = white
    action_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "/workspace/pitches/Carbon_Intelligence_Refined_5Page_MSME_Presentation.pptx"
    prs.save(output_path)
    
    print(f"✅ Refined 5-page MSME presentation created successfully!")
    print(f"📁 File saved: {output_path}")
    print(f"📊 Slides: 5 focused slides with enhanced visuals")
    print(f"🎯 Target: Micro, Small & Medium Enterprises")
    print(f"📱 Features: App mockups, charts, diagrams, ROI analysis")
    print(f"🎨 Visuals: Enhanced with colors, shapes, and visual elements")
    
    return output_path

if __name__ == "__main__":
    create_refined_5page_msme_presentation()