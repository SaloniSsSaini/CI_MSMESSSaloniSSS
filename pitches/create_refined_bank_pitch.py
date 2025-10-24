#!/usr/bin/env python3
"""
Create a refined PowerPoint presentation for Carbon Intelligence Bank Pitch
5 focused slides with enhanced visuals and content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_refined_bank_pitch():
    """Create a refined 5-slide presentation for banks"""
    
    # Create presentation
    prs = Presentation()
    
    # Define color scheme
    primary_green = RGBColor(46, 125, 50)  # #2E7D32
    secondary_orange = RGBColor(255, 111, 0)  # #FF6F00
    accent_blue = RGBColor(25, 118, 210)  # #1976D2
    dark_gray = RGBColor(33, 33, 33)  # #212121
    light_gray = RGBColor(158, 158, 158)  # #9E9E9E
    
    # Slide 1: Title Slide
    create_title_slide(prs, primary_green, secondary_orange, accent_blue)
    
    # Slide 2: Problem & Opportunity
    create_problem_opportunity_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray)
    
    # Slide 3: Solution Overview
    create_solution_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray)
    
    # Slide 4: Key Features & Benefits
    create_features_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray)
    
    # Slide 5: Call to Action & Next Steps
    create_cta_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray)
    
    # Save presentation
    output_path = "/workspace/pitches/RefinedPitch4Banks.pptx"
    prs.save(output_path)
    print(f"Refined presentation saved to: {output_path}")
    
    return output_path

def create_title_slide(prs, primary_green, secondary_orange, accent_blue):
    """Create title slide"""
    dark_gray = RGBColor(33, 33, 33)  # #212121
    light_gray = RGBColor(158, 158, 158)  # #9E9E9E
    
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Carbon Intelligence"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].font.bold = True
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "AI-Powered Green Finance Platform for Banks\n\nTransforming MSME Sustainability into Profitable Banking Opportunities"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = dark_gray
    subtitle.text_frame.paragraphs[1].font.size = Pt(18)
    subtitle.text_frame.paragraphs[1].font.color.rgb = light_gray
    subtitle.text_frame.paragraphs[2].font.size = Pt(18)
    subtitle.text_frame.paragraphs[2].font.color.rgb = light_gray
    
    # Add decorative elements
    add_decorative_elements(slide, primary_green, secondary_orange, accent_blue)

def create_problem_opportunity_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray):
    """Create problem and opportunity slide"""
    slide_layout = prs.slide_layouts[1]  # Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "The Green Finance Opportunity"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Problem section
    p = tf.add_paragraph()
    p.text = "🚨 Current Challenges"
    p.font.size = Pt(24)
    p.font.color.rgb = secondary_orange
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• MSMEs lack access to green finance (only 2% of total lending)"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Manual ESG assessment is costly and time-consuming"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• No real-time carbon footprint tracking for loan decisions"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Regulatory pressure for sustainable finance is increasing"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    # Opportunity section
    p = tf.add_paragraph()
    p.text = "\n💰 Market Opportunity"
    p.font.size = Pt(24)
    p.font.color.rgb = accent_blue
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• $2.5T green finance market by 2025"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• 63M MSMEs in India need sustainable finance"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• 40% lower default rates for green loans"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• New revenue streams through carbon credit trading"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray

def create_solution_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray):
    """Create solution overview slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Carbon Intelligence Platform"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "🤖 AI-Powered Green Finance Solution"
    p.font.size = Pt(24)
    p.font.color.rgb = accent_blue
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Real-time carbon footprint scoring (0-100 scale)"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Automated ESG compliance monitoring"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Smart transaction analysis via SMS/Email"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Predictive risk assessment using ML"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "\n💼 Banking Integration"
    p.font.size = Pt(24)
    p.font.color.rgb = secondary_orange
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Seamless API integration with core banking systems"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Automated loan eligibility assessment"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Dynamic interest rate adjustment based on sustainability"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Carbon credit trading marketplace integration"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray

def create_features_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray):
    """Create key features and benefits slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Key Features & Business Benefits"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Features section
    p = tf.add_paragraph()
    p.text = "🔧 Core Features"
    p.font.size = Pt(24)
    p.font.color.rgb = accent_blue
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• AI Multi-Agent System for comprehensive analysis"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Real-time carbon scoring and ESG compliance"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Automated transaction categorization and analysis"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Professional ESG reporting and documentation"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    # Benefits section
    p = tf.add_paragraph()
    p.text = "\n📈 Business Benefits"
    p.font.size = Pt(24)
    p.font.color.rgb = secondary_orange
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 40% reduction in loan processing time"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• 25% increase in green loan portfolio"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• 30% lower operational costs for ESG assessment"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• New revenue streams through carbon credit trading"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray

def create_cta_slide(prs, primary_green, secondary_orange, accent_blue, dark_gray):
    """Create call to action slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Ready to Transform Your Green Finance?"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "🚀 Next Steps"
    p.font.size = Pt(24)
    p.font.color.rgb = accent_blue
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "1. Pilot Program: 3-month trial with select MSMEs"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "2. Integration: Seamless API integration with your core banking"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "3. Training: Comprehensive team training and support"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "4. Launch: Full-scale deployment with ongoing support"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "\n💡 Why Choose Carbon Intelligence?"
    p.font.size = Pt(24)
    p.font.color.rgb = secondary_orange
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Proven AI technology with 95% accuracy"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Scalable platform supporting 100K+ MSMEs"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Dedicated support team and regular updates"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• ROI within 6 months of implementation"
    p.font.size = Pt(18)
    p.font.color.rgb = dark_gray

def add_decorative_elements(slide, primary_green, secondary_orange, accent_blue):
    """Add decorative elements to slides"""
    # Add some geometric shapes for visual appeal
    left = Inches(0.5)
    top = Inches(6)
    width = Inches(1)
    height = Inches(0.5)
    
    # Add colored rectangles
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = primary_green
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(1.2), top, width, height)
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = secondary_orange
    
    shape3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(2.4), top, width, height)
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = accent_blue

if __name__ == "__main__":
    create_refined_bank_pitch()