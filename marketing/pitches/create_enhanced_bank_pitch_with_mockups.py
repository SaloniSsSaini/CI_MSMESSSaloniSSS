#!/usr/bin/env python3
"""
Enhanced Script to create a PowerPoint presentation for Banks with Application Mockups
Focus on visual mockups of banking interface, loan management, and carbon intelligence features
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_enhanced_bank_pitch_with_mockups():
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme - Professional banking colors with sustainability theme
    primary_color = RGBColor(46, 125, 50)  # Green (#2E7D32)
    secondary_color = RGBColor(255, 111, 0)  # Orange (#FF6F00)
    accent_color = RGBColor(25, 118, 210)  # Blue (#1976D2)
    text_color = RGBColor(51, 51, 51)  # Dark gray text
    light_bg = RGBColor(248, 249, 250)  # Light background
    success_color = RGBColor(76, 175, 80)  # Success green
    warning_color = RGBColor(255, 152, 0)  # Warning orange
    white = RGBColor(255, 255, 255)
    dark_gray = RGBColor(64, 64, 64)
    
    # Slide 1: Title Slide with Visual Elements
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    # Add decorative elements
    # Green circle for visual appeal
    left = Inches(1)
    top = Inches(1)
    width = Inches(0.5)
    height = Inches(0.5)
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    circle.fill.solid()
    circle.fill.fore_color.rgb = primary_color
    circle.line.fill.background()
    
    # Orange circle
    left2 = Inches(8.5)
    top2 = Inches(1)
    circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, left2, top2, width, height)
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = secondary_color
    circle2.line.fill.background()
    
    # Blue circle
    left3 = Inches(4.5)
    top3 = Inches(6)
    circle3 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, left3, top3, width, height)
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = accent_color
    circle3.line.fill.background()
    
    title.text = "Carbon Intelligence for Green Finance"
    subtitle.text = "Revolutionary AI-Powered Risk Assessment & Green Loan Solutions\n\nTransform MSME Sustainability into Profitable Banking Opportunities\n\nPresented to: Financial Institutions | 2024"
    
    # Format title
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Format subtitle
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.bold = True
    
    # Slide 2: Banking Dashboard Mockup
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Banking Dashboard - Carbon Intelligence Integration"
    
    # Create banking dashboard mockup
    # Main dashboard container
    dashboard_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(4))
    dashboard_bg.fill.solid()
    dashboard_bg.fill.fore_color.rgb = RGBColor(240, 240, 240)
    dashboard_bg.line.color.rgb = dark_gray
    dashboard_bg.line.width = Pt(2)
    
    # Header bar
    header = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = primary_color
    header.line.fill.background()
    
    # Header text
    header_text = slide2.shapes.add_textbox(Inches(6.1), Inches(2.1), Inches(3.3), Inches(0.3))
    header_text.text_frame.text = "Carbon Intelligence Banking Dashboard"
    header_text.text_frame.paragraphs[0].font.size = Pt(14)
    header_text.text_frame.paragraphs[0].font.bold = True
    header_text.text_frame.paragraphs[0].font.color.rgb = white
    header_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Dashboard metrics boxes
    # Total Loans
    metric1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(2.7), Inches(1.5), Inches(0.8))
    metric1.fill.solid()
    metric1.fill.fore_color.rgb = white
    metric1.line.color.rgb = RGBColor(200, 200, 200)
    metric1.line.width = Pt(1)
    
    metric1_text = slide2.shapes.add_textbox(Inches(6.2), Inches(2.8), Inches(1.3), Inches(0.6))
    metric1_text.text_frame.text = "Total Green Loans\n₹2,450 Cr"
    metric1_text.text_frame.paragraphs[0].font.size = Pt(10)
    metric1_text.text_frame.paragraphs[0].font.bold = True
    metric1_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Active MSMEs
    metric2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(2.7), Inches(1.5), Inches(0.8))
    metric2.fill.solid()
    metric2.fill.fore_color.rgb = white
    metric2.line.color.rgb = RGBColor(200, 200, 200)
    metric2.line.width = Pt(1)
    
    metric2_text = slide2.shapes.add_textbox(Inches(7.8), Inches(2.8), Inches(1.3), Inches(0.6))
    metric2_text.text_frame.text = "Active MSMEs\n1,247"
    metric2_text.text_frame.paragraphs[0].font.size = Pt(10)
    metric2_text.text_frame.paragraphs[0].font.bold = True
    metric2_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Average Carbon Score
    metric3 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(3.6), Inches(1.5), Inches(0.8))
    metric3.fill.solid()
    metric3.fill.fore_color.rgb = white
    metric3.line.color.rgb = RGBColor(200, 200, 200)
    metric3.line.width = Pt(1)
    
    metric3_text = slide2.shapes.add_textbox(Inches(6.2), Inches(3.7), Inches(1.3), Inches(0.6))
    metric3_text.text_frame.text = "Avg Carbon Score\n78.5"
    metric3_text.text_frame.paragraphs[0].font.size = Pt(10)
    metric3_text.text_frame.paragraphs[0].font.bold = True
    metric3_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Default Rate
    metric4 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(3.6), Inches(1.5), Inches(0.8))
    metric4.fill.solid()
    metric4.fill.fore_color.rgb = white
    metric4.line.color.rgb = RGBColor(200, 200, 200)
    metric4.line.width = Pt(1)
    
    metric4_text = slide2.shapes.add_textbox(Inches(7.8), Inches(3.7), Inches(1.3), Inches(0.6))
    metric4_text.text_frame.text = "Default Rate\n2.3%"
    metric4_text.text_frame.paragraphs[0].font.size = Pt(10)
    metric4_text.text_frame.paragraphs[0].font.bold = True
    metric4_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Recent Applications
    recent_header = slide2.shapes.add_textbox(Inches(6.1), Inches(4.5), Inches(3.3), Inches(0.3))
    recent_header.text_frame.text = "Recent Loan Applications"
    recent_header.text_frame.paragraphs[0].font.size = Pt(12)
    recent_header.text_frame.paragraphs[0].font.bold = True
    recent_header.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Application list
    app1 = slide2.shapes.add_textbox(Inches(6.1), Inches(4.8), Inches(3.3), Inches(0.4))
    app1.text_frame.text = "EcoTech Manufacturing - ₹50L - Score: 85 - Approved"
    app1.text_frame.paragraphs[0].font.size = Pt(9)
    app1.text_frame.paragraphs[0].font.color.rgb = success_color
    
    app2 = slide2.shapes.add_textbox(Inches(6.1), Inches(5.2), Inches(3.3), Inches(0.4))
    app2.text_frame.text = "GreenTextile Ltd - ₹25L - Score: 78 - Under Review"
    app2.text_frame.paragraphs[0].font.size = Pt(9)
    app2.text_frame.paragraphs[0].font.color.rgb = warning_color
    
    content2.text = """🏦 INTEGRATED BANKING DASHBOARD

KEY FEATURES:
• Real-time loan portfolio monitoring
• Carbon Intelligence scoring integration
• Automated risk assessment alerts
• Green loan performance analytics
• MSME sustainability tracking

DASHBOARD METRICS:
• Total Green Loans: ₹2,450 crores
• Active MSMEs: 1,247 customers
• Average Carbon Score: 78.5
• Default Rate: 2.3% (vs 4.1% standard)

AUTOMATED WORKFLOWS:
• Instant loan eligibility assessment
• Real-time carbon score updates
• Automated approval recommendations
• Risk monitoring and alerts
• ESG compliance tracking

BENEFITS FOR BANKS:
• 60% faster loan processing
• 45% reduction in default rates
• Real-time portfolio monitoring
• Automated compliance reporting
• Enhanced customer insights"""
    
    # Slide 3: Carbon Intelligence Scoring Interface
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Carbon Intelligence Scoring Interface"
    
    # Create scoring interface mockup
    # Main interface container
    interface_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(4))
    interface_bg.fill.solid()
    interface_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    interface_bg.line.color.rgb = dark_gray
    interface_bg.line.width = Pt(2)
    
    # Company header
    company_header = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(0.6))
    company_header.fill.solid()
    company_header.fill.fore_color.rgb = primary_color
    company_header.line.fill.background()
    
    company_text = slide3.shapes.add_textbox(Inches(6.1), Inches(2.1), Inches(3.3), Inches(0.4))
    company_text.text_frame.text = "EcoTech Manufacturing Pvt Ltd"
    company_text.text_frame.paragraphs[0].font.size = Pt(14)
    company_text.text_frame.paragraphs[0].font.bold = True
    company_text.text_frame.paragraphs[0].font.color.rgb = white
    company_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Carbon Score Circle
    score_circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(2.8), Inches(1.6), Inches(1.6))
    score_circle.fill.solid()
    score_circle.fill.fore_color.rgb = white
    score_circle.line.color.rgb = primary_color
    score_circle.line.width = Pt(4)
    
    score_text = slide3.shapes.add_textbox(Inches(7.3), Inches(3.2), Inches(1.4), Inches(0.8))
    score_text.text_frame.text = "85\nScore"
    score_text.text_frame.paragraphs[0].font.size = Pt(16)
    score_text.text_frame.paragraphs[0].font.bold = True
    score_text.text_frame.paragraphs[0].font.color.rgb = primary_color
    score_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Score breakdown
    breakdown_text = slide3.shapes.add_textbox(Inches(6.1), Inches(4.5), Inches(3.3), Inches(1.3))
    breakdown_text.text_frame.text = """SCORE BREAKDOWN:
• Energy Efficiency: 22/25 (88%)
• Water Management: 16/20 (80%)
• Waste Management: 18/20 (90%)
• Transportation: 12/15 (80%)
• Materials: 8/10 (80%)
• ESG Compliance: 9/10 (90%)"""
    breakdown_text.text_frame.paragraphs[0].font.size = Pt(9)
    breakdown_text.text_frame.paragraphs[0].font.color.rgb = text_color
    
    content3.text = """🧠 CARBON INTELLIGENCE SCORING INTERFACE

REAL-TIME SCORING:
• Live carbon footprint assessment
• AI-powered transaction analysis
• Automated data collection from SMS/Email
• Continuous monitoring and updates
• Predictive analytics for future performance

SCORING METHODOLOGY (0-100):
• Energy Efficiency (25%): Renewable energy usage
• Water Management (20%): Conservation practices
• Waste Management (20%): Circular economy
• Transportation (15%): Green logistics
• Materials & Supply Chain (10%): Sustainable sourcing
• ESG Compliance (10%): Regulatory adherence

BANKING INTEGRATION:
• Automatic loan eligibility assessment
• Dynamic interest rate calculation
• Risk scoring and monitoring
• Compliance reporting automation
• Portfolio performance tracking

BENEFITS:
• 85% accuracy in sustainability prediction
• 40% reduction in manual assessment time
• Real-time risk monitoring
• Automated compliance reporting
• Industry benchmarking and comparison"""
    
    # Slide 4: Automated Loan Approval Workflow
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Automated Loan Approval Workflow"
    
    # Create workflow mockup
    # Workflow container
    workflow_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(4))
    workflow_bg.fill.solid()
    workflow_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
    workflow_bg.line.color.rgb = dark_gray
    workflow_bg.line.width = Pt(2)
    
    # Workflow steps
    steps = [
        ("Application\nReceived", primary_color, Inches(6.2), Inches(2.3)),
        ("Carbon Score\nAnalysis", secondary_color, Inches(7.5), Inches(2.3)),
        ("Risk Assessment", accent_color, Inches(8.8), Inches(2.3)),
        ("Auto Approval\nDecision", success_color, Inches(6.2), Inches(3.8)),
        ("Loan Disbursement", primary_color, Inches(7.5), Inches(3.8))
    ]
    
    for i, (text, color, x, y) in enumerate(steps):
        # Step box
        step_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.2), Inches(1.2))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = color
        step_box.line.color.rgb = dark_gray
        step_box.line.width = Pt(2)
        
        # Step text
        step_text = slide4.shapes.add_textbox(x + Inches(0.1), y + Inches(0.3), Inches(1), Inches(0.6))
        step_text.text_frame.text = text
        step_text.text_frame.paragraphs[0].font.size = Pt(9)
        step_text.text_frame.paragraphs[0].font.bold = True
        step_text.text_frame.paragraphs[0].font.color.rgb = white
        step_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Step number
        step_num = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), y - Inches(0.2), Inches(0.4), Inches(0.4))
        step_num.fill.solid()
        step_num.fill.fore_color.rgb = white
        step_num.line.color.rgb = color
        step_num.line.width = Pt(2)
        
        num_text = slide4.shapes.add_textbox(x + Inches(0.45), y - Inches(0.15), Inches(0.3), Inches(0.3))
        num_text.text_frame.text = str(i + 1)
        num_text.text_frame.paragraphs[0].font.size = Pt(10)
        num_text.text_frame.paragraphs[0].font.bold = True
        num_text.text_frame.paragraphs[0].font.color.rgb = color
        num_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrows between steps
    arrows = [
        (Inches(7.4), Inches(2.9), Inches(7.3), Inches(2.9)),  # Step 1 to 2
        (Inches(8.7), Inches(2.9), Inches(8.6), Inches(2.9)),  # Step 2 to 3
        (Inches(7.5), Inches(3.5), Inches(6.2), Inches(3.5)),  # Step 3 to 4
        (Inches(7.4), Inches(4.4), Inches(7.3), Inches(4.4))   # Step 4 to 5
    ]
    
    for x1, y1, x2, y2 in arrows:
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x1, y1, Inches(0.2), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = text_color
        arrow.line.fill.background()
    
    content4.text = """⚡ AUTOMATED LOAN APPROVAL WORKFLOW

WORKFLOW STEPS:
1. Application Received - MSME submits loan application
2. Carbon Score Analysis - AI analyzes sustainability metrics
3. Risk Assessment - Automated risk scoring and evaluation
4. Auto Approval Decision - Instant approval based on carbon score
5. Loan Disbursement - Automated fund transfer and documentation

AUTOMATION FEATURES:
• Real-time carbon score integration
• AI-powered risk assessment
• Automated document verification
• Instant approval for high scores (80+)
• Dynamic interest rate calculation

PROCESSING TIME:
• Standard Process: 7-14 days
• Carbon Intelligence: 24-48 hours
• High Score (80+): Instant approval
• Medium Score (60-79): 24 hours
• Low Score (<60): Manual review

BENEFITS:
• 60% faster loan processing
• 40% reduction in manual work
• Consistent decision making
• Reduced human error
• Enhanced customer experience

SCORING-BASED APPROVAL:
• 90-100: Instant approval, 3% rate reduction
• 80-89: Fast track, 2% rate reduction
• 70-79: Standard processing, 1% rate reduction
• 60-69: Basic processing, 0.5% rate reduction
• <60: Manual review required"""
    
    # Slide 5: Risk Assessment Dashboard
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Risk Assessment Dashboard"
    
    # Create risk dashboard mockup
    # Dashboard container
    risk_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(4))
    risk_bg.fill.solid()
    risk_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
    risk_bg.line.color.rgb = dark_gray
    risk_bg.line.width = Pt(2)
    
    # Risk factors visualization
    # Carbon Score (40%)
    carbon_risk = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(2.3), Inches(1.5), Inches(1.2))
    carbon_risk.fill.solid()
    carbon_risk.fill.fore_color.rgb = primary_color
    carbon_risk.line.color.rgb = dark_gray
    carbon_risk.line.width = Pt(1)
    
    carbon_text = slide5.shapes.add_textbox(Inches(6.2), Inches(2.4), Inches(1.3), Inches(1))
    carbon_text.text_frame.text = "Carbon Score\n40%\n85/100"
    carbon_text.text_frame.paragraphs[0].font.size = Pt(10)
    carbon_text.text_frame.paragraphs[0].font.bold = True
    carbon_text.text_frame.paragraphs[0].font.color.rgb = white
    carbon_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Financial Health (30%)
    financial_risk = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(2.3), Inches(1.5), Inches(1.2))
    financial_risk.fill.solid()
    financial_risk.fill.fore_color.rgb = accent_color
    financial_risk.line.color.rgb = dark_gray
    financial_risk.line.width = Pt(1)
    
    financial_text = slide5.shapes.add_textbox(Inches(7.8), Inches(2.4), Inches(1.3), Inches(1))
    financial_text.text_frame.text = "Financial Health\n30%\nGood"
    financial_text.text_frame.paragraphs[0].font.size = Pt(10)
    financial_text.text_frame.paragraphs[0].font.bold = True
    financial_text.text_frame.paragraphs[0].font.color.rgb = white
    financial_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Industry Risk (15%)
    industry_risk = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(3.6), Inches(1.5), Inches(1.2))
    industry_risk.fill.solid()
    industry_risk.fill.fore_color.rgb = secondary_color
    industry_risk.line.color.rgb = dark_gray
    industry_risk.line.width = Pt(1)
    
    industry_text = slide5.shapes.add_textbox(Inches(6.2), Inches(3.7), Inches(1.3), Inches(1))
    industry_text.text_frame.text = "Industry Risk\n15%\nMedium"
    industry_text.text_frame.paragraphs[0].font.size = Pt(10)
    industry_text.text_frame.paragraphs[0].font.bold = True
    industry_text.text_frame.paragraphs[0].font.color.rgb = white
    industry_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Management Quality (10%)
    mgmt_risk = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(3.6), Inches(1.5), Inches(1.2))
    mgmt_risk.fill.solid()
    mgmt_risk.fill.fore_color.rgb = success_color
    mgmt_risk.line.color.rgb = dark_gray
    mgmt_risk.line.width = Pt(1)
    
    mgmt_text = slide5.shapes.add_textbox(Inches(7.8), Inches(3.7), Inches(1.3), Inches(1))
    mgmt_text.text_frame.text = "Management\n10%\nExcellent"
    mgmt_text.text_frame.paragraphs[0].font.size = Pt(10)
    mgmt_text.text_frame.paragraphs[0].font.bold = True
    mgmt_text.text_frame.paragraphs[0].font.color.rgb = white
    mgmt_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Overall Risk Score
    overall_risk = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(4.9), Inches(1.6), Inches(1.6))
    overall_risk.fill.solid()
    overall_risk.fill.fore_color.rgb = white
    overall_risk.line.color.rgb = success_color
    overall_risk.line.width = Pt(3)
    
    overall_text = slide5.shapes.add_textbox(Inches(7.3), Inches(5.3), Inches(1.4), Inches(0.8))
    overall_text.text_frame.text = "LOW RISK\nScore: 78"
    overall_text.text_frame.paragraphs[0].font.size = Pt(11)
    overall_text.text_frame.paragraphs[0].font.bold = True
    overall_text.text_frame.paragraphs[0].font.color.rgb = success_color
    overall_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content5.text = """🛡️ ADVANCED RISK ASSESSMENT DASHBOARD

RISK SCORING MATRIX:
• Carbon Score (40%): Primary sustainability indicator
• Financial Health (30%): Traditional financial metrics
• Industry Risk (15%): Sector-specific regulations
• Management Quality (10%): Leadership commitment
• Market Position (5%): Competitive advantage

PREDICTIVE ANALYTICS:
• Machine learning risk modeling
• Historical data analysis
• Trend prediction algorithms
• Early warning systems
• Automated risk monitoring

RISK MITIGATION:
• Real-time score monitoring
• Automated alerts and notifications
• Proactive intervention programs
• Sustainability improvement plans
• Regular portfolio reviews

PROVEN RESULTS:
• 45% lower default rate for high carbon scores
• 60% faster risk identification
• 35% reduction in loan loss provisions
• 50% improvement in portfolio quality
• 85% accuracy in risk prediction

COMPLIANCE FEATURES:
• Automated ESG reporting
• Regulatory compliance monitoring
• Audit trail maintenance
• Industry benchmark comparison
• Risk documentation automation"""
    
    # Slide 6: Carbon Trading Platform Interface
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Carbon Trading Platform Interface"
    
    # Create carbon trading interface mockup
    # Platform container
    trading_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(4))
    trading_bg.fill.solid()
    trading_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    trading_bg.line.color.rgb = dark_gray
    trading_bg.line.width = Pt(2)
    
    # Header
    trading_header = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(0.5))
    trading_header.fill.solid()
    trading_header.fill.fore_color.rgb = primary_color
    trading_header.line.fill.background()
    
    trading_header_text = slide6.shapes.add_textbox(Inches(6.1), Inches(2.1), Inches(3.3), Inches(0.3))
    trading_header_text.text_frame.text = "Carbon Credit Trading Platform"
    trading_header_text.text_frame.paragraphs[0].font.size = Pt(12)
    trading_header_text.text_frame.paragraphs[0].font.bold = True
    trading_header_text.text_frame.paragraphs[0].font.color.rgb = white
    trading_header_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Portfolio summary
    portfolio = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(2.6), Inches(3.3), Inches(1))
    portfolio.fill.solid()
    portfolio.fill.fore_color.rgb = white
    portfolio.line.color.rgb = RGBColor(200, 200, 200)
    portfolio.line.width = Pt(1)
    
    portfolio_text = slide6.shapes.add_textbox(Inches(6.2), Inches(2.7), Inches(3.1), Inches(0.8))
    portfolio_text.text_frame.text = """PORTFOLIO SUMMARY:
Total Credits: 1,250 tons CO2
Available: 850 tons
Used: 400 tons
Value: ₹2,50,000"""
    portfolio_text.text_frame.paragraphs[0].font.size = Pt(9)
    portfolio_text.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Available projects
    projects_header = slide6.shapes.add_textbox(Inches(6.1), Inches(3.7), Inches(3.3), Inches(0.3))
    projects_header.text_frame.text = "Available Carbon Projects"
    projects_header.text_frame.paragraphs[0].font.size = Pt(11)
    projects_header.text_frame.paragraphs[0].font.bold = True
    projects_header.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Project 1
    project1 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(4), Inches(3.3), Inches(0.6))
    project1.fill.solid()
    project1.fill.fore_color.rgb = RGBColor(240, 248, 255)
    project1.line.color.rgb = RGBColor(200, 200, 200)
    project1.line.width = Pt(1)
    
    project1_text = slide6.shapes.add_textbox(Inches(6.2), Inches(4.1), Inches(3.1), Inches(0.4))
    project1_text.text_frame.text = "Solar Energy Project - 500 tons - ₹200/ton - Verified"
    project1_text.text_frame.paragraphs[0].font.size = Pt(9)
    project1_text.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Project 2
    project2 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(4.7), Inches(3.3), Inches(0.6))
    project2.fill.solid()
    project2.fill.fore_color.rgb = RGBColor(240, 248, 255)
    project2.line.color.rgb = RGBColor(200, 200, 200)
    project2.line.width = Pt(1)
    
    project2_text = slide6.shapes.add_textbox(Inches(6.2), Inches(4.8), Inches(3.1), Inches(0.4))
    project2_text.text_frame.text = "Reforestation Project - 300 tons - ₹180/ton - Gold Standard"
    project2_text.text_frame.paragraphs[0].font.size = Pt(9)
    project2_text.text_frame.paragraphs[0].font.color.rgb = text_color
    
    content6.text = """🌍 CARBON TRADING PLATFORM INTERFACE

PLATFORM FEATURES:
• Real-time carbon credit marketplace
• Verified project portfolio
• Automated trading and settlement
• Portfolio tracking and analytics
• Price monitoring and alerts

CARBON CREDIT TYPES:
• Renewable Energy Projects
• Reforestation & Afforestation
• Energy Efficiency Initiatives
• Waste Management Projects
• Transportation Efficiency

TRADING CAPABILITIES:
• Buy/Sell carbon credits
• Portfolio management
• Automated offset purchasing
• Price alerts and notifications
• Transaction history tracking

BANKING INTEGRATION:
• Carbon credit-backed loans
• Offset portfolio financing
• Trading commission revenue
• Advisory services
• Market making services

FINANCIAL BENEFITS:
• Additional revenue stream for MSMEs
• Carbon credit financing options
• Cost-effective emission offsets
• Enhanced sustainability credentials
• Working capital opportunities

REGULATORY COMPLIANCE:
• VCS (Verified Carbon Standard)
• Gold Standard certification
• CDM project support
• Local carbon market compliance
• International trading standards"""
    
    # Slide 7: Mobile App Interface for Bankers
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Mobile App Interface for Bankers"
    
    # Create mobile app mockup
    # Phone frame
    phone_frame = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(2), Inches(3.5))
    phone_frame.fill.solid()
    phone_frame.fill.fore_color.rgb = RGBColor(50, 50, 50)
    phone_frame.line.color.rgb = dark_gray
    phone_frame.line.width = Pt(3)
    
    # Phone screen
    phone_screen = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(2.2), Inches(1.8), Inches(3.1))
    phone_screen.fill.solid()
    phone_screen.fill.fore_color.rgb = white
    phone_screen.line.fill.background()
    
    # App header
    app_header = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(2.2), Inches(1.8), Inches(0.4))
    app_header.fill.solid()
    app_header.fill.fore_color.rgb = primary_color
    app_header.line.fill.background()
    
    app_title = slide7.shapes.add_textbox(Inches(6.2), Inches(2.3), Inches(1.6), Inches(0.2))
    app_title.text_frame.text = "Carbon Intelligence"
    app_title.text_frame.paragraphs[0].font.size = Pt(10)
    app_title.text_frame.paragraphs[0].font.bold = True
    app_title.text_frame.paragraphs[0].font.color.rgb = white
    app_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Dashboard metrics
    metric1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(2.7), Inches(1.6), Inches(0.4))
    metric1.fill.solid()
    metric1.fill.fore_color.rgb = RGBColor(240, 240, 240)
    metric1.line.color.rgb = RGBColor(200, 200, 200)
    metric1.line.width = Pt(1)
    
    metric1_text = slide7.shapes.add_textbox(Inches(6.3), Inches(2.8), Inches(1.4), Inches(0.2))
    metric1_text.text_frame.text = "Active Loans: 1,247"
    metric1_text.text_frame.paragraphs[0].font.size = Pt(8)
    metric1_text.text_frame.paragraphs[0].font.bold = True
    metric1_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    metric2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(3.2), Inches(1.6), Inches(0.4))
    metric2.fill.solid()
    metric2.fill.fore_color.rgb = RGBColor(240, 240, 240)
    metric2.line.color.rgb = RGBColor(200, 200, 200)
    metric2.line.width = Pt(1)
    
    metric2_text = slide7.shapes.add_textbox(Inches(6.3), Inches(3.3), Inches(1.4), Inches(0.2))
    metric2_text.text_frame.text = "Avg Score: 78.5"
    metric2_text.text_frame.paragraphs[0].font.size = Pt(8)
    metric2_text.text_frame.paragraphs[0].font.bold = True
    metric2_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Quick actions
    actions_header = slide7.shapes.add_textbox(Inches(6.2), Inches(3.8), Inches(1.6), Inches(0.2))
    actions_header.text_frame.text = "Quick Actions"
    actions_header.text_frame.paragraphs[0].font.size = Pt(9)
    actions_header.text_frame.paragraphs[0].font.bold = True
    actions_header.text_frame.paragraphs[0].font.color.rgb = primary_color
    actions_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Action buttons
    action1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(4), Inches(0.7), Inches(0.3))
    action1.fill.solid()
    action1.fill.fore_color.rgb = secondary_color
    action1.line.fill.background()
    
    action1_text = slide7.shapes.add_textbox(Inches(6.25), Inches(4.05), Inches(0.6), Inches(0.2))
    action1_text.text_frame.text = "New Loan"
    action1_text.text_frame.paragraphs[0].font.size = Pt(7)
    action1_text.text_frame.paragraphs[0].font.bold = True
    action1_text.text_frame.paragraphs[0].font.color.rgb = white
    action1_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    action2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4), Inches(0.7), Inches(0.3))
    action2.fill.solid()
    action2.fill.fore_color.rgb = accent_color
    action2.line.fill.background()
    
    action2_text = slide7.shapes.add_textbox(Inches(6.95), Inches(4.05), Inches(0.6), Inches(0.2))
    action2_text.text_frame.text = "Reports"
    action2_text.text_frame.paragraphs[0].font.size = Pt(7)
    action2_text.text_frame.paragraphs[0].font.bold = True
    action2_text.text_frame.paragraphs[0].font.color.rgb = white
    action2_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    action3 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(4.4), Inches(0.7), Inches(0.3))
    action3.fill.solid()
    action3.fill.fore_color.rgb = success_color
    action3.line.fill.background()
    
    action3_text = slide7.shapes.add_textbox(Inches(6.25), Inches(4.45), Inches(0.6), Inches(0.2))
    action3_text.text_frame.text = "Alerts"
    action3_text.text_frame.paragraphs[0].font.size = Pt(7)
    action3_text.text_frame.paragraphs[0].font.bold = True
    action3_text.text_frame.paragraphs[0].font.color.rgb = white
    action3_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    action4 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.4), Inches(0.7), Inches(0.3))
    action4.fill.solid()
    action4.fill.fore_color.rgb = primary_color
    action4.line.fill.background()
    
    action4_text = slide7.shapes.add_textbox(Inches(6.95), Inches(4.45), Inches(0.6), Inches(0.2))
    action4_text.text_frame.text = "Trading"
    action4_text.text_frame.paragraphs[0].font.size = Pt(7)
    action4_text.text_frame.paragraphs[0].font.bold = True
    action4_text.text_frame.paragraphs[0].font.color.rgb = white
    action4_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content7.text = """📱 MOBILE APP INTERFACE FOR BANKERS

KEY FEATURES:
• Real-time portfolio monitoring
• Instant loan approval decisions
• Carbon score tracking and alerts
• Risk assessment on-the-go
• Carbon trading platform access

MOBILE DASHBOARD:
• Active loan portfolio overview
• Average carbon score tracking
• Risk alerts and notifications
• Quick action buttons
• Performance metrics

QUICK ACTIONS:
• New Loan Application - Instant processing
• Reports - Generate sustainability reports
• Alerts - Risk monitoring and notifications
• Trading - Carbon credit marketplace

BENEFITS FOR BANKERS:
• 24/7 portfolio monitoring
• Instant decision making
• Real-time risk alerts
• Mobile-first user experience
• Offline capability for key features

INTEGRATION FEATURES:
• Core banking system sync
• Real-time data updates
• Push notifications
• Biometric authentication
• Secure data transmission

PERFORMANCE METRICS:
• 60% faster loan processing
• 40% reduction in manual work
• 85% accuracy in risk assessment
• 24/7 availability
• 99.9% uptime guarantee"""
    
    # Slide 8: ROI & Financial Projections
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "ROI & Financial Projections"
    
    # Create ROI chart
    left_roi = Inches(6)
    top_roi = Inches(2)
    width_roi = Inches(3)
    height_roi = Inches(3)
    
    # ROI bars for 3 years
    years = ["Year 1", "Year 2", "Year 3"]
    roi_values = [25, 40, 60]  # ROI percentages
    colors = [warning_color, secondary_color, success_color]
    
    for i, (year, roi, color) in enumerate(zip(years, roi_values, colors)):
        bar_width = Inches(0.8)
        bar_height = Inches(roi/20)  # Scale down for visibility
        bar_x = left_roi + i * Inches(1)
        bar_y = top_roi + Inches(2.5) - bar_height
        
        roi_bar = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar_y, bar_width, bar_height)
        roi_bar.fill.solid()
        roi_bar.fill.fore_color.rgb = color
        roi_bar.line.fill.background()
        
        # Year label
        year_label = slide8.shapes.add_textbox(bar_x, top_roi + Inches(2.8), bar_width, Inches(0.3))
        year_label.text_frame.text = year
        year_label.text_frame.paragraphs[0].font.size = Pt(10)
        year_label.text_frame.paragraphs[0].font.bold = True
        year_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # ROI value
        roi_label = slide8.shapes.add_textbox(bar_x, bar_y - Inches(0.3), bar_width, Inches(0.3))
        roi_label.text_frame.text = f"{roi}%"
        roi_label.text_frame.paragraphs[0].font.size = Pt(12)
        roi_label.text_frame.paragraphs[0].font.bold = True
        roi_label.text_frame.paragraphs[0].font.color.rgb = color
        roi_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content8.text = """💰 COMPREHENSIVE FINANCIAL BENEFITS

ROI PROJECTIONS:
• Year 1: 20-25% ROI with pilot program
• Year 2: 35-45% ROI with full rollout
• Year 3: 50-65% ROI with market expansion
• Break-even: 6-9 months
• Payback period: 12-18 months

REVENUE OPPORTUNITIES (3-Year Projection):
• Green Loan Origination: ₹500-800 crores annually
• Interest Rate Premium: 1-3% on green loans
• Carbon Trading Commissions: ₹50-100 crores annually
• Sustainability Advisory: ₹25-50 crores annually
• Platform Subscription Fees: ₹10-20 crores annually

COST SAVINGS:
• 50% reduction in manual ESG assessment costs
• 60% faster loan processing and approval
• 40% reduction in default rates and provisions
• 30% improvement in operational efficiency
• 25% reduction in compliance and reporting costs

RISK REDUCTION BENEFITS:
• 45% lower default rate for green loan portfolio
• 60% faster identification of at-risk accounts
• 35% reduction in loan loss provisions
• 50% improvement in portfolio quality
• 40% reduction in regulatory compliance costs

MARKET EXPANSION:
• 25-40% increase in MSME loan portfolio
• 30-50% growth in green finance market share
• 20-30% improvement in customer retention
• 15-25% increase in average loan size
• 10-20% premium on interest rates"""
    
    # Slide 9: Implementation Roadmap
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Implementation Roadmap"
    
    # Create timeline visualization
    left_timeline = Inches(6)
    top_timeline = Inches(2)
    
    # Timeline phases
    phases = [
        ("Phase 1\nFoundation\n(1-3 months)", primary_color, left_timeline, top_timeline),
        ("Phase 2\nPilot\n(4-9 months)", secondary_color, left_timeline + Inches(1.5), top_timeline),
        ("Phase 3\nRollout\n(10-18 months)", accent_color, left_timeline + Inches(3), top_timeline),
        ("Phase 4\nExpansion\n(19-24 months)", success_color, left_timeline + Inches(4.5), top_timeline)
    ]
    
    for text, color, x, y in phases:
        phase_box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.3), Inches(1.5))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = color
        phase_box.line.color.rgb = RGBColor(100, 100, 100)
        phase_box.line.width = Pt(2)
        
        phase_text = slide9.shapes.add_textbox(x + Inches(0.1), y + Inches(0.2), Inches(1.1), Inches(1.1))
        phase_text.text_frame.text = text
        phase_text.text_frame.paragraphs[0].font.size = Pt(10)
        phase_text.text_frame.paragraphs[0].font.bold = True
        phase_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        phase_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Connect phases with arrows
    for i in range(len(phases) - 1):
        arrow_x = left_timeline + Inches(1.3) + i * Inches(1.5)
        arrow_y = top_timeline + Inches(0.75)
        arrow = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = text_color
        arrow.line.fill.background()
    
    content9.text = """🚀 PHASED IMPLEMENTATION PLAN

PHASE 1: FOUNDATION (Months 1-3)
• Carbon Intelligence platform integration
• API development and testing
• Staff training and certification
• Pilot customer selection (100 MSMEs)
• Regulatory compliance review and approval

PHASE 2: PILOT PROGRAM (Months 4-9)
• Launch with 100 selected MSMEs
• Green loan product testing and optimization
• Carbon trading platform integration
• Performance monitoring and analytics
• Feedback collection and system refinement

PHASE 3: FULL ROLLOUT (Months 10-18)
• Scale to 1,000+ MSMEs
• Complete product suite launch
• Advanced analytics and AI features
• Marketing and customer acquisition
• Performance optimization and scaling

PHASE 4: EXPANSION (Months 19-24)
• Scale to 5,000+ MSMEs
• Advanced AI features and automation
• Carbon trading marketplace expansion
• International market entry
• Innovation and R&D initiatives

SUCCESS METRICS:
• Customer acquisition rate and retention
• Loan portfolio growth and quality
• Revenue per customer and profitability
• Platform adoption and engagement
• Customer satisfaction and NPS scores

SUPPORT & TRAINING:
• Dedicated implementation team
• 24/7 technical support and monitoring
• Regular training and certification programs
• Comprehensive documentation and resources
• Continuous optimization and improvement"""
    
    # Slide 10: Call to Action & Next Steps
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Call to Action & Next Steps"
    
    # Create action items visualization
    left_actions = Inches(6)
    top_actions = Inches(2)
    
    # Action items with checkboxes
    actions = [
        "Partnership Agreement",
        "Technical Integration",
        "Pilot Program Setup",
        "Marketing & Launch"
    ]
    
    for i, action in enumerate(actions):
        # Checkbox
        checkbox = slide10.shapes.add_shape(MSO_SHAPE.OVAL, left_actions, top_actions + i * Inches(0.6), Inches(0.3), Inches(0.3))
        checkbox.fill.solid()
        checkbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        checkbox.line.color.rgb = primary_color
        checkbox.line.width = Pt(2)
        
        # Action text
        action_text = slide10.shapes.add_textbox(left_actions + Inches(0.4), top_actions + i * Inches(0.6), Inches(2.5), Inches(0.3))
        action_text.text_frame.text = action
        action_text.text_frame.paragraphs[0].font.size = Pt(12)
        action_text.text_frame.paragraphs[0].font.bold = True
        action_text.text_frame.paragraphs[0].font.color.rgb = text_color
    
    content10.text = """🎯 IMMEDIATE ACTION ITEMS

1. PARTNERSHIP AGREEMENT:
• Review partnership terms and conditions
• Finalize revenue sharing model (70% Bank, 30% Carbon Intelligence)
• Sign memorandum of understanding
• Establish governance structure and committees

2. TECHNICAL INTEGRATION:
• Schedule technical assessment and planning
• Plan API integration timeline and milestones
• Design custom features and white-label branding
• Set up development and testing environment

3. PILOT PROGRAM SETUP:
• Select pilot MSME customers (100-200)
• Train banking staff on Carbon Intelligence platform
• Launch pilot program with monitoring
• Collect feedback and optimize performance

4. MARKETING & LAUNCH:
• Develop co-marketing strategy and materials
• Create customer acquisition and retention plan
• Launch green finance products and services
• Execute go-to-market strategy and campaigns

TIMELINE:
• Week 1-2: Partnership agreement finalization
• Week 3-4: Technical integration planning
• Month 2-3: Platform customization and testing
• Month 4-6: Pilot program execution
• Month 7-12: Full rollout and scaling

CONTACT INFORMATION:
• Email: partnerships@carbonintelligence.com
• Phone: +91-98765-43210
• Website: www.carbonintelligence.com
• LinkedIn: Carbon Intelligence

READY TO REVOLUTIONIZE GREEN FINANCE?
Let's discuss how Carbon Intelligence can transform your MSME lending business and unlock new opportunities in sustainable finance.

Schedule a detailed discussion today!"""
    
    # Format all slides consistently
    for slide in prs.slides:
        # Format titles
        if slide.shapes.title:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        
        # Format content
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraph.font.color.rgb = text_color
                        paragraph.font.size = Pt(16)
    
    # Save the presentation
    output_path = "/workspace/pitches/Carbon_Intelligence_Enhanced_Bank_Pitch_With_Mockups.pptx"
    prs.save(output_path)
    print(f"Enhanced Bank PowerPoint presentation with mockups created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_enhanced_bank_pitch_with_mockups()