#!/usr/bin/env python3
"""
Create 10-Page MSME Carbon Intelligence Presentation
Comprehensive presentation with app mockups, images, and detailed content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_msme_presentation():
    """Create a comprehensive 10-page MSME presentation"""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    primary_green = RGBColor(76, 175, 80)  # #4CAF50
    secondary_orange = RGBColor(255, 140, 0)  # #FF8C00
    accent_blue = RGBColor(25, 118, 210)  # #1976D2
    dark_green = RGBColor(46, 125, 50)  # #2E7D32
    light_gray = RGBColor(245, 245, 245)  # #F5F5F5
    dark_gray = RGBColor(66, 66, 66)  # #424242
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "🌱 Carbon Intelligence"
    subtitle.text = "AI-Powered Carbon Intelligence Platform for MSMEs\n\nTransform Your Business with Sustainable Intelligence\n\nMicro, Small & Medium Enterprises Carbon Footprint Management"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    title.text_frame.paragraphs[0].font.bold = True
    
    # Style subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = dark_gray
    subtitle.text_frame.paragraphs[1].font.size = Pt(16)
    subtitle.text_frame.paragraphs[1].font.color.rgb = accent_blue
    subtitle.text_frame.paragraphs[2].font.size = Pt(14)
    subtitle.text_frame.paragraphs[2].font.color.rgb = dark_gray
    
    # Add company info
    left = Inches(1)
    top = Inches(6)
    width = Inches(11.33)
    height = Inches(1)
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "📧 contact@carbonintelligence.com | 📱 +91-98765-43210 | 🌐 www.carbonintelligence.com"
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.color.rgb = dark_gray
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "The MSME Sustainability Challenge"
    content.text = """🚨 Critical Problems MSMEs Face:

• High Operational Costs: 60-70% of revenue spent on energy, water, waste management
• Limited Access to Green Finance: Traditional banks don't understand sustainability metrics
• Complex Compliance: ESG reporting requirements increasing annually
• Lack of Data: No systematic way to track carbon footprint
• Competitive Pressure: Customers demanding sustainable practices
• Regulatory Risk: Facing penalties for non-compliance

💰 Financial Impact:
• Average MSME loses ₹2-5 lakhs annually due to inefficient resource usage
• 40% higher operational costs compared to sustainable competitors
• Limited access to low-interest green loans
• Missed opportunities in carbon credit trading

🎯 The Opportunity:
• ₹50,000 crores green finance market in India
• 30% cost reduction potential through sustainability
• Enhanced brand reputation and customer trust
• Future-proof business operations"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(16)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 53, 69)  # Red for problems
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 3: Solution Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    
    title.text = "🌱 Carbon Intelligence Solution"
    content.text = """🤖 AI-Powered Carbon Intelligence Platform

✅ Complete Carbon Footprint Tracking
• Real-time monitoring of all emission sources
• Automated data collection from multiple sources
• AI-powered analysis and insights

✅ Smart Cost Optimization
• 20-40% reduction in operational costs
• Personalized recommendations for efficiency
• ROI tracking and financial projections

✅ Green Finance Access
• Carbon score-based loan eligibility
• 1-3% lower interest rates
• Faster approval process (24-48 hours)

✅ Automated Compliance
• ESG reporting automation
• Regulatory compliance monitoring
• Audit-ready documentation

✅ Carbon Trading Platform
• Buy/sell carbon credits
• Additional revenue stream
• Market price optimization

📱 Multi-Platform Access:
• Web Dashboard for comprehensive management
• Mobile App for on-the-go monitoring
• API integration with existing systems"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 4: App Mockups - Dashboard
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title.text = "📱 Application Interface Mockups"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    # Add mockup description
    desc = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(1))
    desc.text = "Professional Dashboard with Real-time Carbon Metrics and Quick Actions"
    desc.text_frame.paragraphs[0].font.size = Pt(16)
    desc.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Create dashboard mockup
    dashboard_left = Inches(1)
    dashboard_top = Inches(2.5)
    dashboard_width = Inches(5.5)
    dashboard_height = Inches(4)
    
    # Dashboard background
    dashboard_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dashboard_left, dashboard_top, dashboard_width, dashboard_height)
    dashboard_bg.fill.solid()
    dashboard_bg.fill.fore_color.rgb = light_gray
    dashboard_bg.line.color.rgb = primary_green
    dashboard_bg.line.width = Pt(2)
    
    # Dashboard header
    header = slide4.shapes.add_textbox(dashboard_left + Inches(0.2), dashboard_top + Inches(0.2), dashboard_width - Inches(0.4), Inches(0.8))
    header.text = "🌱 Carbon Intelligence Dashboard"
    header.text_frame.paragraphs[0].font.size = Pt(14)
    header.text_frame.paragraphs[0].font.bold = True
    header.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    # Carbon score display
    score_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dashboard_left + Inches(0.3), dashboard_top + Inches(1.2), Inches(2), Inches(1.2))
    score_bg.fill.solid()
    score_bg.fill.fore_color.rgb = primary_green
    
    score_text = slide4.shapes.add_textbox(dashboard_left + Inches(0.4), dashboard_top + Inches(1.3), Inches(1.8), Inches(1))
    score_text.text = "85\nCarbon Score\nGold Tier"
    score_text.text_frame.paragraphs[0].font.size = Pt(16)
    score_text.text_frame.paragraphs[0].font.bold = True
    score_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    score_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Metrics
    metrics = [
        ("CO₂ Saved", "245.6 kg", primary_green),
        ("Cost Savings", "₹15,240", accent_blue),
        ("Reduction", "12.3%", secondary_orange),
        ("Recommendations", "8 Active", dark_green)
    ]
    
    for i, (label, value, color) in enumerate(metrics):
        metric_left = dashboard_left + Inches(2.5) + (i % 2) * Inches(1.4)
        metric_top = dashboard_top + Inches(1.2) + (i // 2) * Inches(0.8)
        
        metric_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, metric_left, metric_top, Inches(1.2), Inches(0.6))
        metric_bg.fill.solid()
        metric_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        metric_bg.line.color.rgb = color
        metric_bg.line.width = Pt(1)
        
        metric_text = slide4.shapes.add_textbox(metric_left + Inches(0.1), metric_top + Inches(0.1), Inches(1), Inches(0.4))
        metric_text.text = f"{value}\n{label}"
        metric_text.text_frame.paragraphs[0].font.size = Pt(10)
        metric_text.text_frame.paragraphs[0].font.bold = True
        metric_text.text_frame.paragraphs[0].font.color.rgb = color
        metric_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Quick actions
    actions_text = slide4.shapes.add_textbox(dashboard_left + Inches(0.3), dashboard_top + Inches(2.8), dashboard_width - Inches(0.6), Inches(0.5))
    actions_text.text = "Quick Actions: Carbon Assessment | View Analytics | Recommendations | Generate Report"
    actions_text.text_frame.paragraphs[0].font.size = Pt(10)
    actions_text.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Mobile app mockup
    mobile_left = Inches(7)
    mobile_top = Inches(2.5)
    mobile_width = Inches(2.5)
    mobile_height = Inches(4)
    
    # Mobile background
    mobile_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mobile_left, mobile_top, mobile_width, mobile_height)
    mobile_bg.fill.solid()
    mobile_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
    mobile_bg.line.color.rgb = dark_gray
    mobile_bg.line.width = Pt(3)
    
    # Mobile screen
    mobile_screen = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mobile_left + Inches(0.1), mobile_top + Inches(0.2), mobile_width - Inches(0.2), mobile_height - Inches(0.4))
    mobile_screen.fill.solid()
    mobile_screen.fill.fore_color.rgb = light_gray
    
    # Mobile header
    mobile_header = slide4.shapes.add_textbox(mobile_left + Inches(0.2), mobile_top + Inches(0.3), mobile_width - Inches(0.4), Inches(0.6))
    mobile_header.text = "🌱 Carbon Intelligence\nMobile App"
    mobile_header.text_frame.paragraphs[0].font.size = Pt(10)
    mobile_header.text_frame.paragraphs[0].font.bold = True
    mobile_header.text_frame.paragraphs[0].font.color.rgb = primary_green
    mobile_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Mobile content
    mobile_content = slide4.shapes.add_textbox(mobile_left + Inches(0.2), mobile_top + Inches(1), mobile_width - Inches(0.4), Inches(2.5))
    mobile_content.text = """📊 Real-time Dashboard
• Carbon Score: 85
• CO₂ Saved: 245 kg
• Cost Savings: ₹15,240

🔍 Quick Actions
• Analyze SMS
• View Analytics
• Check Incentives
• Generate Reports

📱 24/7 Access
• Push Notifications
• Offline Mode
• GPS Tracking
• Voice Commands"""
    mobile_content.text_frame.paragraphs[0].font.size = Pt(8)
    mobile_content.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Slide 5: Key Features & Benefits
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    content = slide5.placeholders[1]
    
    title.text = "🚀 Key Features & Benefits"
    content.text = """🎯 Core Features:

📊 AI-Powered Carbon Intelligence
• Real-time carbon footprint tracking
• Automated data collection from SMS, emails, documents
• Machine learning-based insights and predictions
• Industry benchmarking and performance comparison

💰 Smart Cost Optimization
• 20-40% reduction in operational costs
• Personalized efficiency recommendations
• ROI calculator and financial projections
• Money-back guarantee (5x subscription cost savings)

🏦 Green Finance Integration
• Carbon score-based loan eligibility
• 1-3% lower interest rates than traditional loans
• 24-48 hour approval process
• Access to ₹50L+ green finance market

📈 Carbon Trading Platform
• Buy/sell carbon credits
• Market price optimization
• Additional revenue stream (₹50,000-₹2,00,000 annually)
• Automated trading recommendations

📱 Multi-Platform Access
• Web dashboard for comprehensive management
• React Native mobile app for on-the-go access
• API integration with existing systems
• Offline capability for key features

🤖 Multi-Agent AI System
• Carbon Intelligence Agent for footprint analysis
• Green Finance Agent for loan recommendations
• ESG Compliance Agent for automated reporting
• Analytics Agent for trend analysis and insights"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 6: Financial Impact & ROI
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    content = slide6.placeholders[1]
    
    title.text = "💰 Financial Impact & ROI Analysis"
    content.text = """📈 Cost Savings Breakdown (Annual):

⚡ Energy Costs: 25-40% reduction
• Average savings: ₹75,000 - ₹3,00,000
• LED lighting, smart meters, renewable energy

💧 Water Management: 20-30% reduction
• Average savings: ₹40,000 - ₹1,50,000
• Water recycling, efficient fixtures

🗑️ Waste Management: 35-50% reduction
• Average savings: ₹25,000 - ₹1,00,000
• Recycling programs, waste reduction

🚚 Transportation: 15-25% reduction
• Average savings: ₹30,000 - ₹1,50,000
• Route optimization, fuel efficiency

📦 Materials: 20-30% reduction
• Average savings: ₹50,000 - ₹2,00,000
• Sustainable sourcing, waste reduction

💵 Total Annual Savings: ₹2,20,000 - ₹9,00,000 per MSME

📊 ROI Calculations:
• Platform Cost: ₹2,999 - ₹7,999 per month
• Average Annual Savings: ₹2,20,000 - ₹9,00,000
• ROI: 400-600% in first year
• Payback Period: 2-3 months
• Break-even: Immediate

🎯 Green Finance Benefits:
• 1-3% lower interest rates
• Faster loan approval (24-48 hours vs 7-14 days)
• Reduced collateral requirements
• Government incentives and subsidies
• Carbon credit trading revenue

💎 Money-Back Guarantee:
• 5x subscription cost savings guarantee
• 30-day free trial
• No long-term contracts"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 7: Success Stories & Case Studies
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    content = slide7.placeholders[1]
    
    title.text = "🏆 Success Stories & Case Studies"
    content.text = """📊 Case Study 1: EcoTech Manufacturing
Industry: Electronics manufacturing (75 employees)
Carbon Score: 85 (Gold tier)
Results after 6 months:
• 40% reduction in energy costs (₹2,50,000 saved)
• 30% reduction in water usage (₹1,20,000 saved)
• 45% reduction in waste generation (₹80,000 saved)
• 50% improvement in sustainability score
• Secured green loan at 2.5% lower interest rate
• Earned ₹1,50,000 from carbon credit trading
• Total annual savings: ₹6,00,000

📊 Case Study 2: GreenTextile Ltd
Industry: Textile manufacturing (100 employees)
Carbon Score: 78 (Silver tier)
Results after 4 months:
• 35% reduction in energy consumption
• 25% reduction in water usage
• 30% reduction in waste generation
• 45% improvement in sustainability score
• Generated ₹1,20,000 additional revenue
• Total annual savings: ₹4,20,000

📊 Case Study 3: FreshFoods Pvt Ltd
Industry: Food processing (50 employees)
Carbon Score: 72 (Silver tier)
Results after 3 months:
• 30% reduction in energy costs
• 20% reduction in water usage
• 40% reduction in waste disposal costs
• 35% improvement in carbon score
• Created ₹80,000 new income stream
• Total annual savings: ₹2,50,000

🎯 Average Results Across All MSMEs:
• 35% reduction in operational costs
• 40% improvement in carbon scores
• 25% increase in loan approval rates
• 30% reduction in default rates
• 98% customer satisfaction score"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 8: Technology & AI Features
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    content = slide8.placeholders[1]
    
    title.text = "🤖 Advanced Technology & AI Features"
    content.text = """🧠 AI-Powered Multi-Agent System:

🤖 Carbon Intelligence Agent
• Real-time carbon footprint analysis
• Automated data collection from multiple sources
• Machine learning-based emission predictions
• Industry-specific benchmarking

💰 Green Finance Agent
• Carbon score-based loan recommendations
• Market analysis for optimal financing
• Automated loan application processing
• Risk assessment and credit scoring

📊 ESG Compliance Agent
• Automated ESG reporting generation
• Regulatory compliance monitoring
• Audit trail maintenance
• Stakeholder communication

📈 Analytics Agent
• Trend analysis and forecasting
• Performance optimization recommendations
• Cost-benefit analysis for sustainability initiatives
• Predictive maintenance alerts

🔍 Data Processing Capabilities:
• SMS Analysis: Extract carbon-related insights from business communications
• Email Analysis: Process sustainability-related emails and documents
• Document Upload: AI-powered analysis of bills, receipts, and reports
• GPS Tracking: Monitor transportation and logistics emissions
• Voice Commands: Hands-free operation for mobile users

🌐 Platform Architecture:
• React Web Application for comprehensive management
• React Native Mobile App for on-the-go access
• Node.js Backend with microservices architecture
• PostgreSQL database with real-time synchronization
• AWS cloud infrastructure with 99.9% uptime
• RESTful API for third-party integrations

🔒 Security & Compliance:
• End-to-end encryption for all data
• GDPR and data privacy compliance
• SOC 2 Type II certification
• Regular security audits and penetration testing
• Role-based access control and audit logs"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 9: Pricing & Implementation
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide9.shapes.title
    content = slide9.placeholders[1]
    
    title.text = "💳 Pricing & Implementation Plan"
    content.text = """💰 Flexible Pricing Options:

🥉 Starter Plan - ₹2,999/month
• Up to 25 employees
• Basic carbon tracking
• Standard analytics
• Email support
• Mobile app access

🥈 Professional Plan - ₹4,999/month (Most Popular)
• Up to 100 employees
• Advanced AI analytics
• Green finance integration
• Carbon trading platform
• Priority support
• Custom reporting

🥇 Enterprise Plan - ₹7,999/month
• Unlimited employees
• Full AI multi-agent system
• Custom integrations
• Dedicated account manager
• Advanced compliance features
• White-label options

🎁 Special Launch Offers:
• 30-day free trial (no credit card required)
• 20% discount for annual payment
• First 3 months at 50% off
• Free setup and training
• No long-term contracts
• Money-back guarantee

🚀 Easy Implementation Process:

Step 1: Quick Setup (1 day)
• Download mobile app
• Create company account
• Connect existing accounts (bank, utilities)

Step 2: Data Integration (1 week)
• AI analysis of historical data
• Categorization and initial scoring
• Custom dashboard configuration

Step 3: Start Saving (Immediately)
• Implement AI recommendations
• Track progress in real-time
• Access green finance options

📞 Support Available:
• Phone: +91-98765-43210
• Email: support@carbonintelligence.com
• Live chat on website
• In-app support
• Video tutorials and guides
• Dedicated customer success team"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Slide 10: Call to Action
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide10.shapes.title
    content = slide10.placeholders[1]
    
    title.text = "🚀 Ready to Transform Your Business?"
    content.text = """🎯 Immediate Benefits You'll Get:

✅ Start saving money from day 1
✅ Access to green finance immediately
✅ Improve operational efficiency
✅ Build competitive advantage
✅ Meet customer expectations
✅ Future-proof your business

🌟 Why Choose Carbon Intelligence:

🏆 Proven Results: 35% average cost reduction
🤖 AI-Powered: Advanced multi-agent system
💰 ROI Guaranteed: 5x subscription cost savings
📱 Multi-Platform: Web and mobile access
🔒 Secure: Enterprise-grade security
🌱 Sustainable: Real environmental impact

🎁 Special Launch Offer:

• 30-day free trial
• 50% off first 3 months
• Free setup and training
• Money-back guarantee
• No long-term commitment

📞 Get Started Today:

🌐 Visit: www.carbonintelligence.com
📧 Email: contact@carbonintelligence.com
📱 Call: +91-98765-43210
💬 Live Chat: Available 24/7

🚀 Transform Your MSME Today!

Join 500+ MSMEs already saving money and reducing their carbon footprint with Carbon Intelligence.

The future of sustainable business is here. Are you ready to be part of it?

#CarbonIntelligence #SustainableMSME #GreenFinance #AIAnalytics #CarbonTrading"""
    
    # Style content
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.bold = True
    content.text_frame.paragraphs[0].font.color.rgb = primary_green
    
    for i in range(1, len(content.text_frame.paragraphs)):
        content.text_frame.paragraphs[i].font.size = Pt(14)
        content.text_frame.paragraphs[i].font.color.rgb = dark_gray
    
    # Save presentation
    output_path = "/workspace/pitches/Carbon_Intelligence_10_Page_MSME_Presentation.pptx"
    prs.save(output_path)
    
    print(f"✅ 10-page MSME presentation created successfully!")
    print(f"📁 File saved: {output_path}")
    print(f"📊 Slides: 10 comprehensive slides")
    print(f"🎯 Target: Micro, Small & Medium Enterprises")
    print(f"📱 Features: App mockups, financial analysis, case studies")
    
    return output_path

if __name__ == "__main__":
    create_msme_presentation()